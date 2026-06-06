#!/usr/bin/env python3
"""Generate NiceEze architecture diagram PPTX files."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── colour palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x1a, 0x1f, 0x2e)
ACCENT   = RGBColor(0xEC, 0x6D, 0x74)
BLUE     = RGBColor(0x6E, 0xA7, 0xDB)
GREEN    = RGBColor(0x4C, 0xAF, 0x50)
ORANGE   = RGBColor(0xFF, 0x98, 0x00)
PURPLE   = RGBColor(0x9C, 0x27, 0xB0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x78, 0x78, 0x78)
CARD     = RGBColor(0x22, 0x28, 0x40)
DARK_GRAY= RGBColor(0x55, 0x55, 0x66)
LIGHT_GRAY=RGBColor(0xAA, 0xAA, 0xBB)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(layout)
    return slide


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor, text="", font_size=11,
             bold=False, text_color=WHITE, line_color=None, line_width=Pt(0),
             v_align="middle", h_align=PP_ALIGN.CENTER, rounded=False):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    # rounded corners via XML if requested
    if rounded:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            else:
                avLst.clear()
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', 'val 30000')

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill

    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True

    # vertical alignment
    from pptx.enum.text import MSO_ANCHOR
    if v_align == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif v_align == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if text:
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = h_align
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = text_color
    else:
        tf.paragraphs[0].text = ""

    return shape


def add_text_box(slide, x, y, w, h, text, font_size=11, bold=False,
                 color=WHITE, h_align=PP_ALIGN.LEFT, v_align="top"):
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    if v_align == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = h_align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


def add_arrow(slide, x1, y1, x2, y2, label="", color=LIGHT_GRAY):
    """Draw a simple arrow line (connector) with optional mid-label."""
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    # Use a line connector
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.2)

    # arrowhead via XML
    ln = connector._element.find('.//' + qn('a:ln'))
    if ln is None:
        ln = connector._element.find(qn('p:spPr') + '/' + qn('a:ln'))

    # Try to set end arrowhead
    try:
        spPr = connector._element.find(qn('p:spPr'))
        ln_el = spPr.find(qn('a:ln'))
        if ln_el is None:
            ln_el = etree.SubElement(spPr, qn('a:ln'))
        # remove existing tail elements if any
        for old in ln_el.findall(qn('a:tailEnd')):
            ln_el.remove(old)
        for old in ln_el.findall(qn('a:headEnd')):
            ln_el.remove(old)
        tailEnd = etree.SubElement(ln_el, qn('a:tailEnd'))
        tailEnd.set('type', 'none')
        headEnd = etree.SubElement(ln_el, qn('a:headEnd'))
        headEnd.set('type', 'arrow')
        headEnd.set('w', 'med')
        headEnd.set('len', 'med')
    except Exception:
        pass

    if label:
        mx = (x1 + x2) // 2
        my = (y1 + y2) // 2
        add_text_box(slide, mx - Inches(0.5), my - Inches(0.15),
                     Inches(1.2), Inches(0.3), label,
                     font_size=7, color=LIGHT_GRAY, h_align=PP_ALIGN.CENTER)
    return connector


# ═══════════════════════════════════════════════════════════════════════════════
# FILE 1: 全体システム構成図
# ═══════════════════════════════════════════════════════════════════════════════

def build_file1():
    prs = new_prs()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = blank_slide(prs)
    set_bg(slide, BG)

    # decorative accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), ACCENT)

    # NiceEze logo text
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(3), Inches(0.4),
                 "NiceEze", font_size=13, bold=True, color=ACCENT)

    # title
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.4),
                 "NiceEze 自律経営執行システム v14.2",
                 font_size=36, bold=True, color=WHITE,
                 h_align=PP_ALIGN.CENTER)

    # subtitle
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
                 "全体システム構成図",
                 font_size=22, bold=False, color=BLUE,
                 h_align=PP_ALIGN.CENTER)

    # accent line under subtitle
    add_rect(slide, Inches(4.5), Inches(4.7), Inches(4.3), Inches(0.05), ACCENT)

    # bottom tagline
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.4),
                 "自律型マルチエージェント経営基盤  |  5 Subsystems  |  Gate G0 → G4",
                 font_size=11, color=LIGHT_GRAY, h_align=PP_ALIGN.CENTER)

    add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT)

    # ── Slide 2: システム全体マップ ───────────────────────────────────────────
    slide = blank_slide(prs)
    set_bg(slide, BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), ACCENT)

    add_text_box(slide, Inches(0.3), Inches(0.12), Inches(8), Inches(0.45),
                 "システム全体マップ  —  5 Subsystems & External APIs",
                 font_size=16, bold=True, color=WHITE)
    add_text_box(slide, Inches(9.5), Inches(0.12), Inches(3.5), Inches(0.45),
                 "NiceEze v14.2", font_size=11, bold=True, color=ACCENT,
                 h_align=PP_ALIGN.RIGHT)

    # ── subsystem boxes (row 1: RESEARCH, SURPLUS SHIFT, MARKETING) ──
    sys_y = Inches(0.8)
    sys_h = Inches(1.3)
    sys_w = Inches(2.4)

    # RESEARCH :8080
    add_rect(slide, Inches(0.3), sys_y, sys_w, sys_h, CARD,
             text="RESEARCH\n:8080\n市場調査・価格分析",
             font_size=10, bold=False, text_color=BLUE,
             line_color=BLUE, line_width=Pt(1.5), rounded=True)

    # SURPLUS SHIFT :8081
    add_rect(slide, Inches(3.1), sys_y, sys_w, sys_h, CARD,
             text="SURPLUS SHIFT\n:8081\n不動産余剰スペース商談",
             font_size=10, bold=False, text_color=ACCENT,
             line_color=ACCENT, line_width=Pt(1.5), rounded=True)

    # MARKETING
    add_rect(slide, Inches(5.9), sys_y, sys_w, sys_h, CARD,
             text="MARKETING\n:N/A\nX投稿・コンテンツ生成",
             font_size=10, bold=False, text_color=GREEN,
             line_color=GREEN, line_width=Pt(1.5), rounded=True)

    # GOV :8082 — row 2 left
    gov_y = Inches(2.6)
    add_rect(slide, Inches(2.0), gov_y, sys_w, sys_h, CARD,
             text="GOV\n:8082\nCOOレポート・FinOps・稼働ログ",
             font_size=10, bold=False, text_color=ORANGE,
             line_color=ORANGE, line_width=Pt(1.5), rounded=True)

    # SBDS :8083
    add_rect(slide, Inches(5.9), gov_y, sys_w, sys_h, CARD,
             text="SBDS\n:8083\n手ぶら旅行・配送管理",
             font_size=10, bold=False, text_color=PURPLE,
             line_color=PURPLE, line_width=Pt(1.5), rounded=True)

    # ── external API boxes ──
    ext_x = Inches(9.5)
    ext_w = Inches(3.5)
    ext_h = Inches(0.55)
    ext_gap = Inches(0.7)
    ext_y0 = Inches(0.85)

    ext_labels = [
        "Claude API  (Anthropic)",
        "X (Twitter) API",
        "楽天 / Yahoo! / Keepa API",
        "ヤマト / 佐川  Webhook",
        "SmartLife IoT",
    ]
    ext_ys = []
    for i, lbl in enumerate(ext_labels):
        ey = ext_y0 + i * ext_gap
        ext_ys.append(ey)
        add_rect(slide, ext_x, ey, ext_w, ext_h, DARK_GRAY,
                 text=lbl, font_size=9, text_color=LIGHT_GRAY,
                 line_color=GRAY, line_width=Pt(0.75), rounded=True)

    # ── section label ──
    add_rect(slide, ext_x - Inches(0.05), Inches(0.72), ext_w + Inches(0.1), Inches(0.3),
             BG, text="External Systems", font_size=8, text_color=GRAY)

    # ── arrows between subsystems ──
    # Helper to get center of box
    def cx(lx): return lx + sys_w / 2
    def cy(ly): return ly + sys_h / 2

    # RESEARCH → SURPLUS SHIFT: "市場価格データ"
    add_arrow(slide,
              Inches(0.3) + sys_w, sys_y + sys_h//2,
              Inches(3.1), sys_y + sys_h//2,
              label="市場価格データ", color=BLUE)

    # RESEARCH → MARKETING: "トレンド情報"  (diagonal)
    add_arrow(slide,
              cx(Inches(0.3)), sys_y + sys_h,
              cx(Inches(5.9)), sys_y + sys_h,
              label="トレンド情報", color=BLUE)

    # SURPLUS SHIFT → GOV: "商談KPI"
    add_arrow(slide,
              cx(Inches(3.1)), sys_y + sys_h,
              cx(Inches(2.0)), gov_y,
              label="商談KPI", color=ACCENT)

    # MARKETING → GOV: "投稿実績"
    add_arrow(slide,
              cx(Inches(5.9)), sys_y + sys_h,
              cx(Inches(2.0)) + Inches(0.3), gov_y,
              label="投稿実績", color=GREEN)

    # SBDS → GOV: "配送稼働ログ"
    add_arrow(slide,
              Inches(5.9), gov_y + sys_h//2,
              Inches(2.0) + sys_w, gov_y + sys_h//2,
              label="配送稼働ログ", color=PURPLE)

    # RESEARCH → 楽天/Yahoo!/Keepa: "価格取得"
    add_arrow(slide,
              Inches(0.3) + sys_w, sys_y + Inches(0.3),
              ext_x, ext_ys[2] + ext_h//2,
              label="価格取得", color=BLUE)

    # MARKETING → X API: "OAuth 1.0a投稿"
    add_arrow(slide,
              Inches(5.9) + sys_w, sys_y + Inches(0.5),
              ext_x, ext_ys[1] + ext_h//2,
              label="OAuth 1.0a", color=GREEN)

    # SBDS → Claude API: "14言語AIサポート"
    add_arrow(slide,
              Inches(5.9) + sys_w, gov_y + Inches(0.3),
              ext_x, ext_ys[0] + ext_h//2,
              label="14言語AI", color=PURPLE)

    # SBDS → ヤマト/佐川: "HMAC-SHA256"
    add_arrow(slide,
              Inches(5.9) + sys_w, gov_y + Inches(0.8),
              ext_x, ext_ys[3] + ext_h//2,
              label="HMAC-SHA256", color=PURPLE)

    # SURPLUS SHIFT → SmartLife IoT: "解錠制御"
    add_arrow(slide,
              Inches(3.1) + sys_w, sys_y + Inches(0.7),
              ext_x, ext_ys[4] + ext_h//2,
              label="解錠制御", color=ACCENT)

    # GOV → Claude API: "レポート生成"
    add_arrow(slide,
              Inches(2.0) + sys_w, gov_y + Inches(0.3),
              ext_x, ext_ys[0] + ext_h,
              label="レポート生成", color=ORANGE)

    # legend
    legend_y = Inches(4.3)
    legend_items = [
        (BLUE,   "RESEARCH"),
        (ACCENT, "SURPLUS SHIFT"),
        (GREEN,  "MARKETING"),
        (ORANGE, "GOV"),
        (PURPLE, "SBDS"),
    ]
    add_text_box(slide, Inches(0.3), legend_y - Inches(0.25), Inches(3), Inches(0.25),
                 "Legend", font_size=9, bold=True, color=GRAY)
    for i, (col, lbl) in enumerate(legend_items):
        lx = Inches(0.3) + i * Inches(1.75)
        add_rect(slide, lx, legend_y, Inches(0.18), Inches(0.18), col)
        add_text_box(slide, lx + Inches(0.22), legend_y, Inches(1.4), Inches(0.2),
                     lbl, font_size=8, color=LIGHT_GRAY)

    add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT)

    # ── Slide 3: Gate制フロー ────────────────────────────────────────────────
    slide = blank_slide(prs)
    set_bg(slide, BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), ACCENT)

    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(10), Inches(0.5),
                 "Gate制フロー  G0 → G4", font_size=18, bold=True, color=WHITE)
    add_text_box(slide, Inches(10), Inches(0.15), Inches(3), Inches(0.4),
                 "NiceEze v14.2", font_size=11, bold=True, color=ACCENT,
                 h_align=PP_ALIGN.RIGHT)

    gates = [
        ("G0", "組織設計・ガバナンス",
         "・会社設計（定款・役員・資本政策）\n・SOC2準拠ガバナンス策定\n・情報管理規程"),
        ("G1", "データ基盤・監査",
         "・DB設計 (PostgreSQL + AES-256)\n・監査ログ設計\n・RLS（行レベルセキュリティ）"),
        ("G2", "各システム実装",
         "・5サブシステム開発\n・API実装・単体テスト\n・bandit セキュリティスキャン"),
        ("G3", "UAT・統合テスト",
         "・E2Eテスト 31件\n・HTTPServer ports 18080-18082\n・負荷テスト・PII検証"),
        ("G4", "本番リリース",
         "・GCP Cloud Run デプロイ\n・Artifact Registry\n・監視ダッシュボード稼働"),
    ]

    gate_colors = [BLUE, GREEN, ORANGE, ACCENT, PURPLE]
    box_w = Inches(2.3)
    box_h = Inches(4.5)
    gap   = Inches(0.25)
    start_x = (SLIDE_W - (5 * box_w + 4 * gap)) / 2
    box_y = Inches(1.1)

    for i, (gid, title, detail) in enumerate(gates):
        gx = start_x + i * (box_w + gap)
        col = gate_colors[i]

        # main card
        add_rect(slide, gx, box_y, box_w, box_h, CARD,
                 line_color=col, line_width=Pt(1.5), rounded=True)

        # gate badge
        add_rect(slide, gx + Inches(0.7), box_y - Inches(0.3),
                 Inches(0.9), Inches(0.6), col,
                 text=gid, font_size=16, bold=True, text_color=WHITE, rounded=True)

        # title
        add_text_box(slide, gx + Inches(0.1), box_y + Inches(0.4),
                     box_w - Inches(0.2), Inches(0.7),
                     title, font_size=11, bold=True, color=col,
                     h_align=PP_ALIGN.CENTER)

        # detail
        add_text_box(slide, gx + Inches(0.12), box_y + Inches(1.2),
                     box_w - Inches(0.24), Inches(3.1),
                     detail, font_size=9, color=LIGHT_GRAY)

        # arrow between gates
        if i < 4:
            ax = gx + box_w + Inches(0.02)
            ay = box_y + box_h / 2
            add_arrow(slide, ax, ay, ax + gap - Inches(0.02), ay,
                      color=GRAY)

    add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT)

    # ── Slide 4: データフロー図 ──────────────────────────────────────────────
    slide = blank_slide(prs)
    set_bg(slide, BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), ACCENT)

    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(10), Inches(0.5),
                 "データフロー図", font_size=18, bold=True, color=WHITE)
    add_text_box(slide, Inches(10), Inches(0.15), Inches(3), Inches(0.4),
                 "NiceEze v14.2", font_size=11, bold=True, color=ACCENT,
                 h_align=PP_ALIGN.RIGHT)

    flows = [
        {
            "label": "① 価格調査フロー",
            "color": BLUE,
            "steps": ["外部API\n(楽天/Yahoo!/Keepa)", "→", "RESEARCH\n分析エンジン", "→", "DB / Cache\n価格トレンド"],
            "y": Inches(1.0),
        },
        {
            "label": "② 旅行者配送フロー",
            "color": PURPLE,
            "steps": ["旅行者\nQRコード", "→", "SBDS\n配送管理", "→", "配送ネットワーク\n(ヤマト/佐川)"],
            "y": Inches(2.7),
        },
        {
            "label": "③ 商談・KPIフロー",
            "color": ACCENT,
            "steps": ["商談データ\n(SURPLUS SHIFT)", "→", "Gate D審査\nhuman_approval", "→", "GOV KPI\nダッシュボード"],
            "y": Inches(4.4),
        },
        {
            "label": "④ レポートフロー",
            "color": ORANGE,
            "steps": ["各サービス\n稼働ログ", "→", "GOV\nCOOレポート", "→", "S10\nダッシュボード"],
            "y": Inches(6.1),
        },
    ]

    node_w = Inches(2.4)
    node_h = Inches(0.9)
    arrow_w = Inches(0.6)
    total_w = 3 * node_w + 2 * arrow_w
    start_x = (SLIDE_W - total_w) / 2

    for flow in flows:
        fy = flow["y"]
        col = flow["color"]
        add_text_box(slide, Inches(0.3), fy - Inches(0.02), Inches(2.5), Inches(0.35),
                     flow["label"], font_size=10, bold=True, color=col)
        steps = flow["steps"]
        nx = start_x
        node_idx = 0
        for s in steps:
            if s == "→":
                # draw arrow in the gap
                ax = nx
                ay = fy + Inches(0.35) + node_h / 2
                add_arrow(slide, ax, ay + Inches(0.38), ax + arrow_w, ay + Inches(0.38),
                          color=col)
                nx += arrow_w
            else:
                add_rect(slide, nx, fy + Inches(0.35), node_w, node_h, CARD,
                         text=s, font_size=9, text_color=WHITE,
                         line_color=col, line_width=Pt(1), rounded=True)
                nx += node_w
                node_idx += 1

    add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT)

    out = "/home/user/niceeze-legal/docs/architecture/NiceEze_全体システム構成図.pptx"
    prs.save(out)
    print(f"Saved: {out}")


# ═══════════════════════════════════════════════════════════════════════════════
# FILE 2: 詳細システム構成図
# ═══════════════════════════════════════════════════════════════════════════════

def detail_slide_header(slide, title, subtitle="NiceEze v14.2", color=BLUE):
    set_bg(slide, BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), ACCENT)
    add_text_box(slide, Inches(0.3), Inches(0.12), Inches(9), Inches(0.5),
                 title, font_size=17, bold=True, color=WHITE)
    add_text_box(slide, Inches(10), Inches(0.12), Inches(3), Inches(0.4),
                 subtitle, font_size=11, bold=True, color=ACCENT,
                 h_align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.3), Inches(0.62), Inches(12.7), Inches(0.03), color)
    add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT)


def component_card(slide, x, y, w, h, title, color, items, title_size=10):
    """Draw a component card with title header and bullet items."""
    # header bar
    add_rect(slide, x, y, w, Inches(0.35), color,
             text=title, font_size=title_size, bold=True, text_color=WHITE)
    # body
    add_rect(slide, x, y + Inches(0.35), w, h - Inches(0.35), CARD,
             line_color=color, line_width=Pt(0.75))
    # items text
    body_text = "\n".join(f"• {it}" for it in items)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.38),
                 w - Inches(0.2), h - Inches(0.45),
                 body_text, font_size=8, color=LIGHT_GRAY)


def build_file2():
    prs = new_prs()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = blank_slide(prs)
    set_bg(slide, BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), ACCENT)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(3), Inches(0.4),
                 "NiceEze", font_size=13, bold=True, color=ACCENT)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.4),
                 "NiceEze 詳細システム構成図",
                 font_size=34, bold=True, color=WHITE, h_align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
                 "自律経営執行システム v14.2  —  各サブシステム内部構成",
                 font_size=18, bold=False, color=BLUE, h_align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(4.5), Inches(4.7), Inches(4.3), Inches(0.05), ACCENT)
    add_text_box(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
                 "RESEARCH  |  SURPLUS SHIFT  |  MARKETING  |  GOV  |  SBDS",
                 font_size=11, color=LIGHT_GRAY, h_align=PP_ALIGN.CENTER)
    add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT)

    # ── Slide 2: RESEARCH システム詳細 ───────────────────────────────────────
    slide = blank_slide(prs)
    detail_slide_header(slide, "RESEARCH システム詳細  (:8080)", color=BLUE)

    # api.py card
    component_card(slide, Inches(0.3), Inches(0.75), Inches(4.5), Inches(3.3),
                   "api.py  —  REST API エントリポイント", BLUE,
                   ["GET  /api/v1/research/price-trend",
                    "GET  /api/v1/research/ranking",
                    "GET  /api/v1/research/growth-alert",
                    "GET  /api/v1/research/new-products",
                    "POST /api/v1/research/export  (Bearer認証)"])

    # analytics.py card
    component_card(slide, Inches(5.1), Inches(0.75), Inches(4.3), Inches(3.3),
                   "analytics.py  —  価格分析エンジン", BLUE,
                   ["ResearchAnalytics クラス",
                    "building_types: luxury / family / student / single",
                    "price_trend()  —  トレンド分析",
                    "ranking()  —  ランキング生成",
                    "growth_alerts()  —  急騰アラート",
                    "export_csv()  —  CSVエクスポート"])

    # res_a01 / res_a02
    component_card(slide, Inches(0.3), Inches(4.2), Inches(2.1), Inches(2.0),
                   "res_a01.py", BLUE,
                   ["楽天API連携",
                    "商品検索・価格取得",
                    "レート制限管理"])

    component_card(slide, Inches(2.6), Inches(4.2), Inches(2.1), Inches(2.0),
                   "res_a02.py", BLUE,
                   ["Yahoo! API連携",
                    "Keepa API連携",
                    "履歴価格取得"])

    # architecture note
    add_rect(slide, Inches(5.1), Inches(4.2), Inches(7.9), Inches(2.0), CARD,
             line_color=DARK_GRAY, line_width=Pt(0.75), rounded=True)
    add_text_box(slide, Inches(5.2), Inches(4.25), Inches(7.7), Inches(1.8),
                 "アーキテクチャノート\n"
                 "• REST API → analytics.py → res_a01/02 → 外部API\n"
                 "• Bearer Token 認証 (export エンドポイント)\n"
                 "• 価格データは DB/Cache に永続化\n"
                 "• RESEARCH → SURPLUS SHIFT / MARKETING へデータ供給",
                 font_size=9, color=LIGHT_GRAY)

    # ── Slide 3: SURPLUS SHIFT システム詳細 ──────────────────────────────────
    slide = blank_slide(prs)
    detail_slide_header(slide, "SURPLUS SHIFT システム詳細  (:8081)", color=ACCENT)

    component_card(slide, Inches(0.3), Inches(0.75), Inches(2.8), Inches(1.5),
                   "surplus_api.py  (:8081)", ACCENT,
                   ["商談APIエントリポイント",
                    "6ステータス管理",
                    "PDF / HTML 出力"])

    component_card(slide, Inches(3.3), Inches(0.75), Inches(5.5), Inches(1.5),
                   "Gate制審査モジュール", ACCENT,
                   ["gate_a.py  —  初期審査",
                    "gate_b.py  —  物件評価",
                    "gate_c.py  —  価格交渉",
                    "gate_d.py  —  最終承認  ★ human_approval_required=True（変更禁止）"])

    component_card(slide, Inches(0.3), Inches(2.5), Inches(4.2), Inches(3.7),
                   "negotiation_log.py  —  自律商談ログ", ACCENT,
                   ["STATUS フロー:",
                    "  DRAFT → HUMAN_APPROVED → SENT",
                    "Gate D 制約:",
                    "  human_approval_required = True",
                    "  （変更禁止・監査対象）",
                    "商談履歴の完全監査証跡",
                    "PDF / HTML 出力サポート"])

    component_card(slide, Inches(4.7), Inches(2.5), Inches(8.3), Inches(3.7),
                   "negotiation_api.py  —  商談API", ACCENT,
                   ["6ステータス管理:",
                    "  NEW / DRAFT / HUMAN_APPROVED / SENT / ACCEPTED / CLOSED",
                    "SmartLifeProduct クラス:",
                    "  __setattr__ ガード（不変フィールド保護）",
                    "  SmartLife IoT 解錠制御連携",
                    "PDF生成 / HTML出力",
                    "GOV へ 商談KPI 供給"])

    # ── Slide 4: MARKETING システム詳細 ──────────────────────────────────────
    slide = blank_slide(prs)
    detail_slide_header(slide, "MARKETING システム詳細", color=GREEN)

    component_card(slide, Inches(0.3), Inches(0.75), Inches(4.2), Inches(2.8),
                   "x_poster.py  —  X(Twitter) 投稿", GREEN,
                   ["OAuth 1.0a 署名（HMAC-SHA1）",
                    "140文字トランケート",
                    "mock_mode（環境変数未設定時）",
                    "投稿スケジューリング対応",
                    "エラーリトライロジック"])

    component_card(slide, Inches(4.7), Inches(0.75), Inches(4.0), Inches(2.8),
                   "content_generator.py  —  コンテンツ生成", GREEN,
                   ["Claude API 連携",
                    "不動産・旅行カテゴリ対応",
                    "RESEARCH トレンド情報活用",
                    "ハッシュタグ自動付与"])

    component_card(slide, Inches(0.3), Inches(3.75), Inches(2.6), Inches(2.5),
                   "news_crawler.py", GREEN,
                   ["ニュース収集",
                    "RSSフィード解析",
                    "キーワードフィルタ"])

    component_card(slide, Inches(3.1), Inches(3.75), Inches(2.6), Inches(2.5),
                   "scheduler.py", GREEN,
                   ["定期実行スケジューラ",
                    "cron 設定",
                    "失敗時アラート"])

    component_card(slide, Inches(5.9), Inches(3.75), Inches(3.0), Inches(2.5),
                   "delivery_log.py", GREEN,
                   ["配信ログ管理",
                    "投稿実績 → GOV",
                    "エンゲージメント記録"])

    # flow note
    add_rect(slide, Inches(9.1), Inches(0.75), Inches(4.0), Inches(5.5), CARD,
             line_color=DARK_GRAY, line_width=Pt(0.75), rounded=True)
    add_text_box(slide, Inches(9.2), Inches(0.8), Inches(3.8), Inches(5.3),
                 "投稿フロー\n\n"
                 "1. news_crawler → トレンド取得\n"
                 "2. content_generator → 文章生成\n"
                 "3. x_poster → OAuth 1.0a 署名\n"
                 "4. X API へ POST\n"
                 "5. delivery_log → 実績記録\n"
                 "6. GOV へ 投稿実績 供給\n\n"
                 "mock_mode:\n"
                 "  環境変数 X_API_KEY 未設定時\n"
                 "  実際の投稿なしでテスト可能",
                 font_size=9, color=LIGHT_GRAY)

    # ── Slide 5: GOV システム詳細 ─────────────────────────────────────────────
    slide = blank_slide(prs)
    detail_slide_header(slide, "GOV システム詳細  (:8082)", color=ORANGE)

    component_card(slide, Inches(0.3), Inches(0.75), Inches(5.0), Inches(2.8),
                   "api.py  (:8082)  —  エンドポイント一覧", ORANGE,
                   ["GET/POST  /coo/report   —  COOレポート",
                    "GET/POST  /coo/kpi      —  KPI管理",
                    "GET/POST  /coo/budget   —  予算管理",
                    "GET/POST  /coo/pmo      —  PMO(Gates G0-G4)",
                    "GET  /finops/summary    —  コスト集計",
                    "GET  /finops/alerts     —  コストアラート",
                    "POST /finops/cost       —  コスト登録",
                    "POST /ops/log           —  稼働ログ登録",
                    "GET  /ops/health        —  ヘルスチェック",
                    "GET  /ops/logs/{service}  —  サービス別ログ"])

    component_card(slide, Inches(5.5), Inches(0.75), Inches(4.0), Inches(2.8),
                   "s10_coo_report.py  —  COOレポートエンジン", ORANGE,
                   ["KPIRecord:",
                    "  achievement_rate（達成率）",
                    "BudgetRecord:",
                    "  variance_jpy（予算差異）",
                    "PMOTask:",
                    "  Gates G0 → G4 追跡",
                    "Claude API でレポート文章生成",
                    "S10ダッシュボード向け出力"])

    component_card(slide, Inches(0.3), Inches(3.75), Inches(4.2), Inches(2.5),
                   "finops_monitor.py  —  コスト監視", ORANGE,
                   ["MAX_COST_PER_DELIVERY: ¥0.5",
                    "MONTHLY_BUDGET:        ¥5,000",
                    "ALERT_THRESHOLD:       80%",
                    "超過時アラート → 管理者通知",
                    "Claude API コスト追跡"])

    component_card(slide, Inches(4.7), Inches(3.75), Inches(4.0), Inches(2.5),
                   "ops_log_collector.py  —  稼働ログ", ORANGE,
                   ["監視対象サービス:",
                    "  sbds / surplus_shift / research",
                    "  marketing / gov",
                    "ヘルスチェック集約",
                    "障害検知・アラート"])

    add_rect(slide, Inches(8.9), Inches(3.75), Inches(4.1), Inches(2.5), CARD,
             line_color=DARK_GRAY, line_width=Pt(0.75), rounded=True)
    add_text_box(slide, Inches(9.0), Inches(3.8), Inches(3.9), Inches(2.3),
                 "集約フロー\n"
                 "RESEARCH / SURPLUS SHIFT\n"
                 "MARKETING / SBDS\n"
                 "    ↓\n"
                 "GOV (ops_log_collector)\n"
                 "    ↓\n"
                 "COOレポート + S10ダッシュボード",
                 font_size=9, color=LIGHT_GRAY)

    # ── Slide 6: SBDS システム詳細 ────────────────────────────────────────────
    slide = blank_slide(prs)
    detail_slide_header(slide, "SBDS システム詳細  (:8083)", color=PURPLE)

    component_card(slide, Inches(0.3), Inches(0.75), Inches(2.5), Inches(1.6),
                   "travel_api.py  (:8083)", PURPLE,
                   ["REST APIエントリポイント",
                    "旅行者セッション管理",
                    "QRコード発行"])

    component_card(slide, Inches(3.0), Inches(0.75), Inches(3.3), Inches(1.6),
                   "travel_qr.py  —  QR発行管理", PURPLE,
                   ["TTL: 86400秒（24時間）",
                    "secrets.token_urlsafe(32)",
                    "QR失効・再発行ロジック"])

    component_card(slide, Inches(6.5), Inches(0.75), Inches(6.5), Inches(1.6),
                   "carrier_webhook.py  —  キャリアWebhook", PURPLE,
                   ["ヤマト / 佐川  HMAC-SHA256 署名検証",
                    "PII: tracking_number を SHA-256 ハッシュ化",
                    "auto_assign_locker()  —  ロッカー自動割当",
                    "Webhook受信 → 配送状態更新"])

    component_card(slide, Inches(0.3), Inches(2.55), Inches(4.5), Inches(3.7),
                   "ai_support.py  —  AIサポートセンター", PURPLE,
                   ["14言語対応:",
                    "  ja / en / zh-CN / zh-TW / ko / th",
                    "  fr / it / de / es / pt / id / ar / hi",
                    "Claude API 連携:",
                    "  model: claude-sonnet-4-20250514",
                    "fallback: テンプレートベース",
                    "言語自動検出",
                    "FAQ自動応答"])

    component_card(slide, Inches(5.0), Inches(2.55), Inches(3.0), Inches(1.7),
                   "tms_set_001.py", PURPLE,
                   ["配送セット管理",
                    "荷物グルーピング",
                    "ルート最適化"])

    component_card(slide, Inches(8.2), Inches(2.55), Inches(2.8), Inches(1.7),
                   "tms_drv_001.py", PURPLE,
                   ["ドライバー管理",
                    "配送状態追跡",
                    "完了通知"])

    component_card(slide, Inches(5.0), Inches(4.45), Inches(2.6), Inches(1.85),
                   "travel_pdf.py", PURPLE,
                   ["PDF生成",
                    "旅行者向け",
                    "QR埋め込み"])

    component_card(slide, Inches(7.8), Inches(4.45), Inches(5.2), Inches(1.85),
                   "static/panel.html  —  タッチパネルUI", PURPLE,
                   ["14言語選択画面",
                    "タッチ操作対応",
                    "QRスキャン連携",
                    "多言語AIサポート起動"])

    # ── Slide 7: セキュリティ・インフラ構成 ──────────────────────────────────
    slide = blank_slide(prs)
    detail_slide_header(slide, "セキュリティ・インフラ構成", color=ACCENT)

    # Security section
    add_text_box(slide, Inches(0.3), Inches(0.72), Inches(6), Inches(0.35),
                 "セキュリティレイヤー", font_size=13, bold=True, color=ACCENT)

    sec_items = [
        ("HMAC-SHA256",     "キャリア Webhook 署名検証（ヤマト/佐川）",      BLUE),
        ("OAuth 1.0a",      "X (Twitter) API 認証（HMAC-SHA1署名）",         GREEN),
        ("Bearer Token",    "RESEARCH export エンドポイント認証",             BLUE),
        ("SHA-256 PII Hash","tracking_number のハッシュ化（個人情報保護）",  PURPLE),
        ("Gate D 制約",     "human_approval_required=True（変更禁止）",       ACCENT),
        ("bandit スキャン", "severity High / Medium: 0件",                    GREEN),
    ]

    for i, (name, desc, col) in enumerate(sec_items):
        iy = Inches(1.15) + i * Inches(0.72)
        add_rect(slide, Inches(0.3), iy, Inches(2.1), Inches(0.55), col,
                 text=name, font_size=9, bold=True, text_color=WHITE, rounded=True)
        add_rect(slide, Inches(2.55), iy, Inches(6.1), Inches(0.55), CARD,
                 text=desc, font_size=9, text_color=LIGHT_GRAY,
                 line_color=DARK_GRAY, line_width=Pt(0.5), rounded=True)

    # Infrastructure section
    add_text_box(slide, Inches(9.0), Inches(0.72), Inches(4), Inches(0.35),
                 "インフラ構成", font_size=13, bold=True, color=ORANGE)

    infra_items = [
        ("GCP Cloud Run",        "3サービス デプロイ",             ORANGE),
        ("Artifact Registry",    "コンテナイメージ管理",            ORANGE),
        ("Cloud SQL",            "PostgreSQL + AES-256 + RLS",     ORANGE),
        ("E2Eテスト",            "31テスト / HTTPServer ports\n18080 / 18081 / 18082",
                                                                    GREEN),
    ]

    for i, (name, desc, col) in enumerate(infra_items):
        iy = Inches(1.15) + i * Inches(1.05)
        add_rect(slide, Inches(9.0), iy, Inches(1.9), Inches(0.9), col,
                 text=name, font_size=9, bold=True, text_color=WHITE, rounded=True)
        add_rect(slide, Inches(11.1), iy, Inches(2.0), Inches(0.9), CARD,
                 text=desc, font_size=8, text_color=LIGHT_GRAY,
                 line_color=DARK_GRAY, line_width=Pt(0.5), rounded=True)

    # compliance summary box
    add_rect(slide, Inches(0.3), Inches(5.55), Inches(12.7), Inches(1.2), CARD,
             line_color=ACCENT, line_width=Pt(1.0), rounded=True)
    add_text_box(slide, Inches(0.5), Inches(5.62), Inches(12.3), Inches(1.05),
                 "コンプライアンスサマリー\n"
                 "SOC2準拠ガバナンス  |  RLS（行レベルセキュリティ）  |  AES-256 保存時暗号化  "
                 "|  PII SHA-256ハッシュ化  |  Gate D human_approval  |  bandit 0件  "
                 "|  E2E 31テスト全パス",
                 font_size=9, color=LIGHT_GRAY)

    out = "/home/user/niceeze-legal/docs/architecture/NiceEze_詳細システム構成図.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build_file1()
    build_file2()
    print("Done.")
