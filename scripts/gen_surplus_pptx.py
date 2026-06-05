"""SURPLUS_SHIFT SEQ / UI / PHASE .pptx 生成スクリプト"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_NAVY = RGBColor(0x1a, 0x3a, 0x5c)
ORANGE    = RGBColor(0xf5, 0xa6, 0x23)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
LIGHT_BLUE = RGBColor(0xBF, 0xD7, 0xED)
GREEN     = RGBColor(0x05, 0x96, 0x69)
RED       = RGBColor(0xEF, 0x44, 0x44)

A4_W = Inches(8.27)
A4_H = Inches(11.69)
A4_W_LAND = Inches(11.69)
A4_H_LAND = Inches(8.27)


# ─────────────────────────────────────────────────────────────────
# ユーティリティ
# ─────────────────────────────────────────────────────────────────

def new_prs_portrait():
    prs = Presentation()
    prs.slide_width  = A4_W
    prs.slide_height = A4_H
    return prs


def new_prs_landscape():
    prs = Presentation()
    prs.slide_width  = A4_W_LAND
    prs.slide_height = A4_H_LAND
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=Pt(9), bold=False, color=None,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color or DARK_NAVY
    return txBox


def title_bar(slide, title, subtitle='', w=None, is_landscape=False):
    sw = w or (A4_W_LAND if is_landscape else A4_W)
    add_rect(slide, 0, 0, sw, Inches(0.6), fill_rgb=DARK_NAVY)
    add_text_box(slide, title, Inches(0.15), Inches(0.05), sw - Inches(0.3), Inches(0.35),
                 font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.15), Inches(0.38), sw - Inches(0.3), Inches(0.2),
                     font_size=Pt(8), color=LIGHT_BG, align=PP_ALIGN.LEFT)


def footer_bar(slide, text, is_landscape=False):
    sw = A4_W_LAND if is_landscape else A4_W
    sh = A4_H_LAND if is_landscape else A4_H
    add_rect(slide, 0, sh - Inches(0.25), sw, Inches(0.25), fill_rgb=DARK_NAVY)
    add_text_box(slide, text, Inches(0.1), sh - Inches(0.23), sw - Inches(0.2), Inches(0.2),
                 font_size=Pt(7), color=LIGHT_BG, align=PP_ALIGN.CENTER)


def slide_number_tag(slide, num, total, is_landscape=False):
    sw = A4_W_LAND if is_landscape else A4_W
    sh = A4_H_LAND if is_landscape else A4_H
    add_text_box(slide, f'{num} / {total}',
                 sw - Inches(0.8), sh - Inches(0.23), Inches(0.7), Inches(0.2),
                 font_size=Pt(7), color=LIGHT_BG, align=PP_ALIGN.RIGHT)


def draw_seq_flow(slide, actors, steps, y_start=Inches(0.75), is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W
    margin = Inches(0.3)
    actor_w = Inches(1.4)
    n = len(actors)
    spacing = (sw - 2 * margin - actor_w) / max(n - 1, 1)

    actor_x = [margin + i * spacing for i in range(n)]
    actor_y = y_start

    for i, actor in enumerate(actors):
        add_rect(slide, actor_x[i], actor_y, actor_w, Inches(0.35),
                 fill_rgb=DARK_NAVY, line_rgb=ORANGE)
        add_text_box(slide, actor,
                     actor_x[i], actor_y, actor_w, Inches(0.35),
                     font_size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    line_y_start = actor_y + Inches(0.35)
    line_y_end = actor_y + Inches(0.35) + Inches(0.45) * (len(steps) + 1)
    cx = [ax + actor_w / 2 for ax in actor_x]

    for cx_i in cx:
        line_shape = slide.shapes.add_connector(1, cx_i, line_y_start, cx_i, line_y_end)
        line_shape.line.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
        line_shape.line.width = Pt(0.5)

    step_y = line_y_start + Inches(0.3)
    for step in steps:
        from_idx, to_idx, label, is_async, note = step
        fx = cx[from_idx]
        tx = cx[to_idx]
        arrow = slide.shapes.add_connector(2, fx, step_y, tx, step_y)
        arrow.line.color.rgb = ORANGE if not is_async else RGBColor(0x6b, 0x72, 0x80)
        arrow.line.width = Pt(1.5 if not is_async else 1.0)
        mid_x = min(fx, tx) + abs(tx - fx) / 2 - Inches(0.6)
        add_text_box(slide, label, mid_x, step_y - Inches(0.2), Inches(1.2), Inches(0.18),
                     font_size=Pt(7), color=DARK_NAVY, align=PP_ALIGN.CENTER)
        if note:
            note_x = max(fx, tx) + Inches(0.05)
            add_text_box(slide, note, note_x, step_y - Inches(0.15), Inches(1.5), Inches(0.2),
                         font_size=Pt(6.5), color=RED, align=PP_ALIGN.LEFT)
        step_y += Inches(0.42)


# ─────────────────────────────────────────────────────────────────
# SEQ: シーケンス図集
# ─────────────────────────────────────────────────────────────────

def gen_seq(output='docs/SURPLUS_SHIFT/SURPLUS_SHIFT_SEQ_v1.0.pptx'):
    prs = new_prs_landscape()
    FOOTER = 'SURPLUS_SHIFT  |  シーケンス図 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # ─ Slide 1: タイトル ─
    sl = blank_slide(prs)
    title_bar(sl, 'SURPLUS_SHIFT シーケンス図 v1.0', '余剰在庫転換システム — シーケンス図・フロー', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_rect(sl, Inches(1.5), Inches(1.5), Inches(8.5), Inches(3.5), fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
    add_text_box(sl, 'SURPLUS_SHIFT', Inches(1.7), Inches(1.7), Inches(8.0), Inches(0.8),
                 font_size=Pt(32), bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    add_text_box(sl, '余剰在庫転換システム', Inches(1.7), Inches(2.6), Inches(8.0), Inches(0.5),
                 font_size=Pt(18), bold=False, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'シーケンス図 v1.0  |  2026-06-05', Inches(1.7), Inches(3.2), Inches(8.0), Inches(0.4),
                 font_size=Pt(12), color=RGBColor(0x47, 0x55, 0x69), align=PP_ALIGN.CENTER)

    items = [
        ('Slide 2', '売り手フロー', '在庫登録→自動分析→余剰検知→ダッシュボード表示→マッチング待機'),
        ('Slide 3', '買い手フロー', '条件設定→マッチング候補受信→交渉→契約書生成→決済'),
        ('Slide 4', 'システムシーケンス', 'Cloud Functions→Firestore→Cloud Run→LINE通知'),
        ('Slide 5', 'エラーハンドリング', 'タイムアウト/マッチング失敗/決済エラー 各フォールバック'),
    ]
    for i, (slide_no, title, desc) in enumerate(items):
        y = Inches(5.2) + i * Inches(0.5)
        add_text_box(sl, f'{slide_no}: {title} — {desc}',
                     Inches(0.5), y, Inches(10.5), Inches(0.4),
                     font_size=Pt(9), color=DARK_NAVY)
    slide_number_tag(sl, 1, TOTAL, is_landscape=True)

    # ─ Slide 2: 売り手フロー ─
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-01  売り手フロー', '在庫登録 → 自動分析 → 余剰検知 → ダッシュボード表示 → マッチング待機', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    actors = ['売り手', 'SURPLUS-001\nダッシュボード', 'Cloud Run\nAPI', 'Firestore\nDB', 'Cloud Functions\n定期スキャン']
    steps = [
        (0, 1, '在庫データ登録', False, ''),
        (1, 2, 'POST /inventory', False, ''),
        (2, 3, 'write inventory', False, ''),
        (3, 2, 'ack', False, ''),
        (4, 2, '定期スキャン起動', True, 'cron毎時'),
        (2, 3, '余剰判定クエリ', False, '滞留30日以上'),
        (3, 2, 'return surplus[]', False, ''),
        (2, 1, '余剰検知通知', False, 'SSE push'),
        (1, 0, 'ダッシュボード表示', False, '転換率グラフ更新'),
        (0, 1, 'マッチング待機設定', False, ''),
        (1, 2, 'PUT /matching/wait', True, ''),
        (2, 3, 'update status=WAITING', False, ''),
    ]
    draw_seq_flow(sl, actors, steps, is_landscape=True)
    slide_number_tag(sl, 2, TOTAL, is_landscape=True)

    # ─ Slide 3: 買い手フロー ─
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-02  買い手フロー', '条件設定 → マッチング候補受信 → 交渉 → 契約書生成 → 決済', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    actors_buy = ['買い手', 'SURPLUS-002\n検索・マッチング', 'Cloud Run\nマッチング', 'SURPLUS-003\n価格交渉', '決済\nゲートウェイ']
    steps_buy = [
        (0, 1, '条件フィルタ設定', False, ''),
        (1, 2, 'POST /match/search', False, ''),
        (2, 1, 'マッチング候補返却', False, 'スコア付き'),
        (1, 0, 'マッチングカード表示', False, ''),
        (0, 1, '候補選択', False, ''),
        (1, 3, '交渉画面遷移', False, ''),
        (0, 3, '価格提示', False, ''),
        (3, 2, 'POST /negotiate', True, 'AI推奨価格'),
        (2, 3, '推奨価格返却', False, ''),
        (3, 0, '価格履歴・提案表示', False, ''),
        (0, 3, '交渉合意', False, ''),
        (3, 2, '契約書生成リクエスト', False, ''),
        (2, 4, 'POST /payment', False, ''),
        (4, 0, '決済完了通知', True, 'LINE通知'),
    ]
    draw_seq_flow(sl, actors_buy, steps_buy, is_landscape=True)
    slide_number_tag(sl, 3, TOTAL, is_landscape=True)

    # ─ Slide 4: システムシーケンス ─
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-03  システムシーケンス', 'Cloud Functions(定期スキャン)→Firestore(余剰判定)→Cloud Run(マッチング)→LINE通知', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    actors_sys = ['Cloud Functions\n定期スキャン', 'Firestore\nDB', 'Cloud Run\nマッチング', 'Claude API\nAI判定', 'LINE API\n通知']
    steps_sys = [
        (0, 0, 'cron起動（毎時）', False, 'Cloud Scheduler'),
        (0, 1, 'query: 滞留在庫検索', False, '30日以上未移動'),
        (1, 0, 'return surplus_items[]', False, ''),
        (0, 2, 'POST /match surplus_items', False, ''),
        (2, 3, '余剰度スコアリング', True, 'Claude API'),
        (3, 2, 'score: 0.0-1.0', False, ''),
        (2, 1, 'write match_candidates', False, ''),
        (1, 2, 'ack', False, ''),
        (2, 4, 'LINE Push to buyers', True, 'スコア上位'),
        (4, 2, 'delivery_receipt', False, ''),
        (2, 1, 'update notification_sent', False, ''),
    ]
    draw_seq_flow(sl, actors_sys, steps_sys, is_landscape=True)
    slide_number_tag(sl, 4, TOTAL, is_landscape=True)

    # ─ Slide 5: エラーハンドリング ─
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-04  エラーハンドリング', 'タイムアウト / マッチング失敗 / 決済エラー 各フォールバック', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, '【エラーパターン別フォールバック一覧】',
                 Inches(0.4), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)

    errors = [
        ('タイムアウト（API 30秒超過）',
         'Cloud Run → Firestore タイムアウト',
         'リトライ3回（exponential backoff 1s/2s/4s）→ 失敗時は503返却・アラート送信'),
        ('マッチング失敗（候補0件）',
         'マッチングエンジン → 売り手',
         '条件緩和サジェスト表示 / 30日後再マッチング予約 / 手動マッチング依頼フォーム'),
        ('決済エラー（ゲートウェイ拒否）',
         '決済ゲートウェイ → 買い手',
         '別決済手段提案（振込/後払い） / 交渉状態に戻す / エラーコードをFirestoreに記録'),
        ('LINE通知失敗（Rate Limit 429）',
         'LINE API → Cloud Run',
         'キューに積んで60秒後リトライ / 3回失敗でメール代替通知 / audit_logに記録'),
    ]

    for i, (title, actors_str, fallback) in enumerate(errors):
        y = Inches(1.05) + i * Inches(1.5)
        add_rect(sl, Inches(0.3), y, Inches(10.8), Inches(1.35),
                 fill_rgb=LIGHT_BG, line_rgb=RED)
        add_rect(sl, Inches(0.3), y, Inches(10.8), Inches(0.35),
                 fill_rgb=RGBColor(0x7f, 0x1d, 0x1d), line_rgb=RED)
        add_text_box(sl, f'ERROR: {title}',
                     Inches(0.4), y + Inches(0.04), Inches(10.0), Inches(0.25),
                     font_size=Pt(9), bold=True, color=RGBColor(0xfc, 0xa5, 0xa5))
        add_text_box(sl, f'発生箇所: {actors_str}',
                     Inches(0.4), y + Inches(0.4), Inches(10.0), Inches(0.25),
                     font_size=Pt(8), color=RGBColor(0x47, 0x55, 0x69))
        add_text_box(sl, f'フォールバック: {fallback}',
                     Inches(0.4), y + Inches(0.68), Inches(10.4), Inches(0.55),
                     font_size=Pt(8.5), color=DARK_NAVY, word_wrap=True)

    slide_number_tag(sl, 5, TOTAL, is_landscape=True)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ─────────────────────────────────────────────────────────────────
# UI: 画面設計書
# ─────────────────────────────────────────────────────────────────

def draw_wireframe_box(slide, x, y, w, h, label, fill=LIGHT_BG, border=DARK_NAVY, label_size=Pt(8)):
    add_rect(slide, x, y, w, h, fill_rgb=fill, line_rgb=border)
    add_text_box(slide, label, x + Inches(0.05), y + Inches(0.05),
                 w - Inches(0.1), h - Inches(0.1),
                 font_size=label_size, color=DARK_NAVY, word_wrap=True)


def gen_ui(output='docs/SURPLUS_SHIFT/SURPLUS_SHIFT_UI_v1.0.pptx'):
    prs = new_prs_portrait()
    FOOTER = 'SURPLUS_SHIFT  |  UI設計 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # ─ Slide 1: タイトル ─
    sl = blank_slide(prs)
    title_bar(sl, 'SURPLUS_SHIFT UI設計 v1.0', '余剰在庫転換システム — 画面設計', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_rect(sl, Inches(0.4), Inches(1.0), Inches(7.3), Inches(2.5), fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
    add_text_box(sl, 'SURPLUS_SHIFT UI設計書', Inches(0.6), Inches(1.2), Inches(6.9), Inches(0.7),
                 font_size=Pt(24), bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    add_text_box(sl, '余剰在庫転換システム', Inches(0.6), Inches(2.0), Inches(6.9), Inches(0.4),
                 font_size=Pt(14), color=ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'v1.0  |  2026-06-05  |  © 2026 株式会社NiceEze',
                 Inches(0.6), Inches(2.5), Inches(6.9), Inches(0.3),
                 font_size=Pt(10), color=RGBColor(0x47, 0x55, 0x69), align=PP_ALIGN.CENTER)

    screens = [
        ('SURPLUS-001', '売り手ダッシュボード', '余剰在庫リスト / 転換率グラフ / クイックアクション'),
        ('SURPLUS-002', '買い手検索・マッチング', '条件フィルタ / マッチングカード / スコア表示'),
        ('SURPLUS-003', '価格交渉支援', 'チャットUI / 価格履歴 / 推奨価格AI提案'),
        ('LAYOUT', 'LAYOUT_MASTER準拠指定', 'font-mono tabular-nums / カラーパレット / IndexedDB v142'),
    ]
    for i, (sid, name, desc) in enumerate(screens):
        y = Inches(3.7) + i * Inches(1.0)
        add_rect(sl, Inches(0.4), y, Inches(7.3), Inches(0.85), fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
        add_rect(sl, Inches(0.4), y, Inches(1.3), Inches(0.85), fill_rgb=DARK_NAVY, line_rgb=DARK_NAVY)
        add_text_box(sl, sid, Inches(0.45), y + Inches(0.2), Inches(1.2), Inches(0.4),
                     font_size=Pt(8), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        add_text_box(sl, name, Inches(1.8), y + Inches(0.05), Inches(3.0), Inches(0.3),
                     font_size=Pt(10), bold=True, color=DARK_NAVY)
        add_text_box(sl, desc, Inches(1.8), y + Inches(0.4), Inches(5.7), Inches(0.3),
                     font_size=Pt(8), color=RGBColor(0x47, 0x55, 0x69))
    slide_number_tag(sl, 1, TOTAL)

    # ─ Slide 2: 売り手ダッシュボード SURPLUS-001 ─
    sl = blank_slide(prs)
    title_bar(sl, 'SURPLUS-001  売り手ダッシュボード', '対象ユーザー: 売り手  |  余剰在庫リスト / 転換率グラフ / クイックアクション', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'LAYOUT_MASTER.md 準拠 3ゾーン構成', Inches(0.2), Inches(0.65), Inches(7.5), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)

    # 上部: KPIサマリー
    add_rect(sl, Inches(0.2), Inches(0.9), Inches(7.6), Inches(1.5),
             fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=ORANGE)
    add_text_box(sl, '上部: KPIサマリー', Inches(0.3), Inches(0.92), Inches(4.0), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)
    kpis = [
        ('余剰在庫数', '142件', Inches(0.3)),
        ('転換率', '68.3%', Inches(2.9)),
        ('今月転換額', '¥2.4M', Inches(5.5)),
    ]
    for label, val, sx in kpis:
        add_text_box(sl, label, sx, Inches(1.15), Inches(2.2), Inches(0.2),
                     font_size=Pt(7), color=RGBColor(0x64, 0x74, 0x8b))
        add_text_box(sl, val, sx, Inches(1.35), Inches(2.2), Inches(0.4),
                     font_size=Pt(16), bold=True, color=WHITE)
    add_text_box(sl, 'font-mono tabular-nums tracking-tight（全KPI数値）',
                 Inches(0.3), Inches(1.75), Inches(7.0), Inches(0.18),
                 font_size=Pt(6.5), color=ORANGE)

    # 中央: 余剰在庫リスト + 転換率グラフ
    add_rect(sl, Inches(0.2), Inches(2.45), Inches(4.5), Inches(5.5),
             fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=ORANGE)
    add_text_box(sl, '余剰在庫リスト', Inches(0.3), Inches(2.47), Inches(4.0), Inches(0.22),
                 font_size=Pt(8), bold=True, color=ORANGE)

    headers = ['商品名', '在庫数', '滞留日数', '推奨アクション']
    col_widths = [Inches(1.5), Inches(0.7), Inches(0.8), Inches(1.3)]
    x_offset = Inches(0.25)
    for i, (h, cw) in enumerate(zip(headers, col_widths)):
        add_rect(sl, x_offset, Inches(2.72), cw - Inches(0.02), Inches(0.28),
                 fill_rgb=DARK_NAVY, line_rgb=RGBColor(0x33, 0x41, 0x55))
        add_text_box(sl, h, x_offset + Inches(0.02), Inches(2.74), cw - Inches(0.06), Inches(0.24),
                     font_size=Pt(7), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x_offset += cw

    samples = [
        ['電動ドリルA型', '85個', '45日', 'マッチング推奨'],
        ['工業用テープ', '320m', '62日', '値下げ検討'],
        ['金属部品B', '1,200個', '38日', 'マッチング推奨'],
        ['工具セットC', '15セット', '90日', '処分検討'],
    ]
    for ri, row in enumerate(samples):
        x_offset = Inches(0.25)
        for ci, (val, cw) in enumerate(zip(row, col_widths)):
            bg = RGBColor(0x1e, 0x29, 0x3b) if ri % 2 == 0 else RGBColor(0x0f, 0x17, 0x2a)
            add_rect(sl, x_offset, Inches(3.02) + ri * Inches(0.35),
                     cw - Inches(0.02), Inches(0.32), fill_rgb=bg,
                     line_rgb=RGBColor(0x1e, 0x29, 0x3b))
            color = ORANGE if val == 'マッチング推奨' else (RED if val == '処分検討' else RGBColor(0xe2, 0xe8, 0xf0))
            add_text_box(sl, val, x_offset + Inches(0.02), Inches(3.04) + ri * Inches(0.35),
                         cw - Inches(0.06), Inches(0.28),
                         font_size=Pt(7.5), color=color)
            x_offset += cw

    # 転換率グラフ（簡易）
    add_rect(sl, Inches(4.8), Inches(2.45), Inches(3.0), Inches(5.5),
             fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=ORANGE)
    add_text_box(sl, '転換率グラフ（月次）', Inches(4.9), Inches(2.47), Inches(2.8), Inches(0.22),
                 font_size=Pt(8), bold=True, color=ORANGE)
    months = ['1月', '2月', '3月', '4月', '5月', '6月']
    rates = [0.45, 0.52, 0.58, 0.63, 0.65, 0.68]
    bar_w = Inches(0.35)
    for mi, (month, rate) in enumerate(zip(months, rates)):
        bx = Inches(4.95) + mi * Inches(0.42)
        bar_h = Inches(3.0) * rate
        by = Inches(5.7) - bar_h
        add_rect(sl, bx, by, bar_w, bar_h, fill_rgb=ORANGE, line_rgb=ORANGE)
        add_text_box(sl, month, bx - Inches(0.02), Inches(5.75), bar_w + Inches(0.04), Inches(0.2),
                     font_size=Pt(6.5), color=RGBColor(0x94, 0xa3, 0xb8), align=PP_ALIGN.CENTER)

    # クイックアクション
    add_rect(sl, Inches(0.2), Inches(8.0), Inches(7.6), Inches(0.8),
             fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=ORANGE)
    add_text_box(sl, 'クイックアクション', Inches(0.3), Inches(8.02), Inches(3.0), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)
    actions = [('一括マッチング', GREEN, Inches(0.3)), ('CSV出力', DARK_NAVY, Inches(2.1)), ('条件設定', ORANGE, Inches(3.9))]
    for label, color, ax in actions:
        add_rect(sl, ax, Inches(8.28), Inches(1.5), Inches(0.4), fill_rgb=color, line_rgb=color)
        add_text_box(sl, label, ax, Inches(8.33), Inches(1.5), Inches(0.3),
                     font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number_tag(sl, 2, TOTAL)

    # ─ Slide 3: 買い手検索・マッチング SURPLUS-002 ─
    sl = blank_slide(prs)
    title_bar(sl, 'SURPLUS-002  買い手検索・マッチング', '対象ユーザー: 買い手  |  条件フィルタ / マッチングカード / スコア表示', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'LAYOUT_MASTER.md 準拠', Inches(0.2), Inches(0.65), Inches(7.5), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)

    # 条件フィルタ（左ペイン）
    add_rect(sl, Inches(0.2), Inches(0.9), Inches(2.3), Inches(8.5),
             fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=ORANGE)
    add_text_box(sl, '条件フィルタ', Inches(0.3), Inches(0.92), Inches(2.1), Inches(0.22),
                 font_size=Pt(8), bold=True, color=ORANGE)
    filters = ['カテゴリ', '価格帯', '数量', '地域', '納期']
    for i, f in enumerate(filters):
        fy = Inches(1.2) + i * Inches(0.65)
        add_rect(sl, Inches(0.3), fy, Inches(2.0), Inches(0.5),
                 fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=RGBColor(0x33, 0x41, 0x55))
        add_text_box(sl, f, Inches(0.35), fy + Inches(0.05), Inches(1.9), Inches(0.18),
                     font_size=Pt(7), color=RGBColor(0x64, 0x74, 0x8b))
        add_text_box(sl, '▼ 選択してください', Inches(0.35), fy + Inches(0.25), Inches(1.9), Inches(0.18),
                     font_size=Pt(7), color=RGBColor(0xe2, 0xe8, 0xf0))

    # マッチングカード（右エリア）
    add_rect(sl, Inches(2.6), Inches(0.9), Inches(5.4), Inches(8.5),
             fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=ORANGE)
    add_text_box(sl, 'マッチング候補 (3件)', Inches(2.7), Inches(0.92), Inches(5.0), Inches(0.22),
                 font_size=Pt(8), bold=True, color=ORANGE)

    match_items = [
        ('電動ドリルA型', '85個', '¥1,250/個', '0.92', GREEN),
        ('工業用テープ', '320m', '¥480/m', '0.78', ORANGE),
        ('金属部品B', '1,200個', '¥85/個', '0.71', ORANGE),
    ]
    for i, (name, qty, price, score, score_color) in enumerate(match_items):
        cy = Inches(1.2) + i * Inches(2.4)
        add_rect(sl, Inches(2.7), cy, Inches(5.1), Inches(2.1),
                 fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=DARK_NAVY)
        add_text_box(sl, name, Inches(2.8), cy + Inches(0.08), Inches(3.5), Inches(0.3),
                     font_size=Pt(11), bold=True, color=WHITE)
        add_text_box(sl, f'数量: {qty}  |  単価: {price}',
                     Inches(2.8), cy + Inches(0.45), Inches(3.5), Inches(0.25),
                     font_size=Pt(8.5), color=RGBColor(0xe2, 0xe8, 0xf0))
        # スコア表示
        add_rect(sl, Inches(6.5), cy + Inches(0.08), Inches(1.1), Inches(0.55),
                 fill_rgb=score_color, line_rgb=score_color)
        add_text_box(sl, f'スコア\n{score}', Inches(6.5), cy + Inches(0.1), Inches(1.1), Inches(0.5),
                     font_size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(2.8), cy + Inches(1.55), Inches(2.0), Inches(0.4),
                 fill_rgb=DARK_NAVY, line_rgb=ORANGE)
        add_text_box(sl, '交渉を開始', Inches(2.8), cy + Inches(1.6), Inches(2.0), Inches(0.3),
                     font_size=Pt(9), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    slide_number_tag(sl, 3, TOTAL)

    # ─ Slide 4: 価格交渉支援 SURPLUS-003 ─
    sl = blank_slide(prs)
    title_bar(sl, 'SURPLUS-003  価格交渉支援', '対象ユーザー: 売り手・買い手  |  チャットUI / 価格履歴 / 推奨価格AI提案', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'LAYOUT_MASTER.md 準拠', Inches(0.2), Inches(0.65), Inches(7.5), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)

    # チャットUI（左）
    add_rect(sl, Inches(0.2), Inches(0.9), Inches(4.0), Inches(7.5),
             fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=ORANGE)
    add_text_box(sl, 'チャットUI（価格交渉）', Inches(0.3), Inches(0.92), Inches(3.8), Inches(0.22),
                 font_size=Pt(8), bold=True, color=ORANGE)

    chats = [
        ('買い手', '¥1,100/個でいかがですか？', False),
        ('売り手', '¥1,200/個が最低ラインです', True),
        ('AI推奨', '推奨価格: ¥1,150/個（市場相場±5%）', None),
        ('買い手', '¥1,150/個で合意します', False),
    ]
    for i, (role, msg, is_seller) in enumerate(chats):
        cy = Inches(1.2) + i * Inches(1.45)
        if is_seller is None:  # AI
            add_rect(sl, Inches(0.4), cy, Inches(3.5), Inches(0.9),
                     fill_rgb=RGBColor(0x05, 0x37, 0x2a), line_rgb=GREEN)
            add_text_box(sl, 'AI推奨', Inches(0.45), cy + Inches(0.04), Inches(1.5), Inches(0.2),
                         font_size=Pt(7), bold=True, color=GREEN)
            add_text_box(sl, msg, Inches(0.45), cy + Inches(0.26), Inches(3.3), Inches(0.55),
                         font_size=Pt(8), color=RGBColor(0xe2, 0xe8, 0xf0))
        elif is_seller:
            add_rect(sl, Inches(0.4), cy, Inches(3.5), Inches(0.9),
                     fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=DARK_NAVY)
            add_text_box(sl, role, Inches(0.45), cy + Inches(0.04), Inches(1.5), Inches(0.2),
                         font_size=Pt(7), bold=True, color=ORANGE)
            add_text_box(sl, msg, Inches(0.45), cy + Inches(0.26), Inches(3.3), Inches(0.55),
                         font_size=Pt(8), color=RGBColor(0xe2, 0xe8, 0xf0))
        else:
            add_rect(sl, Inches(0.4), cy, Inches(3.5), Inches(0.9),
                     fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=LIGHT_BLUE)
            add_text_box(sl, role, Inches(0.45), cy + Inches(0.04), Inches(1.5), Inches(0.2),
                         font_size=Pt(7), bold=True, color=LIGHT_BLUE)
            add_text_box(sl, msg, Inches(0.45), cy + Inches(0.26), Inches(3.3), Inches(0.55),
                         font_size=Pt(8), color=RGBColor(0xe2, 0xe8, 0xf0))

    # 価格履歴（右上）
    add_rect(sl, Inches(4.3), Inches(0.9), Inches(3.5), Inches(4.0),
             fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=ORANGE)
    add_text_box(sl, '価格履歴', Inches(4.4), Inches(0.92), Inches(3.2), Inches(0.22),
                 font_size=Pt(8), bold=True, color=ORANGE)
    history = [('2026-06-01', '¥1,300/個'), ('2026-06-02', '¥1,250/個'),
               ('2026-06-03', '¥1,200/個'), ('2026-06-04', '¥1,180/個')]
    for i, (date, price) in enumerate(history):
        hy = Inches(1.2) + i * Inches(0.45)
        add_text_box(sl, date, Inches(4.4), hy, Inches(1.7), Inches(0.3),
                     font_size=Pt(8), color=RGBColor(0x64, 0x74, 0x8b))
        add_text_box(sl, price, Inches(6.2), hy, Inches(1.4), Inches(0.3),
                     font_size=Pt(9), bold=True, color=WHITE)

    # AI推奨価格
    add_rect(sl, Inches(4.3), Inches(5.0), Inches(3.5), Inches(3.4),
             fill_rgb=RGBColor(0x05, 0x37, 0x2a), line_rgb=GREEN)
    add_text_box(sl, 'AI推奨価格提案', Inches(4.4), Inches(5.02), Inches(3.2), Inches(0.22),
                 font_size=Pt(8), bold=True, color=GREEN)
    add_text_box(sl, '¥1,150/個', Inches(4.4), Inches(5.3), Inches(3.2), Inches(0.6),
                 font_size=Pt(24), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, '市場相場: ¥1,100〜¥1,200/個\n類似取引平均: ¥1,145/個\n推奨根拠: 直近30日間の市場データ',
                 Inches(4.4), Inches(6.0), Inches(3.2), Inches(0.9),
                 font_size=Pt(8), color=RGBColor(0x6e, 0xe7, 0xb7))
    add_rect(sl, Inches(4.4), Inches(7.05), Inches(3.1), Inches(0.4),
             fill_rgb=GREEN, line_rgb=GREEN)
    add_text_box(sl, 'この価格で合意する', Inches(4.4), Inches(7.1), Inches(3.1), Inches(0.3),
                 font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number_tag(sl, 4, TOTAL)

    # ─ Slide 5: LAYOUT_MASTER準拠指定 ─
    sl = blank_slide(prs)
    title_bar(sl, 'LAYOUT_MASTER準拠指定', 'font-mono tabular-nums tracking-tight / カラーパレット / IndexedDB v142', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    # フォント仕様
    add_text_box(sl, 'フォント仕様', Inches(0.2), Inches(0.7), Inches(7.6), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    font_rows = [
        ['数値表示全般', 'font-mono tabular-nums tracking-tight', '全KPI・価格・数量・滞留日数'],
        ['見出し', 'font-sans font-bold', '画面タイトル・セクション見出し'],
        ['本文', 'font-sans', '説明文・ラベル'],
        ['エラーメッセージ', 'font-mono font-bold text-red-500', 'バリデーションエラー・アラート'],
    ]
    t = sl.shapes.add_table(len(font_rows) + 1, 3, Inches(0.2), Inches(1.0), Inches(7.6), Inches(2.0)).table
    for i, h in enumerate(['適用箇所', 'CSSクラス', '備考']):
        t.cell(0, i).text = h
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)
    for ri, row in enumerate(font_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8)

    # カラーパレット
    add_text_box(sl, 'カラーパレット', Inches(0.2), Inches(3.2), Inches(7.6), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    colors_def = [
        ('Navy #1A3A5C', DARK_NAVY, 'ヘッダー・プライマリ'),
        ('Orange #F5A623', ORANGE, 'アクション・アクセント'),
        ('White #FFFFFF', WHITE, 'テキスト（ダーク背景）'),
        ('Light #F0F4F8', LIGHT_BG, 'カード背景・セカンダリ'),
        ('Green #059669', GREEN, '合意・成功・AI推奨'),
        ('Red #EF4444', RED, 'エラー・警告・処分'),
    ]
    for i, (name, color, usage) in enumerate(colors_def):
        cx = Inches(0.2) + (i % 3) * Inches(2.55)
        cy = Inches(3.55) + (i // 3) * Inches(0.9)
        add_rect(sl, cx, cy, Inches(0.55), Inches(0.55),
                 fill_rgb=color, line_rgb=DARK_NAVY)
        add_text_box(sl, name, cx + Inches(0.6), cy + Inches(0.02), Inches(1.7), Inches(0.25),
                     font_size=Pt(8), bold=True, color=DARK_NAVY)
        add_text_box(sl, usage, cx + Inches(0.6), cy + Inches(0.28), Inches(1.7), Inches(0.22),
                     font_size=Pt(7), color=RGBColor(0x47, 0x55, 0x69))

    # IndexedDB仕様
    add_text_box(sl, 'IndexedDB v142 仕様', Inches(0.2), Inches(5.6), Inches(7.6), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    idb_rows = [
        ['surplus_inventory', 'SURPLUS_SHIFT在庫データ', '_idb_version: 142'],
        ['match_candidates', 'マッチング候補キャッシュ', '_idb_version: 142'],
        ['negotiation_history', '交渉履歴ローカルキャッシュ', '_idb_version: 142'],
        ['price_suggestions', 'AI価格提案キャッシュ（TTL: 1h）', '_idb_version: 142'],
    ]
    t2 = sl.shapes.add_table(len(idb_rows) + 1, 3, Inches(0.2), Inches(5.9), Inches(7.6), Inches(1.8)).table
    for i, h in enumerate(['ストア名', '用途', 'バージョン']):
        t2.cell(0, i).text = h
        t2.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t2.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)
    for ri, row in enumerate(idb_rows):
        for ci, val in enumerate(row):
            t2.cell(ri + 1, ci).text = val
            t2.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8)

    slide_number_tag(sl, 5, TOTAL)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ─────────────────────────────────────────────────────────────────
# PHASE: フェーズ計画
# ─────────────────────────────────────────────────────────────────

def gen_phase(output='docs/SURPLUS_SHIFT/SURPLUS_SHIFT_PHASE_v1.0.pptx'):
    prs = new_prs_landscape()
    FOOTER = 'SURPLUS_SHIFT  |  フェーズ計画 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # ─ Slide 1: タイトル ─
    sl = blank_slide(prs)
    title_bar(sl, 'SURPLUS_SHIFT フェーズ計画 v1.0', '余剰在庫転換システム — フェーズ計画', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_rect(sl, Inches(1.5), Inches(1.2), Inches(8.5), Inches(3.5), fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
    add_text_box(sl, 'SURPLUS_SHIFT', Inches(1.7), Inches(1.4), Inches(8.0), Inches(0.8),
                 font_size=Pt(32), bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'フェーズ計画 v1.0', Inches(1.7), Inches(2.3), Inches(8.0), Inches(0.5),
                 font_size=Pt(18), color=ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, '2026-06-05  |  © 2026 株式会社NiceEze  Confidential',
                 Inches(1.7), Inches(2.9), Inches(8.0), Inches(0.3),
                 font_size=Pt(11), color=RGBColor(0x47, 0x55, 0x69), align=PP_ALIGN.CENTER)

    meta_rows = [
        ['文書名', 'SURPLUS_SHIFT フェーズ計画 v1.0'],
        ['作成日', '2026-06-05'],
        ['作成者', '自律COO（Claude Code）'],
        ['承認者', '代表取締役CEO 松浦 学'],
    ]
    t = sl.shapes.add_table(len(meta_rows), 2, Inches(0.3), Inches(4.8), Inches(5.0), Inches(1.6)).table
    for ri, (k, v) in enumerate(meta_rows):
        t.cell(ri, 0).text = k
        t.cell(ri, 1).text = v
        for ci in range(2):
            t.cell(ri, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(9)

    slide_number_tag(sl, 1, TOTAL, is_landscape=True)

    # ─ Slide 2: Gate制概要表 ─
    sl = blank_slide(prs)
    title_bar(sl, 'Gate制概要表 G0〜G4', 'SURPLUS_SHIFT 開発フェーズゲート', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    gates = [
        ('G0', '設計完了', '2026/07末', '詳細設計書完成\nDB設計確定\nGCP基盤構築', GREEN),
        ('G1', 'MVP', '2026/09末', 'SURPLUS-001\n売り手ダッシュボード\nSURPLUS-002\nマッチング基本', ORANGE),
        ('G2', '本格運用', '2026/11末', '価格交渉支援\n契約書自動生成\n決済連携', DARK_NAVY),
        ('G3', 'AI高度化', '2027/01末', 'Claude API統合\n予測精度向上\n多言語対応', RGBColor(0x6b, 0x21, 0xa8)),
        ('G4', 'Go-Live', '2027/02末', '本番デプロイ\n全機能リリース\nSLA設定', RED),
    ]

    gate_w = Inches(2.0)
    for i, (gate, milestone, date, content, color) in enumerate(gates):
        gx = Inches(0.3) + i * Inches(2.22)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(5.8), fill_rgb=LIGHT_BG, line_rgb=color)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(0.8), fill_rgb=color, line_rgb=color)
        add_text_box(sl, gate, gx, Inches(0.72), gate_w, Inches(0.35),
                     font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, milestone, gx, Inches(1.08), gate_w, Inches(0.28),
                     font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, date, gx, Inches(1.4), gate_w, Inches(0.25),
                     font_size=Pt(8), color=LIGHT_BG, align=PP_ALIGN.CENTER)
        add_text_box(sl, content, gx + Inches(0.1), Inches(1.7), gate_w - Inches(0.2), Inches(4.5),
                     font_size=Pt(9), color=DARK_NAVY)

    add_text_box(sl, '◀─────────────── SURPLUS_SHIFT 開発ロードマップ（〜2027年2月） ───────────────▶',
                 Inches(0.3), Inches(6.7), Inches(10.8), Inches(0.3),
                 font_size=Pt(9), bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 2, TOTAL, is_landscape=True)

    # ─ Slide 3: G1マイルストーン ─
    sl = blank_slide(prs)
    title_bar(sl, 'G1マイルストーン — MVP（2026/09末）', 'SURPLUS-001売り手ダッシュボード / SURPLUS-002マッチング基本', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'G1 完了条件 / 受け入れ基準', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)

    g1_tasks = [
        ['SURPLUS-001 売り手ダッシュボード実装', '自律COO', '未着手', '2026/08末', '¥0'],
        ['余剰在庫自動検知エンジン（30日閾値）', '自律COO', '未着手', '2026/08末', '¥0'],
        ['SURPLUS-002 買い手検索・マッチング', '自律COO', '未着手', '2026/09末', '¥0'],
        ['マッチングスコアリングロジック', '自律COO', '未着手', '2026/09末', '¥0'],
        ['Cloud Functions 定期スキャン（cron毎時）', '自律COO', '未着手', '2026/09末', '¥0'],
        ['Firestore スキーマ設計（surplus_inventory）', '自律COO', '未着手', '2026/07末', '¥0'],
        ['LINE通知連携（マッチング通知）', '自律COO', '未着手', '2026/09末', '¥0'],
        ['IndexedDB v142 キャッシュ設計', '自律COO', '未着手', '2026/08末', '¥0'],
    ]
    t = sl.shapes.add_table(len(g1_tasks) + 1, 5, Inches(0.3), Inches(1.0), Inches(10.8), Inches(4.5)).table
    for i, h in enumerate(['タスク', '担当', '状態', '期限', '費用']):
        t.cell(0, i).text = h
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)
    for ri, row in enumerate(g1_tasks):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8)

    add_text_box(sl, 'G1完了基準: SURPLUS-001/002がステージング環境で動作確認済 / マッチング精度 ≥70%',
                 Inches(0.3), Inches(5.7), Inches(10.8), Inches(0.3),
                 font_size=Pt(9), bold=True, color=ORANGE)
    slide_number_tag(sl, 3, TOTAL, is_landscape=True)

    # ─ Slide 4: G2マイルストーン ─
    sl = blank_slide(prs)
    title_bar(sl, 'G2マイルストーン — 本格運用（2026/11末）', '価格交渉支援 / 契約書自動生成', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'G2 完了条件 / 受け入れ基準', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)

    g2_tasks = [
        ['SURPLUS-003 価格交渉支援UI実装', '自律COO', '未着手', '2026/10末', '¥0'],
        ['チャットUI + リアルタイム同期（Firestore）', '自律COO', '未着手', '2026/10末', '¥0'],
        ['AI推奨価格エンジン（Claude API統合）', '自律COO', '未着手', '2026/11末', '¥2,250/月〜'],
        ['価格履歴ストレージ設計', '自律COO', '未着手', '2026/10末', '¥0'],
        ['契約書自動生成（テンプレートエンジン）', '自律COO', '未着手', '2026/11末', '¥0'],
        ['PDF出力機能（契約書）', '自律COO', '未着手', '2026/11末', '¥0'],
        ['決済ゲートウェイ連携（Stripe/PayPay）', '自律COO', '未着手', '2026/11末', '¥0+手数料'],
        ['本番環境デプロイ（Cloud Run）', '自律COO', '未着手', '2026/11末', '¥750〜2,250/月'],
    ]
    t = sl.shapes.add_table(len(g2_tasks) + 1, 5, Inches(0.3), Inches(1.0), Inches(10.8), Inches(4.5)).table
    for i, h in enumerate(['タスク', '担当', '状態', '期限', '費用']):
        t.cell(0, i).text = h
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)
    for ri, row in enumerate(g2_tasks):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8)

    add_text_box(sl, 'G2完了基準: SURPLUS-003動作確認 / 契約書生成10秒以内 / 決済フロー疎通確認',
                 Inches(0.3), Inches(5.7), Inches(10.8), Inches(0.3),
                 font_size=Pt(9), bold=True, color=ORANGE)
    slide_number_tag(sl, 4, TOTAL, is_landscape=True)

    # ─ Slide 5: FinOps計画 ─
    sl = blank_slide(prs)
    title_bar(sl, 'FinOps計画 — コスト管理', 'MVP¥0/月 / G3以降¥2,250〜¥4,500/月 / 月額上限¥5,000', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'フェーズ別月次コスト', Inches(0.3), Inches(0.7), Inches(5.5), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    finops_rows = [
        ['G0〜G1（MVP）', '¥0/月', '開発人件費¥0', 'AI自律開発'],
        ['G2（本格運用開始）', '¥750〜2,250/月', 'Cloud Run + Firestore', '従量課金'],
        ['G3（AI高度化）', '¥2,250〜4,500/月', '+ Claude API', '月1,000回呼出想定'],
        ['G4（Go-Live）', '¥3,000〜5,000/月', '全サービス稼働', '月額上限¥5,000'],
    ]
    t = sl.shapes.add_table(len(finops_rows) + 1, 4, Inches(0.3), Inches(1.0), Inches(5.5), Inches(2.2)).table
    for i, h in enumerate(['フェーズ', '月額', '内訳', '備考']):
        t.cell(0, i).text = h
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    for ri, row in enumerate(finops_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)

    # Claude API詳細
    add_text_box(sl, 'Claude API コスト試算（G3以降）', Inches(6.0), Inches(0.7), Inches(5.0), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    api_rows = [
        ['価格交渉AI提案', '月500回', '¥1,125'],
        ['余剰度スコアリング', '月300回', '¥675'],
        ['契約書生成補助', '月200回', '¥450'],
        ['合計', '月1,000回', '¥2,250/月'],
        ['上限設定', '月2,000回', '¥4,500/月'],
    ]
    t2 = sl.shapes.add_table(len(api_rows) + 1, 3, Inches(6.0), Inches(1.0), Inches(5.0), Inches(2.2)).table
    for i, h in enumerate(['用途', '呼出回数', '月額概算']):
        t2.cell(0, i).text = h
        t2.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t2.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    for ri, row in enumerate(api_rows):
        for ci, val in enumerate(row):
            t2.cell(ri + 1, ci).text = val
            t2.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)

    # GCPコスト内訳
    add_text_box(sl, 'GCP月次コスト内訳（G4 本番稼働時）', Inches(0.3), Inches(3.4), Inches(10.8), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    gcp_rows = [
        ['Cloud Run（マッチングエンジン）', '¥750〜2,250', '従量制・リクエスト数に比例'],
        ['Firestore（在庫・マッチングDB）', '¥0〜450', '読取/書込回数に依存'],
        ['Cloud Functions（定期スキャン）', '¥0〜150', '毎時実行・無料枠内'],
        ['BigQuery（分析）', '¥300', '月次バッチ処理'],
        ['LINE Messaging API', '¥750〜1,500', '通知件数に比例'],
        ['月額合計', '¥1,800〜4,650', '月額上限¥5,000で管理'],
    ]
    t3 = sl.shapes.add_table(len(gcp_rows) + 1, 3, Inches(0.3), Inches(3.7), Inches(10.8), Inches(2.5)).table
    for i, h in enumerate(['サービス', '月額概算', '備考']):
        t3.cell(0, i).text = h
        t3.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t3.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    for ri, row in enumerate(gcp_rows):
        for ci, val in enumerate(row):
            t3.cell(ri + 1, ci).text = val
            t3.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)

    add_text_box(sl, 'FinOps原則: AI自律開発により開発人件費¥0 / インフラはG4まで月額¥5,000上限で厳格管理',
                 Inches(0.3), Inches(6.3), Inches(10.8), Inches(0.3),
                 font_size=Pt(9), bold=True, color=GREEN)
    slide_number_tag(sl, 5, TOTAL, is_landscape=True)

    prs.save(output)
    print(f"✅ {output} 生成完了")


if __name__ == '__main__':
    gen_seq()
    gen_ui()
    gen_phase()
