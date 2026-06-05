"""GOV部 5ドキュメント生成スクリプト (BRD/SRS/SEQ/UI/PHASE)"""
from docx import Document
from docx.shared import Pt, Mm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
import os

# ─── カラー定義 ────────────────────────────────────────────────────
DARK_NAVY = RGBColor(0x1a, 0x3a, 0x5c)
ORANGE    = RGBColor(0xf5, 0xa6, 0x23)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)

P_DARK_NAVY = PRGBColor(0x1a, 0x3a, 0x5c)
P_ORANGE    = PRGBColor(0xf5, 0xa6, 0x23)
P_WHITE     = PRGBColor(0xFF, 0xFF, 0xFF)
P_LIGHT_BG  = PRGBColor(0xF0, 0xF4, 0xF8)
P_GREEN     = PRGBColor(0x05, 0x96, 0x69)
P_RED       = PRGBColor(0xEF, 0x44, 0x44)
P_SLATE     = PRGBColor(0x47, 0x55, 0x69)

OUTPUT_DIR = "docs/GOV"

A4_W_LAND = Inches(11.69)
A4_H_LAND = Inches(8.27)
A4_W      = Inches(8.27)
A4_H      = Inches(11.69)


# ═══════════════════════════════════════════════════════════════════
# DOCX ユーティリティ
# ═══════════════════════════════════════════════════════════════════

def set_margins(doc, top=20, bottom=20, left=20, right=20):
    sec = doc.sections[0]
    sec.top_margin    = Mm(top)
    sec.bottom_margin = Mm(bottom)
    sec.left_margin   = Mm(left)
    sec.right_margin  = Mm(right)


def add_heading(doc, text, level=1, color=None):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after  = Pt(2)
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(12) if level == 1 else Pt(10.5)
    run.font.color.rgb = color or DARK_NAVY
    return p


def add_body(doc, text):
    p = doc.add_paragraph(text)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after  = Pt(2)
    for run in p.runs:
        run.font.size = Pt(10.5)
    return p


def set_cell_bg(cell, color_hex):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), color_hex)
    tcPr.append(shd)


def set_table_width(table, width_dxa=9026):
    tbl = table._tbl
    tblPr = tbl.find(qn('w:tblPr'))
    if tblPr is None:
        tblPr = OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    tblW = OxmlElement('w:tblW')
    tblW.set(qn('w:w'), str(width_dxa))
    tblW.set(qn('w:type'), 'dxa')
    tblPr.append(tblW)


def add_full_table(doc, headers, rows, header_bg='1a3a5c', header_fg=WHITE):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    set_table_width(table, 9026)
    table.style = 'Table Grid'

    hdr_row = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr_row.cells[i]
        set_cell_bg(cell, header_bg)
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(h)
        run.bold = True
        run.font.size = Pt(9)
        run.font.color.rgb = header_fg

    for ri, row_data in enumerate(rows):
        tr = table.rows[ri + 1]
        bg = 'F0F4F8' if ri % 2 == 0 else 'FFFFFF'
        for ci, val in enumerate(row_data):
            cell = tr.cells[ci]
            set_cell_bg(cell, bg)
            p = cell.paragraphs[0]
            run = p.add_run(str(val))
            run.font.size = Pt(9)

    return table


def add_footer_docx(doc, text):
    for section in doc.sections:
        footer = section.footer
        p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
        p.clear()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x64, 0x74, 0x8b)


# ═══════════════════════════════════════════════════════════════════
# PPTX ユーティリティ
# ═══════════════════════════════════════════════════════════════════

def new_prs_landscape():
    prs = Presentation()
    prs.slide_width  = A4_W_LAND
    prs.slide_height = A4_H_LAND
    return prs


def new_prs_portrait():
    prs = Presentation()
    prs.slide_width  = A4_W
    prs.slide_height = A4_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=PPt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=PPt(9), bold=False, color=None,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color or P_DARK_NAVY
    return txBox


def title_bar(slide, title, subtitle='', is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W
    add_rect(slide, 0, 0, sw, Inches(0.6), fill_rgb=P_DARK_NAVY)
    add_text_box(slide, title, Inches(0.15), Inches(0.05), sw - Inches(0.3), Inches(0.35),
                 font_size=PPt(13), bold=True, color=P_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.15), Inches(0.38), sw - Inches(0.3), Inches(0.2),
                     font_size=PPt(8), color=P_LIGHT_BG, align=PP_ALIGN.LEFT)


def footer_bar(slide, text, is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W
    sh = A4_H_LAND if is_landscape else A4_H
    add_rect(slide, 0, sh - Inches(0.25), sw, Inches(0.25), fill_rgb=P_DARK_NAVY)
    add_text_box(slide, text, Inches(0.1), sh - Inches(0.23), sw - Inches(0.2), Inches(0.2),
                 font_size=PPt(7), color=P_LIGHT_BG, align=PP_ALIGN.CENTER)


def slide_number_tag(slide, num, total, is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W
    sh = A4_H_LAND if is_landscape else A4_H
    add_text_box(slide, f'{num} / {total}',
                 sw - Inches(0.8), sh - Inches(0.23), Inches(0.7), Inches(0.2),
                 font_size=PPt(7), color=P_LIGHT_BG, align=PP_ALIGN.RIGHT)


def pptx_table(slide, headers, rows, x, y, w, h):
    t = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    for ci, h in enumerate(headers):
        t.cell(0, ci).text = h
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = str(val)
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)
    return t


# ═══════════════════════════════════════════════════════════════════
# 1. GOV_BRD_v1.0.docx
# ═══════════════════════════════════════════════════════════════════

def gen_brd():
    output = f"{OUTPUT_DIR}/GOV_BRD_v1.0.docx"
    doc = Document()
    set_margins(doc)
    add_footer_docx(doc, 'GOV部 — ビジネス要件定義書 v1.0  |  COO経営支援・監査システム部  |  © 2026 株式会社NiceEze  Confidential')

    # 表紙
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('GOV部 — ビジネス要件定義書 (BRD) v1.0')
    run.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_NAVY

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run('COO経営支援・監査システム部　v1.0').font.size = Pt(12)

    doc.add_paragraph()

    # 1. 文書管理表
    add_heading(doc, '1. 文書管理表')
    add_full_table(doc,
        ['項目', '内容'],
        [
            ['文書名', 'GOV部 ビジネス要件定義書'],
            ['バージョン', 'v1.0'],
            ['作成日', '2026-06-05'],
            ['最終更新日', '2026-06-05'],
            ['作成者', 'NiceEze 自律COO'],
            ['承認者（CEO承認日）', '代表取締役CEO 松浦 学（2026-06-05）'],
            ['ステータス', '承認済'],
            ['関連文書', 'GOV_SRS_v1.0.docx / GOV_SEQ_v1.0.pptx / GOV_UI_v1.0.pptx / GOV_PHASE_v1.0.pptx'],
        ]
    )

    # 2. ビジネス背景
    add_heading(doc, '2. ビジネス背景')
    add_body(doc,
        '自律経営執行システムにおけるガバナンス・コンプライアンス要件として、ISMS（ISO/IEC 27001:2022）準拠の証跡管理、'
        'FinOpsによる配送コスト監視（1配送¥0.5円上限・月額¥5,000以内）、AIガバナンス台帳（Claude API使用記録・Gate承認履歴）、'
        'DevSecOps統制（bandit/CI/CD自動実行）が不可欠である。'
        'COO経営支援・監査システム部（GOV部）はこれらの要件を統合し、'
        '自律経営の透明性・説明責任・継続的監査を実現する中枢システムとして機能する。'
    )

    # 3. 主要機能一覧
    add_heading(doc, '3. 主要機能一覧')
    add_full_table(doc,
        ['ID', '機能名', '概要', '優先度'],
        [
            ['GOV-001', 'COO業務報告ダッシュボード',
             'KPI実績/予実差/PMOタスク管理。S10 COO業務報告（月次自動生成）',
             'Must'],
            ['GOV-002', 'ISMS適合自動レポート',
             'ISO/IEC 27001:2022準拠。SHA-256署名付き月次レポート自動生成',
             'Must'],
            ['GOV-003', 'FinOps監視',
             '1配送¥0.5円上限アラート・月額¥5,000上限監視。cost_calculator.py連携',
             'Must'],
            ['GOV-004', 'AIガバナンス台帳',
             'Claude API使用記録・判断ログ・Gate承認履歴。Firestore保存',
             'Must'],
            ['GOV-005', 'DevSecOps統制',
             'bandit -r src/・GitHub Actions CI・lint・セキュリティスキャン自動実行',
             'Must'],
        ]
    )

    # 4. ステークホルダー
    add_heading(doc, '4. ステークホルダー定義')
    add_full_table(doc,
        ['ステークホルダー', '役割', '主要タッチポイント', '優先度'],
        [
            ['代表取締役CEO 松浦 学', '最終承認・経営判断', 'COO業務報告ダッシュボード / チャット通知', 'Must'],
            ['自律COO', 'システム自律実行・KPI報告', 'GOV全機能（自動実行）', 'Must'],
            ['監査担当', 'ISMS・内部監査・証跡確認', 'ISMS自動レポート / AIガバナンス台帳', 'Must'],
            ['GCP管理者', 'インフラ運用・コスト管理', 'Cloud Console / Secret Manager / FinOps監視', 'Should'],
        ]
    )

    # 5. KPI
    add_heading(doc, '5. 成功指標（KPI）')
    add_full_table(doc,
        ['KPI', '目標値', '計測方法', 'Gate'],
        [
            ['ISMS適合率', '100%', 'ISMS月次自動レポート', 'G2'],
            ['FinOps逸脱件数', '0件', 'cost_calculator.py自動監視', 'G1'],
            ['月次監査レポート自動生成率', '100%', 'Cloud Scheduler月次実行確認', 'G2'],
            ['bandit検出件数', '0件継続', 'GitHub Actions CI自動スキャン', '全Gate'],
            ['AIガバナンス台帳記録漏れ', '0件', 'Firestore append-only確認', 'G3'],
            ['FinOps監視コスト', '¥0/月', 'Cloud Functions無料枠内確認', 'G1'],
        ]
    )

    # 6. Gate制
    add_heading(doc, '6. Gate制（G0〜G4）')
    add_full_table(doc,
        ['Gate', '完了条件', 'FinOps上限', '備考'],
        [
            ['G0', '設計完了・GCP環境構築・Secret Manager設定', '¥0（無料枠）', 'インフラ基盤'],
            ['G1', 'FinOps監視稼働（¥0.5円アラート）・bandit CI統合・DevSecOps統制基盤', '¥0', 'Cloud Functions無料枠'],
            ['G2', 'ISMS月次自動レポート・SHA-256署名・BigQuery月次アーカイブ連携', '¥0', 'G2完了時点'],
            ['G3', 'AIガバナンス台帳完成・Claude API使用記録・Gate承認履歴', '¥0', 'G3完了時点'],
            ['G4', 'Go-Live完全自律監査・全KPI達成・CEO最終承認', '¥0', 'Cloud Functions無料枠内維持'],
        ]
    )

    # 7. FinOps
    add_heading(doc, '7. FinOps管理方針')
    add_body(doc,
        'GOVシステム自体の運用コストは¥0/月（Cloud Functions無料枠内）を維持する。'
        'FinOps監視の対象は全社システム月額¥5,000以内（1配送¥0.5円上限）。'
        'G4以降もClaude API監視費用を含む全コストを月額¥5,000上限で管理する。'
    )
    add_full_table(doc,
        ['監視対象', '上限値', 'アラート閾値', '対応'],
        [
            ['1配送コスト', '¥0.5円', '¥0.45円（90%）', '即時アラート・ブロッカー報告'],
            ['月額総コスト', '¥5,000', '¥4,500（90%）', 'Cloud Monitoringアラート'],
            ['GOVシステム自体', '¥0/月', '無料枠超過時', 'Cloud Functions設定見直し'],
        ]
    )

    # 8. 制約条件
    add_heading(doc, '8. 制約条件・前提条件')
    add_full_table(doc,
        ['区分', '内容'],
        [
            ['セキュリティ制約', 'ISMS ISO/IEC 27001:2022準拠必須'],
            ['セキュリティ制約', 'AES-256暗号化 / RLS（Row Level Security）必須'],
            ['セキュリティ制約', 'PII最小化原則の遵守'],
            ['セキュリティ制約', 'bandit 0件必須（GitHub Actions CI強制）'],
            ['データ制約', '監査ログ改ざん不可（Firestore append-only）'],
            ['技術制約', 'GCPサーバーレス構成固定（Cloud Functions / BigQuery / Firestore）'],
            ['前提条件', 'Secret ManagerにAPIキー設定済み（Gate 0）'],
            ['前提条件', 'cost_calculator.pyが既存システムと連携済み'],
        ]
    )

    doc.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 2. GOV_SRS_v1.0.docx
# ═══════════════════════════════════════════════════════════════════

def gen_srs():
    output = f"{OUTPUT_DIR}/GOV_SRS_v1.0.docx"
    doc = Document()
    set_margins(doc)
    add_footer_docx(doc, 'GOV部 — ソフトウェア要件仕様書 v1.0  |  COO経営支援・監査システム部  |  © 2026 株式会社NiceEze  Confidential')

    # 表紙
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('GOV部 — ソフトウェア要件仕様書 (SRS) v1.0')
    run.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_NAVY

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run('COO経営支援・監査システム部　v1.0  |  2026-06-05').font.size = Pt(12)
    doc.add_paragraph()

    # 1. システム概要
    add_heading(doc, '1. システム概要')
    add_body(doc,
        'GOVシステムはCloud Functions（イベント駆動・月次定期実行）+ BigQuery（月次アーカイブ）'
        '+ Firestore（監査ログ・AIガバナンス台帳）+ Secret Manager（APIキー管理）で構成される。'
        'Cloud Run（KPIダッシュボードAPI）を補助的に使用し、全コンポーネントはGCPサーバーレス構成で運用する。'
    )

    # 2. 機能要件
    add_heading(doc, '2. 機能要件')
    add_full_table(doc,
        ['ID', '機能名', '実装技術', '詳細要件'],
        [
            ['GOV-001', 'KPIダッシュボード',
             'Cloud Run / JSON API',
             '実績vs予算差分計算・月次KPIレポート自動生成。GET /kpi エンドポイント提供'],
            ['GOV-002', 'ISMS自動レポート',
             'hashlib.sha256 / Cloud Scheduler',
             'SHA-256署名付きレポート。nosec B303廃止（hashlib使用）。月次自動実行'],
            ['GOV-003', 'FinOps監視',
             'cost_calculator.py / Cloud Monitoring',
             '1配送¥0.5円閾値チェック・月額¥5,000アラート。毎時Cloud Functions実行'],
            ['GOV-004', 'AIガバナンス台帳',
             'Firestore（append-only）',
             'Claude API呼出ログ・Gate承認タイムスタンプ・CEO判断記録保存'],
            ['GOV-005', 'DevSecOps統制',
             'bandit / GitHub Actions',
             'bandit -r src/ 自動実行・lint・失敗時Slack/チャット即時通知'],
        ]
    )

    # 3. 非機能要件
    add_heading(doc, '3. 非機能要件')
    add_full_table(doc,
        ['区分', '要件', '基準値'],
        [
            ['セキュリティ', 'ISMS ISO/IEC 27001:2022準拠', '月次監査レポート100%自動生成'],
            ['セキュリティ', 'AES-256暗号化', '全PII・APIキー対象'],
            ['セキュリティ', 'RLS（Row Level Security）', 'Firestore / BigQuery適用'],
            ['データ整合性', '監査ログ改ざん不可', 'Firestore append-only（削除禁止）'],
            ['プライバシー', 'PII最小化', '必要最小限のみ収集・保存'],
            ['可用性', 'Cloud Functions月次実行成功率', '100%（リトライ3回）'],
            ['コスト', 'GOVシステム月額運用コスト', '¥0（Cloud Functions無料枠内）'],
        ]
    )

    # 4. データモデル
    add_heading(doc, '4. データモデル')
    add_full_table(doc,
        ['コレクション/テーブル', 'ストア', '主要フィールド', '備考'],
        [
            ['AuditLog', 'Firestore', 'id / timestamp / action / actor / sha256 / payload', 'append-only・改ざん不可'],
            ['FinOpsRecord', 'Firestore', 'id / date / cost_per_delivery / monthly_total / alert_fired', '毎時記録'],
            ['ISMSReport', 'BigQuery', 'report_id / period / sha256_signature / compliance_rate / generated_at', '月次アーカイブ'],
            ['GovernanceEntry', 'Firestore', 'entry_id / claude_api_call / gate / ceo_decision / timestamp', 'AIガバナンス台帳'],
            ['CIResult', 'Firestore', 'run_id / bandit_issues / lint_pass / test_pass / triggered_at / status', 'CI/CD結果記録'],
        ]
    )

    # 5. APIエンドポイント
    add_heading(doc, '5. APIエンドポイント')
    add_full_table(doc,
        ['メソッド', 'エンドポイント', '説明', '認証'],
        [
            ['GET', '/health', 'ヘルスチェック', '不要'],
            ['GET', '/kpi', 'KPI実績・予実差取得（JSON）', 'Cloud IAM'],
            ['GET', '/finops/status', 'FinOps監視ステータス取得', 'Cloud IAM'],
            ['POST', '/audit/log', '監査ログ追記（append-only）', 'Cloud IAM + SHA-256'],
            ['GET', '/isms/report', 'ISMS月次レポート取得（SHA-256署名付き）', 'Cloud IAM'],
        ]
    )

    # 6. テスト要件
    add_heading(doc, '6. テスト要件')
    add_full_table(doc,
        ['テスト種別', '対象', '合格基準', '実行タイミング'],
        [
            ['ユニットテスト', 'cost_calculator.py / sha256署名ロジック', 'unittest全Pass', 'git push時'],
            ['セキュリティスキャン', 'src/以下全ファイル', 'bandit 0件', 'GitHub Actions CI必須'],
            ['lint', 'Python / JS全ファイル', '0 error', 'GitHub Actions CI必須'],
            ['PII使用確認', '全APIレスポンス', 'PII不含確認', 'ステージング環境'],
            ['統合テスト', '全APIエンドポイント', '200 OK / 期待レスポンス一致', 'デプロイ前'],
        ]
    )

    doc.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 3. GOV_SEQ_v1.0.pptx
# ═══════════════════════════════════════════════════════════════════

def gen_seq():
    output = f"{OUTPUT_DIR}/GOV_SEQ_v1.0.pptx"
    prs = new_prs_landscape()
    FOOTER = 'GOV部  |  シーケンス図 v1.0  |  COO経営支援・監査システム部  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # ─ S1: タイトル ─
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, A4_W_LAND, A4_H_LAND, fill_rgb=P_DARK_NAVY)
    add_text_box(sl, 'GOV部 シーケンス図 v1.0',
                 Inches(1.5), Inches(2.0), Inches(8.5), Inches(1.0),
                 font_size=PPt(36), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'COO経営支援・監査システム部',
                 Inches(1.5), Inches(3.2), Inches(8.5), Inches(0.5),
                 font_size=PPt(18), color=P_ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, '2026-06-05  |  © 2026 株式会社NiceEze  Confidential',
                 Inches(1.5), Inches(3.9), Inches(8.5), Inches(0.35),
                 font_size=PPt(11), color=P_LIGHT_BG, align=PP_ALIGN.CENTER)
    items = [
        'Slide 2: COO月次報告フロー',
        'Slide 3: FinOps監視フロー',
        'Slide 4: DevSecOps統制フロー',
        'Slide 5: AIガバナンス記録フロー',
    ]
    for i, item in enumerate(items):
        add_text_box(sl, item, Inches(3.5), Inches(4.7) + i * Inches(0.45),
                     Inches(5.0), Inches(0.35), font_size=PPt(10), color=P_LIGHT_BG)
    slide_number_tag(sl, 1, TOTAL)

    # ─ S2: COO月次報告フロー ─
    sl = blank_slide(prs)
    title_bar(sl, 'COO月次報告フロー', 'Cloud Scheduler → BigQuery集計 → KPI計算 → PDF/MD生成 → SHA-256署名 → Drive保存 → CEOチャット通知')
    footer_bar(sl, FOOTER)

    steps = [
        ('Cloud Scheduler\n（月1回）', '月次実行トリガー'),
        ('BigQuery', '月次データ集計・KPI計算（予実差）'),
        ('Cloud Functions\nレポート生成', 'PDF/MD生成・SHA-256署名'),
        ('Google Drive', '署名付きレポート保存'),
        ('CEOチャット\n通知', '月次報告完了通知'),
    ]
    box_w = Inches(1.8)
    spacing = Inches(2.0)
    y_base = Inches(0.75)
    for i, (actor, desc) in enumerate(steps):
        bx = Inches(0.5) + i * spacing
        add_rect(sl, bx, y_base, box_w, Inches(0.5), fill_rgb=P_DARK_NAVY, line_rgb=P_ORANGE)
        add_text_box(sl, actor, bx, y_base, box_w, Inches(0.5),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        # 垂直線
        cx = bx + box_w / 2
        sh = A4_H_LAND
        line = sl.shapes.add_connector(1, cx, y_base + Inches(0.5), cx, Inches(7.0))
        line.line.color.rgb = PRGBColor(0x94, 0xa3, 0xb8)
        line.line.width = PPt(0.5)

    arrow_y = Inches(1.5)
    cxs = [Inches(0.5) + i * spacing + box_w / 2 for i in range(len(steps))]
    flow = [(0,1,'月次トリガー'),(1,2,'集計データ'),(2,2,'KPI計算・差分'),(2,3,'PDF/MD+SHA-256'),(3,4,'保存完了通知')]
    for (fi, ti, label) in flow:
        fx, tx = cxs[fi], cxs[ti]
        if fi == ti:
            add_text_box(sl, f'↺ {label}', fx - Inches(0.5), arrow_y - Inches(0.15),
                         Inches(1.5), Inches(0.25), font_size=PPt(7.5), color=P_ORANGE)
        else:
            arr = sl.shapes.add_connector(2, fx, arrow_y, tx, arrow_y)
            arr.line.color.rgb = P_ORANGE
            arr.line.width = PPt(1.5)
            mid_x = min(fx, tx) + abs(tx - fx) / 2 - Inches(0.5)
            add_text_box(sl, label, mid_x, arrow_y - Inches(0.2), Inches(1.0), Inches(0.18),
                         font_size=PPt(7), color=P_DARK_NAVY, align=PP_ALIGN.CENTER)
        arrow_y += Inches(0.9)

    add_text_box(sl, '✅ 月次監査レポート自動生成100% / SHA-256署名済 / CEOチャット通知',
                 Inches(0.3), Inches(6.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_GREEN, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 2, TOTAL)

    # ─ S3: FinOps監視フロー ─
    sl = blank_slide(prs)
    title_bar(sl, 'FinOps監視フロー', 'Cloud Functions(毎時) → cost_calculator → ¥0.5円/配送チェック → 月額¥5,000チェック → 逸脱 → 即時アラート')
    footer_bar(sl, FOOTER)

    actors3 = ['Cloud Functions\n（毎時）', 'cost_calculator\n.py', '¥0.5円/配送\nチェック', '月額¥5,000\nチェック', '即時アラート\n・ブロッカー']
    for i, actor in enumerate(actors3):
        bx = Inches(0.5) + i * spacing
        add_rect(sl, bx, y_base, box_w, Inches(0.5), fill_rgb=P_DARK_NAVY, line_rgb=P_ORANGE)
        add_text_box(sl, actor, bx, y_base, box_w, Inches(0.5),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        cx = bx + box_w / 2
        line = sl.shapes.add_connector(1, cx, y_base + Inches(0.5), cx, Inches(7.0))
        line.line.color.rgb = PRGBColor(0x94, 0xa3, 0xb8)
        line.line.width = PPt(0.5)

    cxs3 = [Inches(0.5) + i * spacing + box_w / 2 for i in range(len(actors3))]
    flow3 = [(0,1,'毎時起動'),(1,2,'コストデータ'),(2,3,'¥0.5円チェック'),(3,4,'月額チェック')]
    arrow_y3 = Inches(1.5)
    for (fi, ti, label) in flow3:
        arr = sl.shapes.add_connector(2, cxs3[fi], arrow_y3, cxs3[ti], arrow_y3)
        arr.line.color.rgb = P_ORANGE
        arr.line.width = PPt(1.5)
        mid_x = min(cxs3[fi], cxs3[ti]) + abs(cxs3[ti] - cxs3[fi]) / 2 - Inches(0.5)
        add_text_box(sl, label, mid_x, arrow_y3 - Inches(0.2), Inches(1.0), Inches(0.18),
                     font_size=PPt(7), color=P_DARK_NAVY, align=PP_ALIGN.CENTER)
        arrow_y3 += Inches(0.9)

    add_rect(sl, Inches(0.3), Inches(5.4), Inches(10.8), Inches(1.0),
             fill_rgb=PRGBColor(0x7f, 0x1d, 0x1d), line_rgb=P_RED)
    add_text_box(sl, '逸脱検知時: 即時Slack/チャットアラート → ブロッカー報告 → CEO通知',
                 Inches(0.5), Inches(5.6), Inches(10.0), Inches(0.6),
                 font_size=PPt(10), bold=True, color=PRGBColor(0xfc, 0xa5, 0xa5))

    add_text_box(sl, '監視閾値: 1配送¥0.5円上限 / 月額¥5,000上限 / GOVシステム自体¥0/月',
                 Inches(0.3), Inches(6.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_ORANGE, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 3, TOTAL)

    # ─ S4: DevSecOps統制フロー ─
    sl = blank_slide(prs)
    title_bar(sl, 'DevSecOps統制フロー', 'git push → GitHub Actions → bandit -r src/ → lint → テスト → 全Pass → deploy / 失敗 → 即時ブロック')
    footer_bar(sl, FOOTER)

    actors4 = ['git push', 'GitHub\nActions', 'bandit\n-r src/', 'lint\n検査', 'テスト\n実行']
    for i, actor in enumerate(actors4):
        bx = Inches(0.5) + i * spacing
        add_rect(sl, bx, y_base, box_w, Inches(0.5), fill_rgb=P_DARK_NAVY, line_rgb=P_ORANGE)
        add_text_box(sl, actor, bx, y_base, box_w, Inches(0.5),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        cx = bx + box_w / 2
        line = sl.shapes.add_connector(1, cx, y_base + Inches(0.5), cx, Inches(7.0))
        line.line.color.rgb = PRGBColor(0x94, 0xa3, 0xb8)
        line.line.width = PPt(0.5)

    cxs4 = [Inches(0.5) + i * spacing + box_w / 2 for i in range(len(actors4))]
    flow4 = [(0,1,'push'),(1,2,'CI起動'),(2,3,'0件確認'),(3,4,'lint Pass')]
    arrow_y4 = Inches(1.5)
    for (fi, ti, label) in flow4:
        arr = sl.shapes.add_connector(2, cxs4[fi], arrow_y4, cxs4[ti], arrow_y4)
        arr.line.color.rgb = P_ORANGE
        arr.line.width = PPt(1.5)
        mid_x = min(cxs4[fi], cxs4[ti]) + abs(cxs4[ti] - cxs4[fi]) / 2 - Inches(0.5)
        add_text_box(sl, label, mid_x, arrow_y4 - Inches(0.2), Inches(1.0), Inches(0.18),
                     font_size=PPt(7), color=P_DARK_NAVY, align=PP_ALIGN.CENTER)
        arrow_y4 += Inches(0.9)

    add_rect(sl, Inches(0.3), Inches(5.1), Inches(5.0), Inches(0.7),
             fill_rgb=P_GREEN, line_rgb=P_GREEN)
    add_text_box(sl, '✅ 全Pass → deploy', Inches(0.5), Inches(5.2), Inches(4.5), Inches(0.5),
                 font_size=PPt(12), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(5.8), Inches(5.1), Inches(5.0), Inches(0.7),
             fill_rgb=P_RED, line_rgb=P_RED)
    add_text_box(sl, '❌ 失敗 → 即時ブロック → チャット通知', Inches(6.0), Inches(5.2), Inches(4.5), Inches(0.5),
                 font_size=PPt(10), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(sl, 'bandit 0件必須 / lint 0 error必須 / unittest全Pass必須（全Gate共通）',
                 Inches(0.3), Inches(6.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_ORANGE, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 4, TOTAL)

    # ─ S5: AIガバナンス記録フロー ─
    sl = blank_slide(prs)
    title_bar(sl, 'AIガバナンス記録フロー', 'Claude API呼出 → ログ記録(Firestore) → Gate承認タイムスタンプ → CEO判断記録 → 月次台帳生成')
    footer_bar(sl, FOOTER)

    actors5 = ['Claude API\n呼出', 'Firestore\nログ記録', 'Gate承認\nタイムスタンプ', 'CEO判断\n記録', '月次台帳\n生成']
    for i, actor in enumerate(actors5):
        bx = Inches(0.5) + i * spacing
        add_rect(sl, bx, y_base, box_w, Inches(0.5), fill_rgb=P_DARK_NAVY, line_rgb=P_ORANGE)
        add_text_box(sl, actor, bx, y_base, box_w, Inches(0.5),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        cx = bx + box_w / 2
        line = sl.shapes.add_connector(1, cx, y_base + Inches(0.5), cx, Inches(7.0))
        line.line.color.rgb = PRGBColor(0x94, 0xa3, 0xb8)
        line.line.width = PPt(0.5)

    cxs5 = [Inches(0.5) + i * spacing + box_w / 2 for i in range(len(actors5))]
    flow5 = [(0,1,'API呼出ログ'),(1,2,'append-only'),(2,3,'Gate承認TS'),(3,4,'CEO判断記録')]
    arrow_y5 = Inches(1.5)
    for (fi, ti, label) in flow5:
        arr = sl.shapes.add_connector(2, cxs5[fi], arrow_y5, cxs5[ti], arrow_y5)
        arr.line.color.rgb = P_ORANGE
        arr.line.width = PPt(1.5)
        mid_x = min(cxs5[fi], cxs5[ti]) + abs(cxs5[ti] - cxs5[fi]) / 2 - Inches(0.5)
        add_text_box(sl, label, mid_x, arrow_y5 - Inches(0.2), Inches(1.0), Inches(0.18),
                     font_size=PPt(7), color=P_DARK_NAVY, align=PP_ALIGN.CENTER)
        arrow_y5 += Inches(0.9)

    add_text_box(sl, '✅ AIガバナンス台帳: Firestore append-only / 改ざん不可 / 月次BigQueryアーカイブ',
                 Inches(0.3), Inches(6.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_GREEN, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 5, TOTAL)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 4. GOV_UI_v1.0.pptx
# ═══════════════════════════════════════════════════════════════════

def gen_ui():
    output = f"{OUTPUT_DIR}/GOV_UI_v1.0.pptx"
    prs = new_prs_portrait()
    FOOTER = 'GOV部  |  UI設計 v1.0  |  COO経営支援・監査システム部  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5
    IL = False  # is_landscape = False

    # ─ S1: タイトル ─
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, A4_W, A4_H, fill_rgb=P_DARK_NAVY)
    add_text_box(sl, 'GOV部 UI設計 v1.0',
                 Inches(0.4), Inches(2.0), Inches(7.3), Inches(1.0),
                 font_size=PPt(32), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'COO経営支援・監査システム部',
                 Inches(0.4), Inches(3.2), Inches(7.3), Inches(0.5),
                 font_size=PPt(16), color=P_ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'v1.0  |  2026-06-05  |  © 2026 株式会社NiceEze',
                 Inches(0.4), Inches(3.9), Inches(7.3), Inches(0.35),
                 font_size=PPt(10), color=P_LIGHT_BG, align=PP_ALIGN.CENTER)
    screens = [
        ('S2', 'S10 COO業務報告ダッシュボード', 'KPIカード / 予実グラフ / PMOタスク一覧 / 月次レポートボタン'),
        ('S3', 'FinOps監視画面', '1配送コスト / 月額累計 / 警告ライン / リアルタイム更新'),
        ('S4', 'AIガバナンス台帳', 'Claude API使用ログ / Gate承認履歴 / CEO判断記録'),
        ('S5', 'LAYOUT_MASTER準拠', 'font-mono tabular-nums / #1A2B4C/#E8A020 / CI/CDバッジ'),
    ]
    for i, (sn, name, desc) in enumerate(screens):
        y = Inches(4.8) + i * Inches(1.3)
        add_rect(sl, Inches(0.4), y, Inches(7.3), Inches(1.1), fill_rgb=PRGBColor(0x1e,0x29,0x3b), line_rgb=P_ORANGE)
        add_text_box(sl, sn, Inches(0.5), y + Inches(0.1), Inches(0.7), Inches(0.4),
                     font_size=PPt(10), bold=True, color=P_ORANGE)
        add_text_box(sl, name, Inches(1.3), y + Inches(0.05), Inches(6.0), Inches(0.35),
                     font_size=PPt(10), bold=True, color=P_WHITE)
        add_text_box(sl, desc, Inches(1.3), y + Inches(0.45), Inches(6.0), Inches(0.5),
                     font_size=PPt(8), color=P_LIGHT_BG)
    slide_number_tag(sl, 1, TOTAL, is_landscape=IL)

    # ─ S2: S10 COO業務報告ダッシュボード ─
    sl = blank_slide(prs)
    title_bar(sl, 'S10 COO業務報告ダッシュボード', 'KPIカード(目標vs実績) / 予実グラフ / PMOタスク一覧 / 月次レポートボタン', is_landscape=IL)
    footer_bar(sl, FOOTER, is_landscape=IL)

    add_text_box(sl, 'LAYOUT_MASTER準拠 — font-mono tabular-nums tracking-tight（全数値）',
                 Inches(0.2), Inches(0.65), Inches(7.6), Inches(0.2),
                 font_size=PPt(8), bold=True, color=P_ORANGE)

    # KPIカード行
    add_rect(sl, Inches(0.2), Inches(0.9), Inches(7.6), Inches(1.8),
             fill_rgb=PRGBColor(0x1e,0x29,0x3b), line_rgb=P_ORANGE)
    add_text_box(sl, 'KPIカード（目標 vs 実績）', Inches(0.3), Inches(0.92), Inches(4.0), Inches(0.22),
                 font_size=PPt(8), bold=True, color=P_ORANGE)
    kpi_data = [
        ('ISMS適合率', '目標 100%', '実績 100%', P_GREEN),
        ('FinOps逸脱', '目標 0件', '実績 0件', P_GREEN),
        ('bandit件数', '目標 0件', '実績 0件', P_GREEN),
        ('月次レポート', '目標 100%', '実績 100%', P_GREEN),
    ]
    for i, (label, target, actual, color) in enumerate(kpi_data):
        kx = Inches(0.3) + i * Inches(1.85)
        add_text_box(sl, label, kx, Inches(1.15), Inches(1.75), Inches(0.2),
                     font_size=PPt(7), color=PRGBColor(0x64,0x74,0x8b))
        add_text_box(sl, actual, kx, Inches(1.38), Inches(1.75), Inches(0.35),
                     font_size=PPt(14), bold=True, color=color)
        add_text_box(sl, target, kx, Inches(1.75), Inches(1.75), Inches(0.2),
                     font_size=PPt(6.5), color=PRGBColor(0x64,0x74,0x8b))
    add_text_box(sl, 'font-mono tabular-nums tracking-tight（全KPI数値）',
                 Inches(0.3), Inches(2.55), Inches(7.0), Inches(0.18),
                 font_size=PPt(6.5), color=P_ORANGE)

    # 予実グラフ
    add_rect(sl, Inches(0.2), Inches(2.8), Inches(4.5), Inches(3.5),
             fill_rgb=PRGBColor(0x0f,0x17,0x2a), line_rgb=P_ORANGE)
    add_text_box(sl, '予実差グラフ（月次）', Inches(0.3), Inches(2.82), Inches(4.0), Inches(0.22),
                 font_size=PPt(8), bold=True, color=P_ORANGE)
    months = ['1月', '2月', '3月', '4月', '5月', '6月']
    budgets = [5000, 5000, 5000, 5000, 5000, 5000]
    actuals_val = [0, 0, 0, 0, 0, 0]
    bar_w2 = Inches(0.45)
    for mi, month in enumerate(months):
        bx2 = Inches(0.35) + mi * Inches(0.68)
        # 予算バー
        add_rect(sl, bx2, Inches(5.6), bar_w2, Inches(0.5),
                 fill_rgb=PRGBColor(0x33,0x41,0x55), line_rgb=PRGBColor(0x33,0x41,0x55))
        # 実績バー（¥0）
        add_rect(sl, bx2, Inches(5.6), bar_w2, Inches(0.05),
                 fill_rgb=P_GREEN, line_rgb=P_GREEN)
        add_text_box(sl, month, bx2 - Inches(0.02), Inches(5.68), bar_w2 + Inches(0.04), Inches(0.2),
                     font_size=PPt(6.5), color=PRGBColor(0x94,0xa3,0xb8), align=PP_ALIGN.CENTER)
    add_text_box(sl, '■ 予算上限¥5,000  ■ 実績¥0（GOVシステム無料枠）',
                 Inches(0.3), Inches(5.9), Inches(4.2), Inches(0.3),
                 font_size=PPt(7), color=P_LIGHT_BG)

    # PMOタスク一覧
    add_rect(sl, Inches(4.8), Inches(2.8), Inches(3.0), Inches(3.5),
             fill_rgb=PRGBColor(0x0f,0x17,0x2a), line_rgb=P_ORANGE)
    add_text_box(sl, 'PMOタスク一覧', Inches(4.9), Inches(2.82), Inches(2.8), Inches(0.22),
                 font_size=PPt(8), bold=True, color=P_ORANGE)
    pmo_tasks = [
        ('ISMS月次レポート', '完了', P_GREEN),
        ('FinOps監視設定', '完了', P_GREEN),
        ('bandit CI統合', '完了', P_GREEN),
        ('AIガバナンス台帳', '進行中', P_ORANGE),
        ('G2レビュー準備', '未着手', PRGBColor(0x94,0xa3,0xb8)),
    ]
    for i, (task, status, color) in enumerate(pmo_tasks):
        ty = Inches(3.1) + i * Inches(0.55)
        add_rect(sl, Inches(4.9), ty, Inches(2.8), Inches(0.45),
                 fill_rgb=PRGBColor(0x1e,0x29,0x3b) if i % 2 == 0 else PRGBColor(0x0f,0x17,0x2a),
                 line_rgb=PRGBColor(0x1e,0x29,0x3b))
        add_text_box(sl, task, Inches(4.95), ty + Inches(0.05), Inches(1.9), Inches(0.35),
                     font_size=PPt(7.5), color=P_WHITE)
        add_text_box(sl, status, Inches(6.9), ty + Inches(0.05), Inches(0.8), Inches(0.35),
                     font_size=PPt(7.5), bold=True, color=color)

    # 月次レポートボタン
    add_rect(sl, Inches(0.2), Inches(6.4), Inches(7.6), Inches(0.6),
             fill_rgb=PRGBColor(0x1e,0x29,0x3b), line_rgb=P_ORANGE)
    add_rect(sl, Inches(0.4), Inches(6.5), Inches(3.0), Inches(0.38),
             fill_rgb=P_DARK_NAVY, line_rgb=P_ORANGE)
    add_text_box(sl, '📄 月次レポート生成（SHA-256署名付き）',
                 Inches(0.4), Inches(6.55), Inches(3.0), Inches(0.28),
                 font_size=PPt(8), bold=True, color=P_ORANGE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(3.6), Inches(6.5), Inches(2.0), Inches(0.38),
             fill_rgb=P_GREEN, line_rgb=P_GREEN)
    add_text_box(sl, '🔍 ISMS証跡確認',
                 Inches(3.6), Inches(6.55), Inches(2.0), Inches(0.28),
                 font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    slide_number_tag(sl, 2, TOTAL, is_landscape=IL)

    # ─ S3: FinOps監視画面 ─
    sl = blank_slide(prs)
    title_bar(sl, 'FinOps監視画面', '1配送コスト表示 / 月額累計 / 警告ライン可視化 / リアルタイム更新', is_landscape=IL)
    footer_bar(sl, FOOTER, is_landscape=IL)

    add_text_box(sl, 'font-mono tabular-nums tracking-tight（全コスト数値）',
                 Inches(0.2), Inches(0.65), Inches(7.6), Inches(0.2),
                 font_size=PPt(8), bold=True, color=P_ORANGE)

    # 1配送コスト
    add_rect(sl, Inches(0.2), Inches(0.9), Inches(3.5), Inches(2.5),
             fill_rgb=PRGBColor(0x1e,0x29,0x3b), line_rgb=P_ORANGE)
    add_text_box(sl, '1配送コスト', Inches(0.3), Inches(0.92), Inches(3.0), Inches(0.22),
                 font_size=PPt(8), bold=True, color=P_ORANGE)
    add_text_box(sl, '¥0.32', Inches(0.4), Inches(1.25), Inches(3.0), Inches(0.7),
                 font_size=PPt(36), bold=True, color=P_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(sl, '上限: ¥0.50 / 配送', Inches(0.4), Inches(2.05), Inches(3.0), Inches(0.25),
                 font_size=PPt(8), color=PRGBColor(0x64,0x74,0x8b), align=PP_ALIGN.CENTER)
    # 警告ライン
    add_rect(sl, Inches(0.3), Inches(2.35), Inches(3.1), Inches(0.12),
             fill_rgb=PRGBColor(0x33,0x41,0x55))
    add_rect(sl, Inches(0.3), Inches(2.35), Inches(1.98), Inches(0.12),
             fill_rgb=P_GREEN)
    add_text_box(sl, '64% (¥0.32/¥0.50)', Inches(0.3), Inches(2.5), Inches(3.1), Inches(0.2),
                 font_size=PPt(7), color=P_GREEN, align=PP_ALIGN.CENTER)

    # 月額累計
    add_rect(sl, Inches(3.9), Inches(0.9), Inches(3.9), Inches(2.5),
             fill_rgb=PRGBColor(0x1e,0x29,0x3b), line_rgb=P_ORANGE)
    add_text_box(sl, '月額累計', Inches(4.0), Inches(0.92), Inches(3.5), Inches(0.22),
                 font_size=PPt(8), bold=True, color=P_ORANGE)
    add_text_box(sl, '¥2,340', Inches(4.0), Inches(1.25), Inches(3.5), Inches(0.7),
                 font_size=PPt(30), bold=True, color=P_GREEN, align=PP_ALIGN.CENTER)
    add_text_box(sl, '/ ¥5,000 上限', Inches(4.0), Inches(2.05), Inches(3.5), Inches(0.25),
                 font_size=PPt(8), color=PRGBColor(0x64,0x74,0x8b), align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(4.0), Inches(2.35), Inches(3.5), Inches(0.12),
             fill_rgb=PRGBColor(0x33,0x41,0x55))
    add_rect(sl, Inches(4.0), Inches(2.35), Inches(1.638), Inches(0.12),
             fill_rgb=P_ORANGE)
    add_text_box(sl, '46.8% (¥2,340/¥5,000)', Inches(4.0), Inches(2.5), Inches(3.5), Inches(0.2),
                 font_size=PPt(7), color=P_ORANGE, align=PP_ALIGN.CENTER)

    # リアルタイム更新ログ
    add_rect(sl, Inches(0.2), Inches(3.5), Inches(7.6), Inches(3.5),
             fill_rgb=PRGBColor(0x0f,0x17,0x2a), line_rgb=P_ORANGE)
    add_text_box(sl, 'リアルタイム監視ログ（毎時更新）', Inches(0.3), Inches(3.52), Inches(7.0), Inches(0.22),
                 font_size=PPt(8), bold=True, color=P_ORANGE)
    logs = [
        ('2026-06-05 10:00', '¥0.32/配送', '¥2,340/月', '✅ 正常'),
        ('2026-06-05 09:00', '¥0.31/配送', '¥2,278/月', '✅ 正常'),
        ('2026-06-05 08:00', '¥0.33/配送', '¥2,216/月', '✅ 正常'),
        ('2026-06-05 07:00', '¥0.30/配送', '¥2,153/月', '✅ 正常'),
        ('2026-06-05 06:00', '¥0.48/配送', '¥2,091/月', '⚠️ 90%超'),
    ]
    headers_log = ['タイムスタンプ', '1配送コスト', '月額累計', 'ステータス']
    t_log = sl.shapes.add_table(len(logs) + 1, 4, Inches(0.3), Inches(3.8), Inches(7.4), Inches(3.0)).table
    for ci, h in enumerate(headers_log):
        t_log.cell(0, ci).text = h
        t_log.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t_log.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)
    for ri, row in enumerate(logs):
        for ci, val in enumerate(row):
            t_log.cell(ri + 1, ci).text = val
            t_log.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    slide_number_tag(sl, 3, TOTAL, is_landscape=IL)

    # ─ S4: AIガバナンス台帳 ─
    sl = blank_slide(prs)
    title_bar(sl, 'AIガバナンス台帳', 'Claude API使用ログ表 / Gate承認履歴 / CEO判断記録 / SHA-256署名確認ボタン', is_landscape=IL)
    footer_bar(sl, FOOTER, is_landscape=IL)

    # Claude API使用ログ
    add_text_box(sl, 'Claude API使用ログ', Inches(0.2), Inches(0.65), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=P_DARK_NAVY)
    api_logs = [
        ('2026-06-05 10:05', 'KPI集計', 'GOV-001', 'G2', 'sha256:3a7f...'),
        ('2026-06-05 09:30', 'ISMS署名生成', 'GOV-002', 'G2', 'sha256:9c2b...'),
        ('2026-06-04 18:00', 'FinOps評価', 'GOV-003', 'G1', 'sha256:4d8e...'),
        ('2026-06-04 12:00', 'Gate承認', 'GOV-004', 'G2', 'sha256:7f1a...'),
    ]
    t1 = sl.shapes.add_table(len(api_logs) + 1, 5, Inches(0.2), Inches(0.9), Inches(7.6), Inches(2.0)).table
    for ci, h in enumerate(['タイムスタンプ', '操作', '機能ID', 'Gate', 'SHA-256']):
        t1.cell(0, ci).text = h
        t1.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t1.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)
    for ri, row in enumerate(api_logs):
        for ci, val in enumerate(row):
            t1.cell(ri + 1, ci).text = val
            t1.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(7.5)

    # Gate承認履歴
    add_text_box(sl, 'Gate承認履歴', Inches(0.2), Inches(3.1), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=P_DARK_NAVY)
    gate_hist = [
        ('G0', '設計完了', '2026-06-05', '松浦CEO', '✅ 承認済'),
        ('G1', 'FinOps監視+bandit CI', '2026-09-30', '（予定）', '⏳ 未着手'),
        ('G2', 'ISMS自動レポート', '2026-11-30', '（予定）', '⏳ 未着手'),
    ]
    t2 = sl.shapes.add_table(len(gate_hist) + 1, 5, Inches(0.2), Inches(3.35), Inches(7.6), Inches(1.5)).table
    for ci, h in enumerate(['Gate', '内容', '予定日', '承認者', '状態']):
        t2.cell(0, ci).text = h
        t2.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t2.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)
    for ri, row in enumerate(gate_hist):
        for ci, val in enumerate(row):
            t2.cell(ri + 1, ci).text = val
            t2.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    # CEO判断記録
    add_text_box(sl, 'CEO判断記録', Inches(0.2), Inches(5.0), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=P_DARK_NAVY)
    ceo_records = [
        ('2026-06-05', 'G0承認・GOV部設計完了', '松浦 学', '承認'),
        ('2026-06-05', 'BRD v1.0承認', '松浦 学', '承認'),
    ]
    t3 = sl.shapes.add_table(len(ceo_records) + 1, 4, Inches(0.2), Inches(5.25), Inches(7.6), Inches(1.0)).table
    for ci, h in enumerate(['日時', '内容', '判断者', '決定']):
        t3.cell(0, ci).text = h
        t3.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t3.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)
    for ri, row in enumerate(ceo_records):
        for ci, val in enumerate(row):
            t3.cell(ri + 1, ci).text = val
            t3.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    # SHA-256署名確認ボタン
    add_rect(sl, Inches(0.2), Inches(6.4), Inches(7.6), Inches(0.45),
             fill_rgb=PRGBColor(0x1e,0x29,0x3b), line_rgb=P_ORANGE)
    add_rect(sl, Inches(0.4), Inches(6.48), Inches(3.5), Inches(0.3),
             fill_rgb=P_DARK_NAVY, line_rgb=P_ORANGE)
    add_text_box(sl, '🔐 SHA-256署名確認', Inches(0.4), Inches(6.5), Inches(3.5), Inches(0.25),
                 font_size=PPt(9), bold=True, color=P_ORANGE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(4.1), Inches(6.48), Inches(3.3), Inches(0.3),
             fill_rgb=P_GREEN, line_rgb=P_GREEN)
    add_text_box(sl, '📥 台帳エクスポート（CSV/PDF）', Inches(4.1), Inches(6.5), Inches(3.3), Inches(0.25),
                 font_size=PPt(9), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    slide_number_tag(sl, 4, TOTAL, is_landscape=IL)

    # ─ S5: LAYOUT_MASTER準拠 ─
    sl = blank_slide(prs)
    title_bar(sl, 'LAYOUT_MASTER準拠', '全数値font-mono tabular-nums tracking-tight / #1A2B4C/#E8A020 / DevSecOps CI/CDバッジ', is_landscape=IL)
    footer_bar(sl, FOOTER, is_landscape=IL)

    # フォント仕様
    add_text_box(sl, 'フォント仕様', Inches(0.2), Inches(0.7), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=P_DARK_NAVY)
    font_rows = [
        ['全KPI・コスト数値', 'font-mono tabular-nums tracking-tight', '¥0.32 / 100% / 0件'],
        ['見出し', 'font-sans font-bold', '画面タイトル・セクション'],
        ['本文', 'font-sans', '説明文・ラベル'],
        ['エラー/警告', 'font-mono font-bold text-red-500', 'bandit検出・FinOps逸脱'],
    ]
    t_font = sl.shapes.add_table(len(font_rows) + 1, 3, Inches(0.2), Inches(0.95), Inches(7.6), Inches(1.8)).table
    for ci, h in enumerate(['適用箇所', 'CSSクラス', '使用例']):
        t_font.cell(0, ci).text = h
        t_font.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t_font.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(font_rows):
        for ci, val in enumerate(row):
            t_font.cell(ri + 1, ci).text = val
            t_font.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    # カラーパレット
    add_text_box(sl, 'カラーパレット（#1A2B4C / #E8A020 ベース）', Inches(0.2), Inches(2.95), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=P_DARK_NAVY)
    colors_def = [
        ('Navy #1A3A5C', P_DARK_NAVY, 'ヘッダー・プライマリ背景'),
        ('Orange #F5A623', P_ORANGE, 'アクション・アクセント・警告'),
        ('White #FFFFFF', P_WHITE, 'テキスト（ダーク背景）'),
        ('LightBG #F0F4F8', P_LIGHT_BG, 'カード背景・セカンダリ'),
        ('Green #059669', P_GREEN, 'OK・適合・bandit 0件'),
        ('Red #EF4444', P_RED, 'エラー・FinOps逸脱・bandit検出'),
    ]
    for i, (name, color, usage) in enumerate(colors_def):
        cx = Inches(0.2) + (i % 3) * Inches(2.55)
        cy = Inches(3.22) + (i // 3) * Inches(0.8)
        add_rect(sl, cx, cy, Inches(0.5), Inches(0.5),
                 fill_rgb=color, line_rgb=P_DARK_NAVY)
        add_text_box(sl, name, cx + Inches(0.55), cy + Inches(0.02), Inches(1.75), Inches(0.22),
                     font_size=PPt(8), bold=True, color=P_DARK_NAVY)
        add_text_box(sl, usage, cx + Inches(0.55), cy + Inches(0.26), Inches(1.75), Inches(0.2),
                     font_size=PPt(7), color=P_SLATE)

    # DevSecOps CI/CDバッジ
    add_text_box(sl, 'DevSecOps CI/CDバッジ表示', Inches(0.2), Inches(5.05), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=P_DARK_NAVY)
    badges = [
        ('bandit', '0 issues', P_GREEN),
        ('CI', 'passing', P_GREEN),
        ('ISMS', '100%', P_GREEN),
        ('FinOps', '¥0/月', P_GREEN),
    ]
    for i, (label, val, color) in enumerate(badges):
        bx_b = Inches(0.3) + i * Inches(1.85)
        add_rect(sl, bx_b, Inches(5.32), Inches(0.8), Inches(0.35),
                 fill_rgb=PRGBColor(0x33,0x41,0x55))
        add_text_box(sl, label, bx_b, Inches(5.37), Inches(0.8), Inches(0.25),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, bx_b + Inches(0.8), Inches(5.32), Inches(0.9), Inches(0.35),
                 fill_rgb=color)
        add_text_box(sl, val, bx_b + Inches(0.8), Inches(5.37), Inches(0.9), Inches(0.25),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    slide_number_tag(sl, 5, TOTAL, is_landscape=IL)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 5. GOV_PHASE_v1.0.pptx
# ═══════════════════════════════════════════════════════════════════

def gen_phase():
    output = f"{OUTPUT_DIR}/GOV_PHASE_v1.0.pptx"
    prs = new_prs_landscape()
    FOOTER = 'GOV部  |  フェーズ計画 v1.0  |  COO経営支援・監査システム部  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # ─ S1: タイトル ─
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, A4_W_LAND, A4_H_LAND, fill_rgb=P_DARK_NAVY)
    add_text_box(sl, 'GOV部 フェーズ計画 v1.0',
                 Inches(1.5), Inches(1.8), Inches(8.5), Inches(1.0),
                 font_size=PPt(36), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'COO経営支援・監査システム部',
                 Inches(1.5), Inches(3.0), Inches(8.5), Inches(0.5),
                 font_size=PPt(18), color=P_ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, '2026-06-05  |  © 2026 株式会社NiceEze  Confidential',
                 Inches(1.5), Inches(3.65), Inches(8.5), Inches(0.35),
                 font_size=PPt(11), color=P_LIGHT_BG, align=PP_ALIGN.CENTER)
    meta = [
        ['文書名', 'GOV部 フェーズ計画 v1.0'],
        ['作成日', '2026-06-05'],
        ['作成者', '自律COO（Claude Code）'],
        ['承認者', '代表取締役CEO 松浦 学'],
    ]
    t_meta = sl.shapes.add_table(4, 2, Inches(3.5), Inches(4.5), Inches(5.0), Inches(1.6)).table
    for ri, (k, v) in enumerate(meta):
        t_meta.cell(ri, 0).text = k
        t_meta.cell(ri, 1).text = v
        for ci in range(2):
            t_meta.cell(ri, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(9)
    slide_number_tag(sl, 1, TOTAL)

    # ─ S2: Gate制概要 ─
    sl = blank_slide(prs)
    title_bar(sl, 'Gate制概要', 'G0〜G4 GOV部開発フェーズゲート')
    footer_bar(sl, FOOTER)

    gates = [
        ('G0', '設計完了', '2026/07末', '詳細設計書完成\nGCP基盤構築\nSecret Manager設定', P_GREEN),
        ('G1', 'FinOps監視\n+bandit CI', '2026/09末', 'FinOps監視稼働\n(¥0.5円アラート)\nbandit CI統合\nDevSecOps統制基盤', P_ORANGE),
        ('G2', 'ISMS自動\nレポート', '2026/11末', 'ISMS月次自動レポート\nSHA-256署名\nBigQuery月次アーカイブ', P_DARK_NAVY),
        ('G3', 'AIガバナンス\n台帳', '2027/01末', 'Claude API使用記録\nGate承認履歴\nCEO判断記録完備', PRGBColor(0x6b, 0x21, 0xa8)),
        ('G4', 'Go-Live\n完全自律監査', '2027/02末', '完全自律監査稼働\n全KPI達成\nCEO最終承認', P_RED),
    ]

    gate_w = Inches(2.0)
    for i, (gate, milestone, date, content, color) in enumerate(gates):
        gx = Inches(0.3) + i * Inches(2.22)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(5.8), fill_rgb=P_LIGHT_BG, line_rgb=color)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(0.85), fill_rgb=color, line_rgb=color)
        add_text_box(sl, gate, gx, Inches(0.72), gate_w, Inches(0.4),
                     font_size=PPt(16), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, milestone, gx, Inches(1.12), gate_w, Inches(0.3),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, date, gx, Inches(1.5), gate_w, Inches(0.22),
                     font_size=PPt(8), color=P_LIGHT_BG, align=PP_ALIGN.CENTER)
        add_text_box(sl, content, gx + Inches(0.1), Inches(1.8), gate_w - Inches(0.2), Inches(4.5),
                     font_size=PPt(9), color=P_DARK_NAVY)

    add_text_box(sl, '◀─── GOV部 開発ロードマップ（2026/07〜2027/02） ───▶',
                 Inches(0.3), Inches(6.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_DARK_NAVY, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 2, TOTAL)

    # ─ S3: G1マイルストーン 2026/09末 ─
    sl = blank_slide(prs)
    title_bar(sl, 'G1マイルストーン — 2026/09末', 'FinOps監視稼働(¥0.5円アラート) / bandit CI統合 / DevSecOps統制基盤')
    footer_bar(sl, FOOTER)

    add_text_box(sl, 'G1 完了条件 / 受け入れ基準', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)

    g1_tasks = [
        ['FinOps監視システム構築（Cloud Functions毎時）', '自律COO', '未着手', '2026/08末', '¥0'],
        ['cost_calculator.py連携・¥0.5円閾値チェック実装', '自律COO', '未着手', '2026/08末', '¥0'],
        ['月額¥5,000アラート設定（Cloud Monitoring）', '自律COO', '未着手', '2026/09末', '¥0'],
        ['bandit -r src/ GitHub Actions CI統合', '自律COO', '未着手', '2026/08末', '¥0'],
        ['DevSecOps統制基盤（lint + テスト自動実行）', '自律COO', '未着手', '2026/09末', '¥0'],
        ['失敗時Slack/チャット通知実装', '自律COO', '未着手', '2026/09末', '¥0'],
        ['Firestore スキーマ設計（AuditLog・FinOpsRecord）', '自律COO', '未着手', '2026/07末', '¥0'],
        ['Secret Manager設定完了確認（Gate 0前提）', '自律COO', '未着手', '2026/07末', '¥0'],
    ]
    t = sl.shapes.add_table(len(g1_tasks) + 1, 5, Inches(0.3), Inches(1.0), Inches(10.8), Inches(4.5)).table
    for ci, h in enumerate(['タスク', '担当', '状態', '期限', '費用']):
        t.cell(0, ci).text = h
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(g1_tasks):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    add_text_box(sl, 'G1完了基準: FinOps監視稼働確認 / bandit 0件CI通過 / GOVシステム運用コスト¥0/月確認',
                 Inches(0.3), Inches(5.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_ORANGE)
    slide_number_tag(sl, 3, TOTAL)

    # ─ S4: G2マイルストーン 2026/11末 ─
    sl = blank_slide(prs)
    title_bar(sl, 'G2マイルストーン — 2026/11末', 'ISMS月次自動レポート / SHA-256署名 / BigQuery月次アーカイブ連携')
    footer_bar(sl, FOOTER)

    add_text_box(sl, 'G2 完了条件 / 受け入れ基準', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)

    g2_tasks = [
        ['ISMS月次自動レポート生成（ISO/IEC 27001:2022）', '自律COO', '未着手', '2026/10末', '¥0'],
        ['SHA-256署名実装（hashlib.sha256 / nosec B303廃止）', '自律COO', '未着手', '2026/10末', '¥0'],
        ['Cloud Scheduler月次トリガー設定', '自律COO', '未着手', '2026/10末', '¥0'],
        ['BigQuery月次アーカイブ連携（ISMSReport テーブル）', '自律COO', '未着手', '2026/11末', '¥0'],
        ['Google Drive保存・CEOチャット通知実装', '自律COO', '未着手', '2026/11末', '¥0'],
        ['KPIダッシュボード API実装（GET /kpi）', '自律COO', '未着手', '2026/10末', '¥0'],
        ['監査ログ改ざん防止確認（Firestore append-only）', '自律COO', '未着手', '2026/11末', '¥0'],
        ['PII最小化・AES-256・RLS設定確認', '自律COO', '未着手', '2026/11末', '¥0'],
    ]
    t2 = sl.shapes.add_table(len(g2_tasks) + 1, 5, Inches(0.3), Inches(1.0), Inches(10.8), Inches(4.5)).table
    for ci, h in enumerate(['タスク', '担当', '状態', '期限', '費用']):
        t2.cell(0, ci).text = h
        t2.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t2.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(g2_tasks):
        for ci, val in enumerate(row):
            t2.cell(ri + 1, ci).text = val
            t2.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    add_text_box(sl, 'G2完了基準: ISMS月次レポート自動生成100% / SHA-256署名確認 / BigQueryアーカイブ疎通確認',
                 Inches(0.3), Inches(5.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_ORANGE)
    slide_number_tag(sl, 4, TOTAL)

    # ─ S5: FinOps計画 ─
    sl = blank_slide(prs)
    title_bar(sl, 'FinOps計画 — GOV部コスト管理', 'GOVシステム¥0/月(無料枠) / 監視対象全システム月額上限¥5,000 / G4以降Claude API監視費用含む')
    footer_bar(sl, FOOTER)

    add_text_box(sl, 'GOV部フェーズ別コスト計画', Inches(0.3), Inches(0.7), Inches(5.5), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)
    finops_rows = [
        ['G0〜G1', '¥0/月', 'Cloud Functions無料枠', 'GOVシステム自体の運用コスト'],
        ['G2（ISMS自動化）', '¥0/月', 'Cloud Functions + BigQuery無料枠', 'ISMS月次レポート含む'],
        ['G3（AIガバナンス）', '¥0/月', 'Firestore無料枠内', 'Claude API監視ログのみ'],
        ['G4（Go-Live）', '¥0/月', '全Cloud Functions無料枠維持', 'Claude API監視費用含む'],
    ]
    t_f = sl.shapes.add_table(len(finops_rows) + 1, 4, Inches(0.3), Inches(1.0), Inches(5.5), Inches(2.2)).table
    for ci, h in enumerate(['フェーズ', 'GOV月額', '内訳', '備考']):
        t_f.cell(0, ci).text = h
        t_f.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t_f.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(9)
    for ri, row in enumerate(finops_rows):
        for ci, val in enumerate(row):
            t_f.cell(ri + 1, ci).text = val
            t_f.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)

    # 監視対象全システム上限
    add_text_box(sl, '監視対象: 全社システム月額上限', Inches(6.0), Inches(0.7), Inches(5.0), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)
    monitor_rows = [
        ['1配送コスト', '¥0.5円/配送', '逸脱 → 即時アラート'],
        ['月額総コスト', '¥5,000/月', '90%超 → 警告アラート'],
        ['GOVシステム', '¥0/月', 'Cloud Functions無料枠維持'],
        ['G4以降Claude API', '月額¥5,000上限内', 'API使用量監視含む'],
    ]
    t_m = sl.shapes.add_table(len(monitor_rows) + 1, 3, Inches(6.0), Inches(1.0), Inches(5.0), Inches(2.2)).table
    for ci, h in enumerate(['監視項目', '上限値', '対応']):
        t_m.cell(0, ci).text = h
        t_m.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t_m.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(9)
    for ri, row in enumerate(monitor_rows):
        for ci, val in enumerate(row):
            t_m.cell(ri + 1, ci).text = val
            t_m.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)

    # FinOps原則まとめ
    add_text_box(sl, 'FinOps原則サマリー', Inches(0.3), Inches(3.4), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)
    principles = [
        ['原則', '内容', '確認方法'],
        ['GOVシステム自体¥0', 'Cloud Functions無料枠内（月200万回以内）で全機能を実現', 'GCP Billing確認（月次）'],
        ['1配送¥0.5円上限', 'cost_calculator.pyによる毎時監視・即時アラート', 'FinOps監視画面リアルタイム確認'],
        ['月額¥5,000上限', '全社システム合算の月額上限（Cloud Monitoring設定）', 'Cloud Monitoringアラート'],
        ['G4以降Claude API含む', 'Claude API監視費用を含めても月額¥5,000上限を維持', 'AIガバナンス台帳コスト記録'],
    ]
    t_p = sl.shapes.add_table(len(principles), 3, Inches(0.3), Inches(3.7), Inches(10.8), Inches(2.8)).table
    for ri, row in enumerate(principles):
        for ci, val in enumerate(row):
            t_p.cell(ri, ci).text = val
            run = t_p.cell(ri, ci).text_frame.paragraphs[0].runs[0]
            run.font.size = PPt(9 if ri == 0 else 8.5)
            if ri == 0:
                run.font.bold = True

    add_text_box(sl, 'FinOps原則: AI自律開発により開発人件費¥0 / GOVシステム運用¥0 / 全社月額¥5,000上限厳格管理',
                 Inches(0.3), Inches(6.65), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_GREEN, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 5, TOTAL)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# メイン
# ═══════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    gen_brd()
    gen_srs()
    gen_seq()
    gen_ui()
    gen_phase()
    print("\n✅ GOV部 5ドキュメント生成完了")
