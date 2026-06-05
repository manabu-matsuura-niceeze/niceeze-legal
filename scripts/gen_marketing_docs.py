"""MARKETING部 5文書生成スクリプト"""
import os

# ─── 共通カラー ───────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, Mm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

DARK_NAVY = RGBColor(0x1a, 0x3a, 0x5c)
ORANGE    = RGBColor(0xf5, 0xa6, 0x23)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)

P_DARK_NAVY = PRGBColor(0x1a, 0x3a, 0x5c)
P_ORANGE    = PRGBColor(0xf5, 0xa6, 0x23)
P_WHITE     = PRGBColor(0xFF, 0xFF, 0xFF)
P_LIGHT_BG  = PRGBColor(0xF0, 0xF4, 0xF8)

OUT_DIR = "docs/MARKETING"

# ═══════════════════════════════════════════════════════════════════
# DOCX ユーティリティ
# ═══════════════════════════════════════════════════════════════════

def set_margins(doc, top=20, bottom=20, left=20, right=20):
    sec = doc.sections[0]
    sec.top_margin    = Mm(top)
    sec.bottom_margin = Mm(bottom)
    sec.left_margin   = Mm(left)
    sec.right_margin  = Mm(right)


def add_heading(doc, text, level=1, color=None):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after  = Pt(2)
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(12) if level == 1 else Pt(10.5)
    run.font.color.rgb = color or DARK_NAVY
    return p


def add_body(doc, text):
    p = doc.add_paragraph(text)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after  = Pt(2)
    for run in p.runs:
        run.font.size = Pt(10.5)
    return p


def set_cell_bg(cell, color_hex):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), color_hex)
    tcPr.append(shd)


def set_table_width(table, width_dxa=9026):
    tbl = table._tbl
    tblPr = tbl.find(qn('w:tblPr'))
    if tblPr is None:
        tblPr = OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    tblW = OxmlElement('w:tblW')
    tblW.set(qn('w:w'), str(width_dxa))
    tblW.set(qn('w:type'), 'dxa')
    tblPr.append(tblW)


def add_full_table(doc, headers, rows, header_bg='1a3a5c', header_fg=WHITE):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    set_table_width(table, 9026)
    table.style = 'Table Grid'

    hdr_row = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr_row.cells[i]
        set_cell_bg(cell, header_bg)
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(h)
        run.bold = True
        run.font.size = Pt(9)
        run.font.color.rgb = header_fg

    for ri, row_data in enumerate(rows):
        tr = table.rows[ri + 1]
        bg = 'F0F4F8' if ri % 2 == 0 else 'FFFFFF'
        for ci, val in enumerate(row_data):
            cell = tr.cells[ci]
            set_cell_bg(cell, bg)
            p = cell.paragraphs[0]
            run = p.add_run(str(val))
            run.font.size = Pt(9)

    return table


def add_footer_docx(doc, text):
    for section in doc.sections:
        footer = section.footer
        p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
        p.clear()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x64, 0x74, 0x8b)


# ═══════════════════════════════════════════════════════════════════
# PPTX ユーティリティ
# ═══════════════════════════════════════════════════════════════════

A4_W_LAND = Inches(11.69)
A4_H_LAND = Inches(8.27)
A4_W = Inches(8.27)
A4_H = Inches(11.69)


def new_prs_landscape():
    prs = Presentation()
    prs.slide_width  = A4_W_LAND
    prs.slide_height = A4_H_LAND
    return prs


def new_prs_portrait():
    prs = Presentation()
    prs.slide_width  = A4_W
    prs.slide_height = A4_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=PPt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=PPt(9), bold=False, color=None,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color or P_DARK_NAVY
    return txBox


def title_slide_cover(slide, title, subtitle='', is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W
    sh = A4_H_LAND if is_landscape else A4_H
    # Full navy background
    add_rect(slide, 0, 0, sw, sh, fill_rgb=P_DARK_NAVY)
    # Orange accent bar
    add_rect(slide, 0, sh * 0.55, sw, PPt(4), fill_rgb=P_ORANGE)
    add_text_box(slide, title, Inches(0.8), sh * 0.2, sw - Inches(1.6), sh * 0.35,
                 font_size=PPt(28), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.8), sh * 0.6, sw - Inches(1.6), sh * 0.15,
                     font_size=PPt(13), color=P_ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(slide, '2026-06-05  |  © 2026 株式会社NiceEze  Confidential',
                 Inches(0.8), sh * 0.75, sw - Inches(1.6), Inches(0.3),
                 font_size=PPt(10), color=P_LIGHT_BG, align=PP_ALIGN.CENTER)


def title_bar(slide, title, subtitle='', is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W
    add_rect(slide, 0, 0, sw, Inches(0.6), fill_rgb=P_DARK_NAVY)
    add_text_box(slide, title, Inches(0.15), Inches(0.05), sw - Inches(0.3), Inches(0.35),
                 font_size=PPt(13), bold=True, color=P_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.15), Inches(0.38), sw - Inches(0.3), Inches(0.2),
                     font_size=PPt(8), color=P_LIGHT_BG, align=PP_ALIGN.LEFT)


def footer_bar(slide, text, is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W
    sh = A4_H_LAND if is_landscape else A4_H
    add_rect(slide, 0, sh - Inches(0.25), sw, Inches(0.25), fill_rgb=P_DARK_NAVY)
    add_text_box(slide, text, Inches(0.1), sh - Inches(0.23), sw - Inches(0.2), Inches(0.2),
                 font_size=PPt(7), color=P_LIGHT_BG, align=PP_ALIGN.CENTER)


def slide_number_tag(slide, num, total, is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W
    sh = A4_H_LAND if is_landscape else A4_H
    add_text_box(slide, f'{num} / {total}',
                 sw - Inches(0.8), sh - Inches(0.23), Inches(0.7), Inches(0.2),
                 font_size=PPt(7), color=P_LIGHT_BG, align=PP_ALIGN.RIGHT)


def pptx_table(slide, headers, rows, x, y, w, h):
    t = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    for ci, hdr in enumerate(headers):
        cell = t.cell(0, ci)
        cell.text = hdr
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.bold = True
        run.font.size = PPt(9)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri + 1, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].font.size = PPt(8.5)
    return t


# ═══════════════════════════════════════════════════════════════════
# 1. MARKETING_BRD_v1.0.docx
# ═══════════════════════════════════════════════════════════════════

def gen_brd():
    output = f"{OUT_DIR}/MARKETING_BRD_v1.0.docx"
    doc = Document()
    set_margins(doc)
    add_footer_docx(doc, 'MARKETING部 ビジネス要件定義書 v1.0  |  © 2026 株式会社NiceEze  Confidential')

    # 表紙
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('MARKETING部 — ビジネス要件定義書 (BRD) v1.0')
    run.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_NAVY

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run('マーケティング自動化システム（SMART-MKT）　v1.0').font.size = Pt(12)
    doc.add_paragraph()

    # 1. 文書管理表
    add_heading(doc, '1. 文書管理表')
    add_full_table(doc,
        ['項目', '内容'],
        [
            ['文書名', 'MARKETING部 ビジネス要件定義書'],
            ['バージョン', 'v1.0'],
            ['作成日', '2026-06-05'],
            ['最終更新日', '2026-06-05'],
            ['作成者', 'NiceEze 自律COO'],
            ['承認者', '代表取締役CEO 松浦 学'],
            ['ステータス', '承認済（Gate 1スプリント適用中）'],
            ['関連文書', 'MARKETING_SRS_v1.0.docx / MARKETING_SEQ_v1.0.pptx'],
        ]
    )

    # 2. ビジネス背景
    add_heading(doc, '2. ビジネス背景（コンテンツ制作工数削減）')
    add_body(doc,
        'EC運営者・マーケティング担当者において、SNS投稿・メルマガ配信・Note記事・YouTube台本の制作が'
        '手作業で行われており、週あたり10〜20時間の工数が発生している。'
        'Google News RSS（8カテゴリ）を活用したニュースクローラーと、'
        '4フォーマット自動生成（X投稿/メルマガ/Note/YouTube台本）により、'
        'コンテンツ生成時間を90%削減し、週2回の安定投稿体制を実現する。'
        '配信ログの自動記録により、投稿履歴の完全トレーサビリティを確保する。'
    )
    add_full_table(doc,
        ['課題', '現状', 'SMART-MKT目標値'],
        [
            ['SNS投稿制作工数', '週5〜10時間（手動）', '週0時間（完全自動）'],
            ['メルマガ作成時間', '1通あたり2〜4時間', '自動生成（5分以内）'],
            ['投稿頻度', '月2〜4回（不定期）', '週2回以上（スケジューラー管理）'],
            ['配信ログ記録', 'なし（属人管理）', '100%自動記録'],
        ]
    )

    # 3. 主要機能表
    add_heading(doc, '3. 主要機能一覧')
    add_full_table(doc,
        ['ID', '機能名', '概要', '優先度'],
        [
            ['MKT-001', 'ニュースクローラー',
             'Google News RSS 8カテゴリ（テクノロジー/ビジネス/経済/EC/マーケティング/AI/スタートアップ/トレンド）から記事を収集。SHA-256でarticle_id生成（64文字）。nosec B310/B314/B405対応。',
             'Must'],
            ['MKT-002', '4フォーマット生成',
             'X投稿（140文字上限）/ メルマガHTML / Note Markdown / YouTube台本の4形式を自動生成。カテゴリ・トーン指定対応。',
             'Must'],
            ['MKT-003', 'コンテンツスケジューラー',
             'Cloud Scheduler cron "0 23 * * *"（朝8:00 JST）/ "0 10 * * *"（夜19:00 JST）の2回自動実行。',
             'Must'],
            ['MKT-004', '配信ログ',
             'SHA-256 id付与 / get_recent(days=7) / summary() / to_json()の4メソッドでログを完全管理。',
             'Must'],
        ]
    )

    # 4. ステークホルダー表
    add_heading(doc, '4. ステークホルダー定義')
    add_full_table(doc,
        ['ステークホルダー', '役割', '主要タッチポイント', '優先度'],
        [
            ['マーケ担当', 'コンテンツ戦略・承認・カテゴリ設定', 'SMART-MKT UI / スケジュール管理画面', 'Must'],
            ['EC運営者', 'コンテンツ活用・SNS運用確認', '4フォーマット表示 / 配信ログ確認', 'Must'],
            ['SNS管理者', 'X投稿・YouTube台本の手動アップロード（G4まで）', 'コピー機能 / YouTube台本ダウンロード', 'Should'],
        ]
    )

    # 5. KPI表
    add_heading(doc, '5. 成功指標（KPI）')
    add_full_table(doc,
        ['KPI', '目標値', '計測方法', 'Gate'],
        [
            ['週2回投稿達成率', '≥90%', 'スケジューラー実行ログ', 'G1'],
            ['コンテンツ生成時間削減率', '90%削減', '生成時間ログ / 工数比較', 'G1'],
            ['配信ログ記録率', '100%', 'DeliveryLog.summary()', '全Gate'],
            ['bandit検出件数', '0件', 'bandit自動スキャン', '全Gate'],
            ['月額インフラコスト（MVP）', '¥0（RSS無料）', 'GCP請求 / 月次', 'G1'],
            ['テスト合格率', '30テスト全Pass', 'unittest自動実行', 'G1'],
        ]
    )

    # 6. Gate制表
    add_heading(doc, '6. Gate制（G0〜G4）/ FinOps')
    add_full_table(doc,
        ['Gate', '完了条件', 'FinOps上限', '備考'],
        [
            ['G0', 'GCP環境構築・Cloud Scheduler設定完了', '¥0（無料枠）', 'インフラ基盤のみ'],
            ['G1', '4フォーマット安定生成・30テスト全Pass・bandit 0件', '¥0/月（RSS無料）', 'Google News RSS無料活用'],
            ['G2', 'RSS品質向上・記事選定ロジック強化', '¥0/月', 'ニュースクローラー精度向上'],
            ['G3', 'Claude API品質向上・YouTube台本精度向上', '¥2,250〜¥4,500/月', 'Claude API統合 月額上限¥5,000'],
            ['G4', 'note.com自動投稿（判断①C）・全SNS自動配信', '¥5,000/月上限', 'note.com G4判断①C適用'],
        ]
    )

    add_body(doc,
        'FinOps詳細: MVP（G1）はGoogle News RSS無料活用により¥0/月。'
        'G3以降はClaude API¥2,250〜¥4,500/月。月額上限¥5,000を厳守。'
        'GCP予算アラートを¥4,000（80%）・¥5,000（100%）に設定。'
    )

    # 7. 制約条件
    add_heading(doc, '7. 制約条件・前提条件')
    add_full_table(doc,
        ['区分', '内容'],
        [
            ['技術制約', 'bandit 0件（B310/B314/B405は nosec コメントで対応）'],
            ['技術制約', 'PII（個人情報）不使用 — ニュース記事のみ処理'],
            ['技術制約', 'モックフォールバック実装必須（RSS障害時）'],
            ['業務判断①C', 'note.comへの自動投稿はG4まで保留（手動投稿導線のみ提供）'],
            ['業務判断②B', 'YouTubeは台本生成のみ。動画アップロードは手動（G4以降検討）'],
            ['前提条件', 'Cloud Scheduler・Cloud Functions設定済み（Gate 0）'],
        ]
    )

    doc.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 2. MARKETING_SRS_v1.0.docx
# ═══════════════════════════════════════════════════════════════════

def gen_srs():
    output = f"{OUT_DIR}/MARKETING_SRS_v1.0.docx"
    doc = Document()
    set_margins(doc)
    add_footer_docx(doc, 'MARKETING部 ソフトウェア要件仕様書 v1.0  |  © 2026 株式会社NiceEze  Confidential')

    # 表紙
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('MARKETING部 — ソフトウェア要件仕様書 (SRS) v1.0')
    run.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_NAVY

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run('マーケティング自動化システム（SMART-MKT）　v1.0').font.size = Pt(12)
    doc.add_paragraph()

    # 1. 文書管理表
    add_heading(doc, '1. 文書管理表')
    add_full_table(doc,
        ['項目', '内容'],
        [
            ['文書名', 'MARKETING部 ソフトウェア要件仕様書'],
            ['バージョン', 'v1.0'],
            ['作成日', '2026-06-05'],
            ['関連BRD', 'MARKETING_BRD_v1.0.docx'],
            ['作成者', 'NiceEze 自律COO'],
        ]
    )

    # 2. システム概要
    add_heading(doc, '2. システム概要')
    add_body(doc,
        'SMART-MKTは Cloud Functions + Google News RSS（無料）+ Cloud Scheduler を基盤とする'
        'マーケティングコンテンツ自動生成システムである。'
        '朝夜2回のスケジューラートリガーにより、8カテゴリのニュースから記事を収集・選定し、'
        '4フォーマット（X/メルマガ/Note/YouTube台本）のコンテンツを自動生成・ログ記録する。'
        'PII不使用・bandit 0件・モックフォールバック実装を必須要件とする。'
    )
    add_full_table(doc,
        ['コンポーネント', '技術', '役割'],
        [
            ['Cloud Functions', 'Python 3.11', 'コンテンツ生成メインロジック'],
            ['Google News RSS', '無料 / 8カテゴリ', 'ニュース記事収集ソース（G1〜G2）'],
            ['Cloud Scheduler', 'cron式', '朝夜2回の定期実行トリガー'],
            ['Claude API', 'anthropic SDK（G3〜）', 'AI品質向上（G3以降）'],
            ['DeliveryLog', 'Python dataclass', '配信ログ管理・JSON出力'],
        ]
    )

    # 3. 機能要件表
    add_heading(doc, '3. 機能要件')

    add_heading(doc, '3.1 NewsCrawler', level=2)
    add_full_table(doc,
        ['項目', '仕様'],
        [
            ['対象カテゴリ数', '8カテゴリ（テクノロジー/ビジネス/経済/EC/マーケティング/AI/スタートアップ/トレンド）'],
            ['article_id生成', 'SHA-256ハッシュ 64文字（URL+タイトルから生成）'],
            ['bandit対応', 'nosec B310（URLオープン）/ nosec B314（XML解析）/ nosec B405（xml.etree）'],
            ['フォールバック', 'RSS障害時はモックデータで動作継続'],
            ['重複排除', 'article_idによる重複チェック実装'],
        ]
    )

    add_heading(doc, '3.2 ContentGenerator（4フォーマット）', level=2)
    add_full_table(doc,
        ['フォーマット', '仕様', '備考'],
        [
            ['X投稿', '140文字上限・ハッシュタグ付与', '文字数カウンター表示'],
            ['メルマガ', 'HTML形式・件名+本文+CTA', 'HTMLプレビュー対応'],
            ['Note記事', 'Markdown形式・見出し+本文+まとめ', 'Markdownレンダリング表示'],
            ['YouTube台本', '構成（導入/本編/まとめ）形式', '手動アップロード（判断②B）'],
        ]
    )

    add_heading(doc, '3.3 ContentScheduler', level=2)
    add_full_table(doc,
        ['スケジュール', 'cron式', '実行時刻（JST）', '実行内容'],
        [
            ['朝実行', '"0 23 * * *"（UTC）', '朝8:00 JST', 'ニュース収集→コンテンツ生成→ログ記録'],
            ['夜実行', '"0 10 * * *"（UTC）', '夜19:00 JST', 'ニュース収集→コンテンツ生成→ログ記録'],
        ]
    )

    add_heading(doc, '3.4 DeliveryLog', level=2)
    add_full_table(doc,
        ['メソッド', '仕様', '説明'],
        [
            ['id生成', 'SHA-256ハッシュ', '配信ID（64文字）を自動付与'],
            ['get_recent(days=7)', '直近N日間のログ取得', 'デフォルト7日'],
            ['summary()', '配信統計サマリー生成', '件数・フォーマット別集計'],
            ['to_json()', 'JSON形式でエクスポート', '監査ログ・外部連携用'],
        ]
    )

    # 4. 非機能要件
    add_heading(doc, '4. 非機能要件')
    add_full_table(doc,
        ['区分', '要件', '計測方法'],
        [
            ['セキュリティ', 'bandit 0件（nosec B310/B314/B405は許容）', 'bandit自動スキャン'],
            ['プライバシー', 'PII（個人情報）一切不使用', 'コードレビュー'],
            ['可用性', 'RSS障害時はモックフォールバックで継続', '統合テスト'],
            ['コスト', 'G1: ¥0/月（RSS無料）/ G3上限: ¥5,000/月', 'GCP請求確認'],
        ]
    )

    # 5. APIエンドポイント表
    add_heading(doc, '5. APIエンドポイント')
    add_full_table(doc,
        ['メソッド', 'エンドポイント', '説明', 'レスポンス'],
        [
            ['GET', '/health', 'ヘルスチェック', '200 OK + {"status":"ok"}'],
            ['POST', '/generate', 'コンテンツ生成（カテゴリ/トーン指定）', '4フォーマットJSON'],
            ['GET', '/log/summary', '配信ログサマリー取得', 'summary JSON'],
            ['POST', '/log/add', '配信ログ追加', '201 Created + log_id'],
        ]
    )

    # 6. テスト要件
    add_heading(doc, '6. テスト要件')
    add_body(doc,
        'unittest による自動テストを30テスト実装し、全Pass（100%合格）を必須とする。'
        'テスト対象: NewsCrawler（RSS解析/重複排除）/ ContentGenerator（4フォーマット生成）/ '
        'ContentScheduler（cron実行）/ DeliveryLog（CRUD/JSON出力）/ APIエンドポイント（正常系/異常系）'
    )
    add_full_table(doc,
        ['テスト区分', 'テスト数', '対象'],
        [
            ['単体テスト', '20件', 'NewsCrawler/ContentGenerator/DeliveryLog各クラス'],
            ['統合テスト', '8件', 'APIエンドポイント・スケジューラー連携'],
            ['セキュリティテスト', '2件', 'bandit 0件確認・PII不使用確認'],
            ['合計', '30件（全Pass必須）', 'CI/CD自動実行'],
        ]
    )

    doc.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 3. MARKETING_SEQ_v1.0.pptx
# ═══════════════════════════════════════════════════════════════════

def gen_seq():
    output = f"{OUT_DIR}/MARKETING_SEQ_v1.0.pptx"
    prs = new_prs_landscape()
    FOOTER = 'MARKETING部  |  シーケンス図 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # S1: タイトル
    sl = blank_slide(prs)
    title_slide_cover(sl, 'MARKETING部 シーケンス図 v1.0', 'SMART-MKT マーケティング自動化システム', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    slide_number_tag(sl, 1, TOTAL, is_landscape=True)

    # S2: スケジューラートリガーフロー
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-01  スケジューラートリガーフロー',
              'Cloud Scheduler → Cloud Functions → NewsCrawler(8カテゴリ) → 記事選定 → ContentGenerator',
              is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    actors = ['Cloud\nScheduler', 'Cloud\nFunctions', 'NewsCrawler\n(8カテゴリ)', '記事選定\nエンジン', 'Content\nGenerator']
    sw = A4_W_LAND
    margin = Inches(0.3)
    actor_w = Inches(1.6)
    n = len(actors)
    spacing = (sw - 2 * margin - actor_w) / max(n - 1, 1)
    actor_x = [margin + i * spacing for i in range(n)]
    actor_y = Inches(0.75)

    for i, actor in enumerate(actors):
        add_rect(sl, actor_x[i], actor_y, actor_w, Inches(0.38), fill_rgb=P_DARK_NAVY, line_rgb=P_ORANGE)
        add_text_box(sl, actor, actor_x[i], actor_y, actor_w, Inches(0.38),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    cx = [ax + actor_w / 2 for ax in actor_x]
    line_y_start = actor_y + Inches(0.38)
    steps = [
        (0, 1, 'cron trigger\n朝8:00/夜19:00 JST'),
        (1, 2, 'crawl() 呼出'),
        (2, 2, 'RSS fetch\n8カテゴリ'),
        (2, 3, '記事リスト返却\n(SHA-256 id付)'),
        (3, 3, 'スコアリング\n重複排除'),
        (3, 4, '選定記事\n渡し'),
        (4, 4, '4フォーマット\n生成'),
        (4, 1, 'コンテンツ\n返却'),
    ]
    step_y = line_y_start + Inches(0.35)
    line_y_end = step_y + len(steps) * Inches(0.7) + Inches(0.3)
    for ci_i in cx:
        line_shape = sl.shapes.add_connector(1, ci_i, line_y_start, ci_i, line_y_end)
        line_shape.line.color.rgb = PRGBColor(0x94, 0xa3, 0xb8)
        line_shape.line.width = PPt(0.5)

    for (fi, ti, label) in steps:
        fx = cx[fi]
        tx = cx[ti]
        arrow = sl.shapes.add_connector(2, fx, step_y, tx, step_y)
        arrow.line.color.rgb = P_ORANGE
        arrow.line.width = PPt(1.5)
        mid_x = min(fx, tx) + abs(tx - fx) / 2 - Inches(0.6) if fi != ti else fx + Inches(0.1)
        add_text_box(sl, label, mid_x, step_y - Inches(0.28), Inches(1.5), Inches(0.25),
                     font_size=PPt(7), color=P_DARK_NAVY, align=PP_ALIGN.CENTER)
        step_y += Inches(0.7)

    slide_number_tag(sl, 2, TOTAL, is_landscape=True)

    # S3: コンテンツ生成フロー
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-02  コンテンツ生成フロー',
              'トピック選定 → 4フォーマット出力(X/メルマガ/Note/YouTube台本) → DeliveryLog記録 → 完了',
              is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    flow_items = [
        ('トピック選定', 'カテゴリ・トーン指定\n記事スコアリング', P_DARK_NAVY),
        ('X投稿生成', '140文字以内\nハッシュタグ付与', P_ORANGE),
        ('メルマガ生成', 'HTML形式\n件名+本文+CTA', P_DARK_NAVY),
        ('Note生成', 'Markdown形式\n見出し+本文', P_DARK_NAVY),
        ('YouTube台本', '構成（導入/本編/まとめ）\n手動アップロード', P_DARK_NAVY),
        ('DeliveryLog', 'SHA-256 id付与\nsummary()記録', PRGBColor(0x05, 0x96, 0x69)),
        ('完了', '全フォーマット\n生成・記録完了', PRGBColor(0x05, 0x96, 0x69)),
    ]
    box_w = Inches(1.45)
    box_h = Inches(1.0)
    arrow_w = Inches(0.25)
    total_items = len(flow_items)
    total_w = total_items * box_w + (total_items - 1) * arrow_w
    start_x = (A4_W_LAND - total_w) / 2
    y_center = Inches(3.5)

    for i, (label, desc, color) in enumerate(flow_items):
        bx = start_x + i * (box_w + arrow_w)
        add_rect(sl, bx, y_center - box_h / 2, box_w, box_h, fill_rgb=color, line_rgb=P_ORANGE)
        add_text_box(sl, label, bx, y_center - box_h / 2 + Inches(0.05), box_w, Inches(0.35),
                     font_size=PPt(8.5), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, desc, bx, y_center - box_h / 2 + Inches(0.4), box_w, Inches(0.55),
                     font_size=PPt(7), color=P_LIGHT_BG, align=PP_ALIGN.CENTER)
        if i < total_items - 1:
            ax = bx + box_w
            ay = y_center
            arrow = sl.shapes.add_connector(2, ax, ay, ax + arrow_w, ay)
            arrow.line.color.rgb = P_ORANGE
            arrow.line.width = PPt(2)

    slide_number_tag(sl, 3, TOTAL, is_landscape=True)

    # S4: 手動生成フロー
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-03  手動生成フロー',
              'SMART-MKT UI → ユーザー入力(カテゴリ/トーン) → POST /generate → 4フォーマット表示 → コピー',
              is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    manual_actors = ['ユーザー\n(マーケ担当)', 'SMART-MKT\nUI', 'Cloud\nFunctions\nAPI', 'Content\nGenerator', 'DeliveryLog']
    n2 = len(manual_actors)
    actor_x2 = [margin + i * spacing for i in range(n2)]

    for i, actor in enumerate(manual_actors):
        add_rect(sl, actor_x2[i], actor_y, actor_w, Inches(0.38), fill_rgb=P_DARK_NAVY, line_rgb=P_ORANGE)
        add_text_box(sl, actor, actor_x2[i], actor_y, actor_w, Inches(0.38),
                     font_size=PPt(7.5), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    cx2 = [ax + actor_w / 2 for ax in actor_x2]
    manual_steps = [
        (0, 1, 'カテゴリ/トーン\n入力'),
        (1, 2, 'POST /generate'),
        (2, 3, 'generate()呼出'),
        (3, 3, '4フォーマット\n生成処理'),
        (3, 2, 'JSON返却'),
        (2, 1, '4フォーマット\n表示'),
        (1, 4, 'log/add記録'),
        (4, 1, 'log_id返却'),
        (1, 0, 'コピー / DL\n完了'),
    ]
    step_y2 = line_y_start + Inches(0.35)
    line_y_end2 = step_y2 + len(manual_steps) * Inches(0.65) + Inches(0.3)
    for ci_i in cx2:
        ls = sl.shapes.add_connector(1, ci_i, line_y_start, ci_i, line_y_end2)
        ls.line.color.rgb = PRGBColor(0x94, 0xa3, 0xb8)
        ls.line.width = PPt(0.5)
    for (fi, ti, label) in manual_steps:
        fx = cx2[fi]
        tx = cx2[ti]
        arrow = sl.shapes.add_connector(2, fx, step_y2, tx, step_y2)
        arrow.line.color.rgb = P_ORANGE
        arrow.line.width = PPt(1.5)
        mid_x = min(fx, tx) + abs(tx - fx) / 2 - Inches(0.6) if fi != ti else fx + Inches(0.1)
        add_text_box(sl, label, mid_x, step_y2 - Inches(0.28), Inches(1.5), Inches(0.25),
                     font_size=PPt(7), color=P_DARK_NAVY, align=PP_ALIGN.CENTER)
        step_y2 += Inches(0.65)

    slide_number_tag(sl, 4, TOTAL, is_landscape=True)

    # S5: エラーハンドリング
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-04  エラーハンドリング',
              'RSS障害→モックフォールバック / 生成失敗→デフォルトテンプレート / スケジューラー失敗→次回自動リトライ',
              is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, '【エラーパターン別フォールバック一覧】',
                 Inches(0.4), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)

    RED = PRGBColor(0xEF, 0x44, 0x44)
    errors = [
        ('RSS障害（Google News RSS接続失敗）',
         'NewsCrawler → Cloud Functions',
         'モックフォールバック起動 → 前回キャッシュ記事で4フォーマット生成継続 / 障害ログ記録 / 次回クロールで自動復旧'),
        ('コンテンツ生成失敗（Generator例外）',
         'ContentGenerator → Cloud Functions',
         'デフォルトテンプレートで代替コンテンツ生成 / エラー詳細をDeliveryLogに記録 / アラート送信'),
        ('スケジューラー失敗（Cloud Scheduler実行失敗）',
         'Cloud Scheduler → Cloud Functions',
         '次回スケジュール（朝/夜）で自動リトライ / 連続3回失敗時にアラート送信 / 手動実行エンドポイント提供'),
    ]
    for i, (title_e, actors_str, fallback) in enumerate(errors):
        y = Inches(1.0) + i * Inches(2.1)
        add_rect(sl, Inches(0.3), y, Inches(10.8), Inches(1.9),
                 fill_rgb=P_LIGHT_BG, line_rgb=RED)
        add_rect(sl, Inches(0.3), y, Inches(10.8), Inches(0.38),
                 fill_rgb=PRGBColor(0x7f, 0x1d, 0x1d), line_rgb=RED)
        add_text_box(sl, f'ERROR: {title_e}',
                     Inches(0.4), y + Inches(0.04), Inches(10.0), Inches(0.28),
                     font_size=PPt(9), bold=True, color=PRGBColor(0xfc, 0xa5, 0xa5))
        add_text_box(sl, f'発生箇所: {actors_str}',
                     Inches(0.4), y + Inches(0.45), Inches(10.0), Inches(0.25),
                     font_size=PPt(8), color=PRGBColor(0x47, 0x55, 0x69))
        add_text_box(sl, f'フォールバック: {fallback}',
                     Inches(0.4), y + Inches(0.72), Inches(10.4), Inches(0.9),
                     font_size=PPt(8.5), color=P_DARK_NAVY, word_wrap=True)

    slide_number_tag(sl, 5, TOTAL, is_landscape=True)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 4. MARKETING_UI_v1.0.pptx
# ═══════════════════════════════════════════════════════════════════

def gen_ui():
    output = f"{OUT_DIR}/MARKETING_UI_v1.0.pptx"
    prs = new_prs_portrait()
    FOOTER = 'MARKETING部  |  UI設計 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # S1: タイトル
    sl = blank_slide(prs)
    title_slide_cover(sl, 'MARKETING部 UI設計 v1.0', 'SMART-MKT マーケティング自動化システム', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)
    slide_number_tag(sl, 1, TOTAL, is_landscape=False)

    # S2: SMART-MKT メイン画面
    sl = blank_slide(prs)
    title_bar(sl, 'SMART-MKT  メイン画面',
              'トレンドニュース表示カード / 8カテゴリ選択チップ / トーン選択(Professional/Casual/Urgent)',
              is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    # カテゴリチップ
    add_text_box(sl, 'カテゴリ選択（8カテゴリ）', Inches(0.2), Inches(0.68), Inches(7.6), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    cats = ['テクノロジー', 'ビジネス', '経済', 'EC', 'マーケ', 'AI', 'スタートアップ', 'トレンド']
    chip_colors = [P_DARK_NAVY, P_ORANGE, P_DARK_NAVY, P_ORANGE, P_DARK_NAVY, P_DARK_NAVY, P_DARK_NAVY, P_ORANGE]
    for i, (cat, cc) in enumerate(zip(cats, chip_colors)):
        cx = Inches(0.2) + (i % 4) * Inches(1.93)
        cy = Inches(0.95) + (i // 4) * Inches(0.38)
        add_rect(sl, cx, cy, Inches(1.75), Inches(0.3), fill_rgb=cc, line_rgb=P_ORANGE)
        add_text_box(sl, cat, cx, cy + Inches(0.04), Inches(1.75), Inches(0.22),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    # トーン選択
    add_text_box(sl, 'トーン選択', Inches(0.2), Inches(1.82), Inches(7.6), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    tones = [('Professional', P_DARK_NAVY), ('Casual', P_ORANGE), ('Urgent', PRGBColor(0xEF, 0x44, 0x44))]
    for i, (tone, tc) in enumerate(tones):
        tx = Inches(0.2) + i * Inches(2.5)
        add_rect(sl, tx, Inches(2.1), Inches(2.3), Inches(0.35), fill_rgb=tc, line_rgb=P_ORANGE)
        add_text_box(sl, tone, tx, Inches(2.13), Inches(2.3), Inches(0.28),
                     font_size=PPt(9), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    # トレンドニュースカード
    add_text_box(sl, 'トレンドニュース（Top 3記事）', Inches(0.2), Inches(2.6), Inches(7.6), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    news_items = [
        ('生成AIがEC業界を変革 — 2026年最新動向', 'AI / テクノロジー', '30分前'),
        ('日本のEC市場 2026年上半期レポート', 'EC / ビジネス', '2時間前'),
        ('TikTokショッピング機能 日本展開開始', 'SNS / マーケ', '4時間前'),
    ]
    for i, (headline, category, time_str) in enumerate(news_items):
        ny = Inches(2.88) + i * Inches(1.1)
        add_rect(sl, Inches(0.2), ny, Inches(7.6), Inches(0.95),
                 fill_rgb=P_LIGHT_BG, line_rgb=P_DARK_NAVY)
        add_rect(sl, Inches(0.2), ny, Inches(0.08), Inches(0.95), fill_rgb=P_ORANGE, line_rgb=P_ORANGE)
        add_text_box(sl, headline, Inches(0.35), ny + Inches(0.07), Inches(6.5), Inches(0.3),
                     font_size=PPt(9.5), bold=True, color=P_DARK_NAVY)
        add_text_box(sl, f'{category}  |  {time_str}', Inches(0.35), ny + Inches(0.42), Inches(5.0), Inches(0.22),
                     font_size=PPt(8), color=PRGBColor(0x47, 0x55, 0x69))
        add_rect(sl, Inches(6.0), ny + Inches(0.55), Inches(1.6), Inches(0.3),
                 fill_rgb=P_DARK_NAVY, line_rgb=P_ORANGE)
        add_text_box(sl, '生成する →', Inches(6.0), ny + Inches(0.58), Inches(1.6), Inches(0.24),
                     font_size=PPt(8), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    slide_number_tag(sl, 2, TOTAL, is_landscape=False)

    # S3: 4フォーマットタブ
    sl = blank_slide(prs)
    title_bar(sl, '4フォーマットタブ',
              'X投稿(140字カウンター) / メルマガHTMLプレビュー / Note Markdown表示 / YouTube台本',
              is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    tabs = ['X投稿', 'メルマガ', 'Note', 'YouTube']
    tab_colors = [P_ORANGE, P_DARK_NAVY, P_DARK_NAVY, P_DARK_NAVY]
    for i, (tab, tc) in enumerate(zip(tabs, tab_colors)):
        tx = Inches(0.2) + i * Inches(1.95)
        add_rect(sl, tx, Inches(0.68), Inches(1.85), Inches(0.35), fill_rgb=tc, line_rgb=P_ORANGE)
        add_text_box(sl, tab, tx, Inches(0.71), Inches(1.85), Inches(0.28),
                     font_size=PPt(9), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)

    # X投稿プレビュー
    add_rect(sl, Inches(0.2), Inches(1.08), Inches(7.6), Inches(2.2),
             fill_rgb=P_LIGHT_BG, line_rgb=P_DARK_NAVY)
    add_text_box(sl, 'X投稿プレビュー', Inches(0.3), Inches(1.1), Inches(4.0), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    add_text_box(sl, '140文字カウンター', Inches(5.5), Inches(1.1), Inches(2.0), Inches(0.22),
                 font_size=PPt(8), color=PRGBColor(0x47, 0x55, 0x69), align=PP_ALIGN.RIGHT)
    add_rect(sl, Inches(6.8), Inches(1.1), Inches(0.8), Inches(0.22),
             fill_rgb=P_ORANGE, line_rgb=P_ORANGE)
    add_text_box(sl, '87/140', Inches(6.8), Inches(1.12), Inches(0.8), Inches(0.18),
                 font_size=PPt(8.5), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl,
                 '【速報】生成AIがEC業界を変革 — 2026年最新動向🤖\n\n'
                 '日本のEC市場でAI活用が急加速。コンテンツ制作工数90%削減事例も。\n\n'
                 '#AI #EC #マーケティング #NiceEze',
                 Inches(0.3), Inches(1.38), Inches(7.2), Inches(0.8),
                 font_size=PPt(9), color=P_DARK_NAVY)

    # メルマガHTMLプレビュー
    add_rect(sl, Inches(0.2), Inches(3.38), Inches(7.6), Inches(1.5),
             fill_rgb=P_LIGHT_BG, line_rgb=P_DARK_NAVY)
    add_text_box(sl, 'メルマガ HTMLプレビュー', Inches(0.3), Inches(3.4), Inches(4.0), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    add_text_box(sl,
                 '件名: 【週刊NiceEze】AIがEC業界を変える — 最新トレンド2026\n'
                 '本文: <h1>今週のトレンド</h1><p>生成AIの活用でEC運営が...</p>\n'
                 'CTA: <a href="#">詳しく見る →</a>',
                 Inches(0.3), Inches(3.68), Inches(7.2), Inches(1.1),
                 font_size=PPt(8.5), color=P_DARK_NAVY)

    # Note/YouTube
    add_rect(sl, Inches(0.2), Inches(4.98), Inches(3.7), Inches(1.4),
             fill_rgb=P_LIGHT_BG, line_rgb=P_DARK_NAVY)
    add_text_box(sl, 'Note Markdown', Inches(0.3), Inches(5.0), Inches(3.5), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    add_text_box(sl, '# AIがEC業界を変える\n## はじめに\nEC業界では...\n## まとめ\n...',
                 Inches(0.3), Inches(5.28), Inches(3.5), Inches(0.95),
                 font_size=PPt(8), color=P_DARK_NAVY)

    add_rect(sl, Inches(4.1), Inches(4.98), Inches(3.7), Inches(1.4),
             fill_rgb=P_LIGHT_BG, line_rgb=P_DARK_NAVY)
    add_text_box(sl, 'YouTube台本（手動アップロード）', Inches(4.2), Inches(5.0), Inches(3.5), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    add_text_box(sl, '【導入】こんにちは！今回は...\n【本編】AIがEC業界にもたらす...\n【まとめ】今週のポイントは...',
                 Inches(4.2), Inches(5.28), Inches(3.5), Inches(0.95),
                 font_size=PPt(8), color=P_DARK_NAVY)

    slide_number_tag(sl, 3, TOTAL, is_landscape=False)

    # S4: スケジュール状態
    sl = blank_slide(prs)
    title_bar(sl, 'スケジュール状態',
              '⏰朝8:00/夜19:00実行状況 / 直近配信ログ一覧 / 配信件数統計',
              is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    # スケジューラー状態
    add_text_box(sl, 'スケジューラー実行状態', Inches(0.2), Inches(0.68), Inches(7.6), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    GREEN = PRGBColor(0x05, 0x96, 0x69)
    sched_items = [
        ('⏰ 朝8:00 (cron "0 23 * * *")', '本日 08:00 実行済', GREEN),
        ('⏰ 夜19:00 (cron "0 10 * * *")', '本日 19:00 予定', P_ORANGE),
    ]
    for i, (label, status, sc) in enumerate(sched_items):
        sy = Inches(0.95) + i * Inches(0.65)
        add_rect(sl, Inches(0.2), sy, Inches(7.6), Inches(0.55), fill_rgb=P_LIGHT_BG, line_rgb=P_DARK_NAVY)
        add_rect(sl, Inches(0.2), sy, Inches(0.12), Inches(0.55), fill_rgb=sc, line_rgb=sc)
        add_text_box(sl, label, Inches(0.4), sy + Inches(0.06), Inches(5.0), Inches(0.22),
                     font_size=PPt(9), bold=True, color=P_DARK_NAVY)
        add_text_box(sl, status, Inches(5.5), sy + Inches(0.06), Inches(2.1), Inches(0.22),
                     font_size=PPt(9), bold=True, color=sc, align=PP_ALIGN.RIGHT)

    # 直近配信ログ
    add_text_box(sl, '直近配信ログ（get_recent(days=7)）', Inches(0.2), Inches(2.35), Inches(7.6), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    log_rows = [
        ['2026-06-05 08:00', 'AI/テクノロジー', 'X+メルマガ+Note', '✅ 完了'],
        ['2026-06-04 19:00', 'EC/ビジネス', 'X+メルマガ+Note+YouTube', '✅ 完了'],
        ['2026-06-04 08:00', 'マーケ/トレンド', 'X+メルマガ', '✅ 完了'],
        ['2026-06-03 19:00', 'AI/スタートアップ', 'X+Note+YouTube', '✅ 完了'],
    ]
    t = sl.shapes.add_table(len(log_rows) + 1, 4, Inches(0.2), Inches(2.62), Inches(7.6), Inches(2.2)).table
    for ci, h in enumerate(['実行日時', 'カテゴリ', '生成フォーマット', 'ステータス']):
        t.cell(0, ci).text = h
        if t.cell(0, ci).text_frame.paragraphs[0].runs:
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(log_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            if t.cell(ri + 1, ci).text_frame.paragraphs[0].runs:
                t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    # 配信件数統計
    add_text_box(sl, '配信件数統計（summary()）', Inches(0.2), Inches(4.95), Inches(7.6), Inches(0.22),
                 font_size=PPt(8.5), bold=True, color=P_DARK_NAVY)
    stats = [('今週配信数', '14回'), ('X投稿', '14件'), ('メルマガ', '12件'), ('Note', '10件'), ('YouTube台本', '6件')]
    for i, (label, val) in enumerate(stats):
        sx = Inches(0.2) + i * Inches(1.52)
        add_rect(sl, sx, Inches(5.22), Inches(1.42), Inches(0.9),
                 fill_rgb=P_DARK_NAVY if i == 0 else P_LIGHT_BG,
                 line_rgb=P_ORANGE)
        add_text_box(sl, val, sx, Inches(5.28), Inches(1.42), Inches(0.4),
                     font_size=PPt(18), bold=True,
                     color=P_WHITE if i == 0 else P_ORANGE, align=PP_ALIGN.CENTER)
        add_text_box(sl, label, sx, Inches(5.7), Inches(1.42), Inches(0.3),
                     font_size=PPt(7.5),
                     color=P_LIGHT_BG if i == 0 else P_DARK_NAVY, align=PP_ALIGN.CENTER)

    slide_number_tag(sl, 4, TOTAL, is_landscape=False)

    # S5: LAYOUT_MASTER準拠
    sl = blank_slide(prs)
    title_bar(sl, 'LAYOUT_MASTER準拠',
              '全数値font-mono tabular-nums / NiceEzeブランドカラー#1A2B4C/#E8A020',
              is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'フォント仕様', Inches(0.2), Inches(0.7), Inches(7.6), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)
    font_rows = [
        ['数値表示全般', 'font-mono tabular-nums tracking-tight', '全KPI・文字数カウンター・配信件数'],
        ['見出し', 'font-sans font-bold', 'タイトル・セクション見出し'],
        ['本文', 'font-sans', '説明文・ラベル・プレビュー'],
        ['エラー表示', 'font-mono font-bold text-red-500', 'バリデーション・アラート'],
    ]
    t = sl.shapes.add_table(len(font_rows) + 1, 3, Inches(0.2), Inches(1.0), Inches(7.6), Inches(1.9)).table
    for ci, h in enumerate(['適用箇所', 'CSSクラス', '備考']):
        t.cell(0, ci).text = h
        if t.cell(0, ci).text_frame.paragraphs[0].runs:
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(font_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            if t.cell(ri + 1, ci).text_frame.paragraphs[0].runs:
                t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    add_text_box(sl, 'ブランドカラーパレット', Inches(0.2), Inches(3.05), Inches(7.6), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)
    colors_def = [
        ('Dark Navy #1A3A5C', P_DARK_NAVY, 'ヘッダー・プライマリ'),
        ('Orange #F5A623', P_ORANGE, 'アクション・アクセント'),
        ('White #FFFFFF', P_WHITE, 'テキスト（ダーク背景）'),
        ('Light #F0F4F8', P_LIGHT_BG, 'カード背景・セカンダリ'),
        ('Brand Navy #1A2B4C', PRGBColor(0x1A, 0x2B, 0x4C), 'NiceEzeブランドカラー'),
        ('Brand Orange #E8A020', PRGBColor(0xE8, 0xA0, 0x20), 'NiceEzeブランドアクセント'),
    ]
    for i, (name, color, usage) in enumerate(colors_def):
        cx = Inches(0.2) + (i % 3) * Inches(2.55)
        cy = Inches(3.35) + (i // 3) * Inches(0.9)
        add_rect(sl, cx, cy, Inches(0.55), Inches(0.55), fill_rgb=color, line_rgb=P_DARK_NAVY)
        add_text_box(sl, name, cx + Inches(0.6), cy + Inches(0.02), Inches(1.75), Inches(0.25),
                     font_size=PPt(8), bold=True, color=P_DARK_NAVY)
        add_text_box(sl, usage, cx + Inches(0.6), cy + Inches(0.28), Inches(1.75), Inches(0.22),
                     font_size=PPt(7), color=PRGBColor(0x47, 0x55, 0x69))

    add_text_box(sl, 'LAYOUT_MASTER準拠チェックリスト',
                 Inches(0.2), Inches(5.25), Inches(7.6), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)
    checks = [
        '✅ 全数値: font-mono tabular-nums tracking-tight 適用',
        '✅ ブランドカラー #1A2B4C / #E8A020 / #1A3A5C / #F5A623 統一',
        '✅ X投稿: 140文字カウンター font-mono で常時表示',
        '✅ 配信件数統計: tabular-nums で桁揃え表示',
    ]
    for i, check in enumerate(checks):
        add_text_box(sl, check, Inches(0.3), Inches(5.55) + i * Inches(0.4), Inches(7.4), Inches(0.35),
                     font_size=PPt(9), color=P_DARK_NAVY)

    slide_number_tag(sl, 5, TOTAL, is_landscape=False)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 5. MARKETING_PHASE_v1.0.pptx
# ═══════════════════════════════════════════════════════════════════

def gen_phase():
    output = f"{OUT_DIR}/MARKETING_PHASE_v1.0.pptx"
    prs = new_prs_landscape()
    FOOTER = 'MARKETING部  |  フェーズ計画 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5
    GREEN = PRGBColor(0x05, 0x96, 0x69)

    # S1: タイトル
    sl = blank_slide(prs)
    title_slide_cover(sl, 'MARKETING部 フェーズ計画 v1.0', 'SMART-MKT マーケティング自動化システム', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    slide_number_tag(sl, 1, TOTAL, is_landscape=True)

    # S2: Gate制概要
    sl = blank_slide(prs)
    title_bar(sl, 'Gate制概要  G0〜G4', 'SMART-MKT 開発フェーズゲート', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    gates = [
        ('G0', '設計', '2026/07末', 'GCP環境構築\nCloud Scheduler設定\nRSSカテゴリ定義', GREEN),
        ('G1', 'テンプレートMVP', '2026/09末', '4フォーマット安定生成\n30テスト全Pass\nスケジューラー稼働\nbandit 0件', P_ORANGE),
        ('G2', 'RSS品質向上', '2026/11末', 'ニュース選定精度向上\n重複排除強化\n記事スコアリング', P_DARK_NAVY),
        ('G3', 'Claude API生成', '2027/01末', 'Claude API品質向上\nYouTube台本精度向上\nnote.com手動投稿導線', PRGBColor(0x6b, 0x21, 0xa8)),
        ('G4', '自動投稿', '2027/03末', 'note.com自動投稿（判断①C）\n全SNS自動配信\nSLA設定', PRGBColor(0xEF, 0x44, 0x44)),
    ]
    gate_w = Inches(2.0)
    for i, (gate, milestone, date, content, color) in enumerate(gates):
        gx = Inches(0.3) + i * Inches(2.22)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(5.8), fill_rgb=P_LIGHT_BG, line_rgb=color)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(0.8), fill_rgb=color, line_rgb=color)
        add_text_box(sl, gate, gx, Inches(0.72), gate_w, Inches(0.35),
                     font_size=PPt(16), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, milestone, gx, Inches(1.08), gate_w, Inches(0.28),
                     font_size=PPt(9), bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, date, gx, Inches(1.4), gate_w, Inches(0.25),
                     font_size=PPt(8), color=P_LIGHT_BG, align=PP_ALIGN.CENTER)
        add_text_box(sl, content, gx + Inches(0.1), Inches(1.7), gate_w - Inches(0.2), Inches(4.5),
                     font_size=PPt(9), color=P_DARK_NAVY)

    add_text_box(sl, '◀─────────────── SMART-MKT 開発ロードマップ（〜2027年3月） ───────────────▶',
                 Inches(0.3), Inches(6.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_DARK_NAVY, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 2, TOTAL, is_landscape=True)

    # S3: G1マイルストーン
    sl = blank_slide(prs)
    title_bar(sl, 'G1マイルストーン — テンプレートMVP（2026/09末）',
              '4フォーマット安定生成 / 30テスト全Pass / スケジューラー稼働 / bandit 0件',
              is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'G1 完了条件 / 受け入れ基準', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)

    g1_tasks = [
        ['NewsCrawler実装（8カテゴリ RSS）', '自律COO', '未着手', '2026/08末', '¥0'],
        ['SHA-256 article_id生成（64文字）', '自律COO', '未着手', '2026/08末', '¥0'],
        ['nosec B310/B314/B405対応', '自律COO', '未着手', '2026/08末', '¥0'],
        ['4フォーマット生成（X/メルマガ/Note/YouTube台本）', '自律COO', '未着手', '2026/09末', '¥0'],
        ['X: 140文字上限カウンター実装', '自律COO', '未着手', '2026/09末', '¥0'],
        ['Cloud Scheduler cron設定（朝夜2回）', '自律COO', '未着手', '2026/09末', '¥0'],
        ['DeliveryLog実装（SHA-256 id/get_recent/summary/to_json）', '自律COO', '未着手', '2026/09末', '¥0'],
        ['unittest 30テスト実装・全Pass確認', '自律COO', '未着手', '2026/09末', '¥0'],
    ]
    t = sl.shapes.add_table(len(g1_tasks) + 1, 5, Inches(0.3), Inches(1.0), Inches(10.8), Inches(4.5)).table
    for ci, h in enumerate(['タスク', '担当', '状態', '期限', '費用']):
        t.cell(0, ci).text = h
        if t.cell(0, ci).text_frame.paragraphs[0].runs:
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(g1_tasks):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            if t.cell(ri + 1, ci).text_frame.paragraphs[0].runs:
                t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    add_text_box(sl,
                 'G1完了基準: 4フォーマット安定生成 / 30テスト全Pass / bandit 0件 / スケジューラー朝夜2回稼働確認',
                 Inches(0.3), Inches(5.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_ORANGE)
    slide_number_tag(sl, 3, TOTAL, is_landscape=True)

    # S4: G3マイルストーン
    sl = blank_slide(prs)
    title_bar(sl, 'G3マイルストーン — Claude API生成（2027/01末）',
              'Claude API品質向上 / note.com手動投稿導線 / YouTube台本精度向上',
              is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'G3 完了条件 / 受け入れ基準', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)

    g3_tasks = [
        ['Claude API統合（コンテンツ品質向上）', '自律COO', '未着手', '2026/12末', '¥2,250〜¥4,500/月'],
        ['X投稿: Claude API活用（より魅力的な文章）', '自律COO', '未着手', '2026/12末', 'API費用に含む'],
        ['メルマガ: パーソナライズ強化', '自律COO', '未着手', '2027/01末', 'API費用に含む'],
        ['YouTube台本: 構成精度向上（導入/本編/まとめ高品質化）', '自律COO', '未着手', '2027/01末', 'API費用に含む'],
        ['note.com手動投稿導線実装（判断①C）', '自律COO', '未着手', '2027/01末', '¥0'],
        ['Note Markdown → note.com用フォーマット変換', '自律COO', '未着手', '2027/01末', '¥0'],
        ['月額コスト監視（上限¥5,000厳守）', '自律COO', '未着手', '継続', '¥0（監視）'],
    ]
    t = sl.shapes.add_table(len(g3_tasks) + 1, 5, Inches(0.3), Inches(1.0), Inches(10.8), Inches(4.2)).table
    for ci, h in enumerate(['タスク', '担当', '状態', '期限', '費用']):
        t.cell(0, ci).text = h
        if t.cell(0, ci).text_frame.paragraphs[0].runs:
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(g3_tasks):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            if t.cell(ri + 1, ci).text_frame.paragraphs[0].runs:
                t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    add_text_box(sl,
                 'G3完了基準: Claude API品質確認 / YouTube台本構成精度向上 / note.com手動投稿導線動作確認 / 月額¥5,000以内',
                 Inches(0.3), Inches(5.4), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=P_ORANGE)
    slide_number_tag(sl, 4, TOTAL, is_landscape=True)

    # S5: FinOps計画
    sl = blank_slide(prs)
    title_bar(sl, 'FinOps計画 — コスト管理',
              'G1:¥0/月 / G3:Claude API¥2,250〜¥4,500/月 / 月額上限¥5,000厳守',
              is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'フェーズ別月次コスト', Inches(0.3), Inches(0.7), Inches(5.5), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)
    finops_rows = [
        ['G0〜G1（テンプレートMVP）', '¥0/月', 'Google News RSS無料', 'AI自律開発・インフラ無料枠'],
        ['G2（RSS品質向上）', '¥0/月', 'Cloud Functions無料枠', 'RSS継続無料'],
        ['G3（Claude API生成）', '¥2,250〜¥4,500/月', 'Claude API費用', '月1,000〜2,000回呼出想定'],
        ['G4（自動投稿）', '¥3,000〜¥5,000/月', '全サービス稼働', '月額上限¥5,000厳守'],
    ]
    t = sl.shapes.add_table(len(finops_rows) + 1, 4, Inches(0.3), Inches(1.0), Inches(5.5), Inches(2.2)).table
    for ci, h in enumerate(['フェーズ', '月額', '内訳', '備考']):
        t.cell(0, ci).text = h
        if t.cell(0, ci).text_frame.paragraphs[0].runs:
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
            t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(9)
    for ri, row in enumerate(finops_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            if t.cell(ri + 1, ci).text_frame.paragraphs[0].runs:
                t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)

    # Claude API詳細
    add_text_box(sl, 'Claude API コスト試算（G3以降）', Inches(6.0), Inches(0.7), Inches(5.0), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)
    api_rows = [
        ['X投稿生成', '月400回', '¥900'],
        ['メルマガ生成', '月300回', '¥675'],
        ['Note記事生成', '月200回', '¥450'],
        ['YouTube台本生成', '月100回', '¥225'],
        ['合計（標準）', '月1,000回', '¥2,250/月'],
        ['上限設定', '月2,000回', '¥4,500/月'],
    ]
    t2 = sl.shapes.add_table(len(api_rows) + 1, 3, Inches(6.0), Inches(1.0), Inches(5.0), Inches(2.2)).table
    for ci, h in enumerate(['用途', '呼出回数', '月額概算']):
        t2.cell(0, ci).text = h
        if t2.cell(0, ci).text_frame.paragraphs[0].runs:
            t2.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
            t2.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(9)
    for ri, row in enumerate(api_rows):
        for ci, val in enumerate(row):
            t2.cell(ri + 1, ci).text = val
            if t2.cell(ri + 1, ci).text_frame.paragraphs[0].runs:
                t2.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)

    # GCPコスト内訳
    add_text_box(sl, 'GCP月次コスト内訳（G4 本番稼働時）', Inches(0.3), Inches(3.4), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=P_DARK_NAVY)
    gcp_rows = [
        ['Cloud Functions（コンテンツ生成）', '¥0〜¥750', '無料枠200万呼出/月内'],
        ['Cloud Scheduler（朝夜2回）', '¥0', '無料枠内（3ジョブまで）'],
        ['Google News RSS', '¥0', '完全無料'],
        ['Claude API（G3以降）', '¥2,250〜¥4,500', '月1,000〜2,000回呼出'],
        ['Cloud Storage（ログ保存）', '¥0〜¥150', '5GB無料枠内'],
        ['月額合計', '¥2,250〜¥5,000', '月額上限¥5,000で管理'],
    ]
    t3 = sl.shapes.add_table(len(gcp_rows) + 1, 3, Inches(0.3), Inches(3.7), Inches(10.8), Inches(2.5)).table
    for ci, h in enumerate(['サービス', '月額概算', '備考']):
        t3.cell(0, ci).text = h
        if t3.cell(0, ci).text_frame.paragraphs[0].runs:
            t3.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
            t3.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(9)
    for ri, row in enumerate(gcp_rows):
        for ci, val in enumerate(row):
            t3.cell(ri + 1, ci).text = val
            if t3.cell(ri + 1, ci).text_frame.paragraphs[0].runs:
                t3.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)

    add_text_box(sl,
                 'FinOps原則: AI自律開発により開発人件費¥0 / G1まで完全¥0 / G3以降Claude API費用のみ / 月額上限¥5,000厳格管理',
                 Inches(0.3), Inches(6.3), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=GREEN)
    slide_number_tag(sl, 5, TOTAL, is_landscape=True)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# メイン
# ═══════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    os.makedirs(OUT_DIR, exist_ok=True)
    gen_brd()
    gen_srs()
    gen_seq()
    gen_ui()
    gen_phase()
    print("\n✅ MARKETING部 5文書 生成完了")
