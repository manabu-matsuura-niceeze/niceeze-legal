"""SBDS SEQ / UI / PHASE .pptx 生成スクリプト"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import math

DARK_NAVY = RGBColor(0x1a, 0x3a, 0x5c)
ORANGE    = RGBColor(0xf5, 0xa6, 0x23)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
LIGHT_BLUE = RGBColor(0xBF, 0xD7, 0xED)
GREEN     = RGBColor(0x05, 0x96, 0x69)
RED       = RGBColor(0xEF, 0x44, 0x44)

A4_W = Inches(8.27)   # 210mm
A4_H = Inches(11.69)  # 297mm
A4_W_LAND = Inches(11.69)
A4_H_LAND = Inches(8.27)


# ─────────────────────────────────────────────────────────────────
# ユーティリティ
# ─────────────────────────────────────────────────────────────────

def new_prs_portrait():
    prs = Presentation()
    prs.slide_width  = A4_W
    prs.slide_height = A4_H
    return prs


def new_prs_landscape():
    prs = Presentation()
    prs.slide_width  = A4_W_LAND
    prs.slide_height = A4_H_LAND
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=Pt(9), bold=False, color=None,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color or DARK_NAVY
    return txBox


def title_bar(slide, title, subtitle='', w=None, is_landscape=False):
    sw = w or (A4_W_LAND if is_landscape else A4_W)
    bar = add_rect(slide, 0, 0, sw, Inches(0.6), fill_rgb=DARK_NAVY)
    add_text_box(slide, title, Inches(0.15), Inches(0.05), sw - Inches(0.3), Inches(0.35),
                 font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.15), Inches(0.38), sw - Inches(0.3), Inches(0.2),
                     font_size=Pt(8), color=LIGHT_BG, align=PP_ALIGN.LEFT)


def footer_bar(slide, text, is_landscape=False):
    sw = A4_W_LAND if is_landscape else A4_W
    sh = A4_H_LAND if is_landscape else A4_H
    add_rect(slide, 0, sh - Inches(0.25), sw, Inches(0.25), fill_rgb=DARK_NAVY)
    add_text_box(slide, text, Inches(0.1), sh - Inches(0.23), sw - Inches(0.2), Inches(0.2),
                 font_size=Pt(7), color=LIGHT_BG, align=PP_ALIGN.CENTER)


def slide_number_tag(slide, num, total, is_landscape=False):
    sw = A4_W_LAND if is_landscape else A4_W
    sh = A4_H_LAND if is_landscape else A4_H
    add_text_box(slide, f'{num} / {total}',
                 sw - Inches(0.8), sh - Inches(0.23), Inches(0.7), Inches(0.2),
                 font_size=Pt(7), color=LIGHT_BG, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────
# SEQ: シーケンス図集
# ─────────────────────────────────────────────────────────────────

def draw_seq_flow(slide, actors, steps, y_start=Inches(0.75), is_landscape=True):
    """シンプルなシーケンス図を描画する"""
    sw = A4_W_LAND if is_landscape else A4_W
    margin = Inches(0.3)
    actor_w = Inches(1.4)
    n = len(actors)
    spacing = (sw - 2 * margin - actor_w) / max(n - 1, 1)

    actor_x = [margin + i * spacing for i in range(n)]
    actor_y = y_start

    # アクターボックス
    for i, actor in enumerate(actors):
        add_rect(slide, actor_x[i], actor_y, actor_w, Inches(0.35),
                 fill_rgb=DARK_NAVY, line_rgb=ORANGE)
        add_text_box(slide, actor,
                     actor_x[i], actor_y, actor_w, Inches(0.35),
                     font_size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ライフライン（縦線）
    from pptx.util import Pt as PtU
    line_y_start = actor_y + Inches(0.35)
    line_y_end = actor_y + Inches(0.35) + Inches(0.45) * (len(steps) + 1)
    cx = [ax + actor_w / 2 for ax in actor_x]

    for cx_i in cx:
        line_shape = slide.shapes.add_connector(1, cx_i, line_y_start, cx_i, line_y_end)
        line_shape.line.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
        line_shape.line.width = Pt(0.5)

    # ステップ（矢印 + ラベル）
    step_y = line_y_start + Inches(0.3)
    for step in steps:
        from_idx, to_idx, label, is_async, note = step
        fx = cx[from_idx]
        tx = cx[to_idx]
        arrow = slide.shapes.add_connector(2, fx, step_y, tx, step_y)
        arrow.line.color.rgb = ORANGE if not is_async else RGBColor(0x6b, 0x72, 0x80)
        arrow.line.width = Pt(1.5 if not is_async else 1.0)
        mid_x = min(fx, tx) + abs(tx - fx) / 2 - Inches(0.6)
        add_text_box(slide, label, mid_x, step_y - Inches(0.2), Inches(1.2), Inches(0.18),
                     font_size=Pt(7), color=DARK_NAVY, align=PP_ALIGN.CENTER)
        if note:
            note_x = max(fx, tx) + Inches(0.05)
            add_text_box(slide, note, note_x, step_y - Inches(0.15), Inches(1.5), Inches(0.2),
                         font_size=Pt(6.5), color=RED, align=PP_ALIGN.LEFT)
        step_y += Inches(0.42)


def gen_seq(output='docs/SBDS/SBDS_SEQ_v1.0.pptx'):
    prs = new_prs_landscape()
    FOOTER = 'SBDS  |  シーケンス図集 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    actors_main = ['配送スタッフ', 'TMS-DRV-001\n(PWA/LIFF)', 'Cloud Run\nAPI', 'Firestore\nDB', 'LINE API']

    # ─ Slide 1: 目次 ─
    sl = blank_slide(prs)
    title_bar(sl, 'SBDS シーケンス図集 v1.0 — 目次', 'TMS-SET-001 / TMS-DRV-001', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    add_text_box(sl, '収録フロー一覧', Inches(0.4), Inches(0.75), Inches(10.5), Inches(0.3),
                 font_size=Pt(11), bold=True, color=DARK_NAVY)
    items = [
        ('Slide 2', '正常系メインフロー', '配送スタッフ操作→TMS-DRV-001→Firestore→完了登録'),
        ('Slide 3', '異常系・エラーハンドリング', '労働法ロック / 冷凍警告 / 1分前通知失敗'),
        ('Slide 4', 'LINE Webhook連携フロー', 'LINE→Cloud Run→Redis Consumer Group→Firestore'),
        ('Slide 5', 'DXF/AR 計測フロー（TMS-SET-001）', 'DXFインポート / WebXR Hit Test → IndexedDB v142'),
    ]
    for i, (slide_no, title, desc) in enumerate(items):
        y = Inches(1.2) + i * Inches(0.9)
        add_rect(sl, Inches(0.4), y, Inches(10.5), Inches(0.75),
                 fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
        add_text_box(sl, slide_no, Inches(0.5), y + Inches(0.05), Inches(0.9), Inches(0.3),
                     font_size=Pt(9), bold=True, color=ORANGE)
        add_text_box(sl, title, Inches(1.4), y + Inches(0.05), Inches(4.0), Inches(0.3),
                     font_size=Pt(10), bold=True, color=DARK_NAVY)
        add_text_box(sl, desc, Inches(1.4), y + Inches(0.38), Inches(9.0), Inches(0.3),
                     font_size=Pt(8), color=RGBColor(0x47, 0x55, 0x69))
    slide_number_tag(sl, 1, TOTAL, is_landscape=True)

    # ─ Slide 2: 正常系メインフロー ─
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-01  正常系メインフロー', '配送スタッフ → TMS-DRV-001 → Firestore → 完了登録', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    steps = [
        (0, 1, 'アプリ起動', False, ''),
        (1, 2, 'GET /packages', False, ''),
        (2, 3, 'query packages', False, ''),
        (3, 2, 'return []', False, ''),
        (2, 1, 'IndexedDB v142保存', True, '0.7秒以下'),
        (1, 0, '配送リスト表示', False, 'ルーティングソート済'),
        (0, 1, '配送完了タップ', False, ''),
        (1, 2, 'POST /complete', False, ''),
        (2, 3, 'update status=COMPLETED', False, ''),
        (3, 2, 'ack', False, ''),
        (2, 4, 'LINE PUSH（次配送先）', True, '1分前のみ'),
        (2, 1, '完了確認', False, ''),
        (1, 0, 'UI更新', False, '完了バッジ表示'),
    ]
    draw_seq_flow(sl, actors_main, steps, is_landscape=True)
    slide_number_tag(sl, 2, TOTAL, is_landscape=True)

    # ─ Slide 3: 異常系 ─
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-02  異常系・エラーハンドリングフロー', '労働法ロック / 冷凍警告 / PULL通知失敗', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    add_text_box(sl, '【Alt: 4時間超労働法ロック】', Inches(0.4), Inches(0.7), Inches(5.0), Inches(0.25),
                 font_size=Pt(9), bold=True, color=RED)
    steps_err = [
        (0, 1, '操作（4h経過）', False, ''),
        (1, 1, 'check_labor_law()', False, 'elapsed≥240min'),
        (1, 0, 'STATUS_LOCKED_BY_LABOR_LAW', False, '全入力無効・バナー表示'),
        (0, 1, '冷凍荷物スキャン', False, ''),
        (1, 0, '赤字警告表示', False, '手渡し必須'),
        (1, 2, 'POST /notify/pull（1分前）', True, ''),
        (2, 4, 'LINE Push', True, ''),
        (4, 2, 'Error 429', False, 'レート制限'),
        (2, 3, 'log error', False, 'Firestore audit_log'),
        (2, 1, 'retry 3回', True, 'exponential backoff'),
    ]
    draw_seq_flow(sl, actors_main, steps_err, y_start=Inches(1.05), is_landscape=True)
    slide_number_tag(sl, 3, TOTAL, is_landscape=True)

    # ─ Slide 4: LINE Webhook ─
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-03  LINE Webhook連携フロー', 'Redis Streams Consumer Group によるデデュープ', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    actors_line = ['LINE Platform', 'Cloud Run\nWebhook', 'Redis\nStreams', 'Firestore\nDB', 'TMS-DRV-001']
    steps_line = [
        (0, 1, 'POST /webhook\n(X-Line-Signature)', False, 'HMAC-SHA256検証'),
        (1, 1, 'verify signature', False, ''),
        (1, 2, 'XADD events stream', False, ''),
        (2, 1, 'XREADGROUP consumer1', False, 'Consumer Group'),
        (1, 1, 'check duplicate', False, '処理済みIDチェック'),
        (1, 3, 'write package', False, ''),
        (3, 1, 'ack', False, ''),
        (1, 2, 'XACK', False, '処理完了'),
        (1, 4, 'SSE push', True, 'リアルタイム更新'),
        (4, 4, 'IndexedDB v142更新', False, '0.7秒以下'),
    ]
    draw_seq_flow(sl, actors_line, steps_line, is_landscape=True)
    slide_number_tag(sl, 4, TOTAL, is_landscape=True)

    # ─ Slide 5: DXF/AR ─
    sl = blank_slide(prs)
    title_bar(sl, 'SEQ-04  DXFインポート / AR計測フロー（TMS-SET-001）',
              'dxf-parser.js（ブラウザ完結）/ WebXR Hit Test', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    actors_dxf = ['管理者', 'TMS-SET-001\n(Browser)', 'dxf-parser.js', 'IndexedDB\nv142', 'Firestore']
    steps_dxf = [
        (0, 1, 'DXFファイル選択', False, ''),
        (1, 2, 'parse(dxfText)', False, 'FileReader API'),
        (2, 2, 'ENTITIES解析\nPOLYLINE→面積計算', False, 'Shoelace公式'),
        (2, 1, 'return RoomRecord[]', False, ''),
        (1, 0, 'プレビュー表示', False, '手動修正可能'),
        (0, 1, '保存ボタン', False, ''),
        (1, 3, 'idbPut(building_masters)', False, '_idb_version:142'),
        (1, 4, 'POST /buildings', True, 'Cloud Run経由'),
        (4, 1, 'ack', False, ''),
        (1, 0, 'Toast通知「保存しました」', False, ''),
    ]
    draw_seq_flow(sl, actors_dxf, steps_dxf, is_landscape=True)
    add_text_box(sl, '※ AR計測（WebXR Hit Test）は判断③（iOS/Android）CEO決定後に追加',
                 Inches(0.4), A4_H_LAND - Inches(0.5), Inches(10.0), Inches(0.25),
                 font_size=Pt(8), color=RED)
    slide_number_tag(sl, 5, TOTAL, is_landscape=True)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ─────────────────────────────────────────────────────────────────
# UI: 画面設計書
# ─────────────────────────────────────────────────────────────────

def draw_wireframe_box(slide, x, y, w, h, label, fill=LIGHT_BG, border=DARK_NAVY, label_size=Pt(8)):
    add_rect(slide, x, y, w, h, fill_rgb=fill, line_rgb=border)
    add_text_box(slide, label, x + Inches(0.05), y + Inches(0.05),
                 w - Inches(0.1), h - Inches(0.1),
                 font_size=label_size, color=DARK_NAVY, word_wrap=True)


def gen_ui(output='docs/SBDS/SBDS_UI_v1.0.pptx'):
    prs = new_prs_portrait()
    FOOTER = 'SBDS  |  画面設計書 v1.0  |  © 2026 株式会社NiceEze  Confidential  ※Gemini関与なし'
    TOTAL = 4

    # ─ Slide 1: 画面遷移図（サマリー）─
    sl = blank_slide(prs)
    title_bar(sl, 'SBDS 画面設計書 v1.0 — 画面遷移図（サマリー）', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    # 遷移図
    boxes = [
        (Inches(0.4), Inches(1.0), Inches(3.2), Inches(1.2), 'TMS-SET-001\n初期設定・マスタ管理\n（管理者）'),
        (Inches(4.4), Inches(1.0), Inches(3.2), Inches(1.2), 'TMS-DRV-001\n配送員スマホ画面\n（配送スタッフ）'),
        (Inches(0.4), Inches(3.2), Inches(3.2), Inches(0.8), 'DXFインポート\n（モーダル）'),
        (Inches(0.4), Inches(4.4), Inches(3.2), Inches(0.8), 'AR計測\n（WebXRセッション）'),
        (Inches(4.4), Inches(3.2), Inches(3.2), Inches(0.8), '完了確認\n（ダイアログ）'),
        (Inches(4.4), Inches(4.4), Inches(3.2), Inches(0.8), 'ロック画面\nSTATUS_LOCKED_BY_LABOR_LAW'),
    ]
    for bx, by, bw, bh, label in boxes:
        add_rect(sl, bx, by, bw, bh, fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
        add_text_box(sl, label, bx + Inches(0.1), by + Inches(0.1),
                     bw - Inches(0.2), bh - Inches(0.2),
                     font_size=Pt(9), bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)

    # 矢印（接続線）
    arrows = [
        (Inches(2.0), Inches(2.2), Inches(2.0), Inches(3.2)),  # SET-001 → DXF
        (Inches(2.0), Inches(2.2), Inches(2.0), Inches(4.4)),  # SET-001 → AR
        (Inches(6.0), Inches(2.2), Inches(6.0), Inches(3.2)),  # DRV-001 → 完了
        (Inches(6.0), Inches(2.2), Inches(6.0), Inches(4.4)),  # DRV-001 → ロック
    ]
    for ax, ay, bx2, by2 in arrows:
        conn = sl.shapes.add_connector(1, ax, ay, bx2, by2)
        conn.line.color.rgb = ORANGE
        conn.line.width = Pt(1.5)

    add_text_box(sl, '画面ID: TMS-SET-001', Inches(0.4), Inches(0.8), Inches(3.2), Inches(0.2),
                 font_size=Pt(7), color=RGBColor(0x47, 0x55, 0x69))
    add_text_box(sl, '画面ID: TMS-DRV-001', Inches(4.4), Inches(0.8), Inches(3.2), Inches(0.2),
                 font_size=Pt(7), color=RGBColor(0x47, 0x55, 0x69))
    slide_number_tag(sl, 1, TOTAL)

    # ─ Slide 2: TMS-SET-001 ─
    sl = blank_slide(prs)
    title_bar(sl, 'TMS-SET-001  初期設定・マスタ管理画面', '対象ユーザー: 管理者  |  アクセス権: 管理者ロール', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    # ゾーン説明
    add_text_box(sl, 'LAYOUT_MASTER.md 準拠 3ゾーン構成', Inches(0.2), Inches(0.65), Inches(7.5), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)

    # 上部ゾーン
    add_rect(sl, Inches(0.2), Inches(0.9), Inches(7.6), Inches(1.8),
             fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=ORANGE)
    add_text_box(sl, '上部: 建物基本スペック入力エリア', Inches(0.3), Inches(0.92), Inches(7.0), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)
    input_fields = [
        ('棟数 (1-20)', Inches(0.3), Inches(1.15)),
        ('階数 (1-100)', Inches(2.2), Inches(1.15)),
        ('居住者用EV (0-20基)', Inches(4.1), Inches(1.15)),
        ('業務用EV (最低4基)', Inches(0.3), Inches(1.7)),
    ]
    for label, fx, fy in input_fields:
        add_rect(sl, fx, fy, Inches(1.7), Inches(0.45),
                 fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=RGBColor(0x33, 0x41, 0x55))
        add_text_box(sl, label, fx + Inches(0.05), fy + Inches(0.05), Inches(1.6), Inches(0.15),
                     font_size=Pt(6.5), color=RGBColor(0x64, 0x74, 0x8b))
        add_text_box(sl, 'font-mono tabular-nums tracking-tight',
                     fx + Inches(0.05), fy + Inches(0.22), Inches(1.6), Inches(0.18),
                     font_size=Pt(6), color=ORANGE)
    add_rect(sl, Inches(6.0), Inches(1.15), Inches(1.6), Inches(0.45),
             fill_rgb=RGBColor(0x1d, 0x4e, 0xd8), line_rgb=RGBColor(0x25, 0x63, 0xeb))
    add_text_box(sl, 'CAD/DXFインポート', Inches(6.05), Inches(1.25), Inches(1.5), Inches(0.25),
                 font_size=Pt(7.5), color=WHITE, align=PP_ALIGN.CENTER)

    # 中央ゾーン（グリッド）
    add_rect(sl, Inches(0.2), Inches(2.75), Inches(7.6), Inches(5.0),
             fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=ORANGE)
    add_text_box(sl, '中央: フロアスプレッドシート型グリッドエディタ', Inches(0.3), Inches(2.77), Inches(6.0), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)
    headers = ['棟名', '部屋番号', '専有面積(㎡)', '家賃(円)', 'EV出口距離(m)', '階']
    col_w = Inches(1.18)
    for i, h in enumerate(headers):
        add_rect(sl, Inches(0.3) + i * col_w, Inches(3.0), col_w - Inches(0.02), Inches(0.28),
                 fill_rgb=DARK_NAVY, line_rgb=RGBColor(0x33, 0x41, 0x55))
        add_text_box(sl, h, Inches(0.32) + i * col_w, Inches(3.02), col_w - Inches(0.06), Inches(0.24),
                     font_size=Pt(7), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for row in range(3):
        for i in range(6):
            add_rect(sl, Inches(0.3) + i * col_w, Inches(3.3) + row * Inches(0.35),
                     col_w - Inches(0.02), Inches(0.32),
                     fill_rgb=RGBColor(0x1e, 0x29, 0x3b) if row % 2 == 0 else RGBColor(0x0f, 0x17, 0x2a),
                     line_rgb=RGBColor(0x1e, 0x29, 0x3b))
            sample = ['A棟', '0101', '25.00', '80,000', '15.0', '1'][i]
            add_text_box(sl, sample, Inches(0.32) + i * col_w,
                         Inches(3.32) + row * Inches(0.35), col_w - Inches(0.06), Inches(0.28),
                         font_size=Pt(8), color=RGBColor(0xe2, 0xe8, 0xf0))
    add_text_box(sl, '※ 専有面積・家賃・EV距離は font-mono tabular-nums tracking-tight 必須',
                 Inches(0.3), Inches(4.4), Inches(7.0), Inches(0.2),
                 font_size=Pt(7), color=ORANGE)

    # 下部ゾーン
    add_rect(sl, Inches(0.2), Inches(7.8), Inches(7.6), Inches(0.7),
             fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=ORANGE)
    add_text_box(sl, '下部: 保存 / リセット', Inches(0.3), Inches(7.82), Inches(3.0), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)
    for label, color, ox in [('保存', GREEN, Inches(5.0)), ('リセット', RED, Inches(6.3))]:
        add_rect(sl, ox, Inches(7.88), Inches(1.1), Inches(0.45), fill_rgb=color, line_rgb=color)
        add_text_box(sl, label, ox, Inches(7.94), Inches(1.1), Inches(0.3),
                     font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number_tag(sl, 2, TOTAL)

    # ─ Slide 3: TMS-DRV-001 ─
    sl = blank_slide(prs)
    title_bar(sl, 'TMS-DRV-001  配送員専用スマホ画面', '対象ユーザー: 配送スタッフ  |  アクセス権: 認証済みスタッフ', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'LAYOUT_MASTER.md 準拠 3ゾーン構成（スマホ最適化）', Inches(0.2), Inches(0.65), Inches(7.5), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)

    # 上部（固定ヘッダー）
    add_rect(sl, Inches(0.2), Inches(0.9), Inches(7.6), Inches(1.6),
             fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=ORANGE)
    add_text_box(sl, '上部（固定ヘッダー）', Inches(0.3), Inches(0.92), Inches(4.0), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)
    stats = [
        ('本日配送総数', '12件', Inches(0.3)),
        ('現在フロア', '2F', Inches(2.8)),
        ('完了 / 未完了', '3 / 9', Inches(5.1)),
    ]
    for label, val, sx in stats:
        add_text_box(sl, label, sx, Inches(1.15), Inches(2.2), Inches(0.2),
                     font_size=Pt(7), color=RGBColor(0x64, 0x74, 0x8b))
        add_text_box(sl, val, sx, Inches(1.35), Inches(2.2), Inches(0.35),
                     font_size=Pt(16), bold=True, color=WHITE)
    add_text_box(sl, 'font-mono tabular-nums tracking-tight（3統計値すべて）',
                 Inches(0.3), Inches(1.75), Inches(7.0), Inches(0.18),
                 font_size=Pt(6.5), color=ORANGE)

    # 労働法ロックバナー
    add_rect(sl, Inches(0.2), Inches(2.55), Inches(7.6), Inches(0.5),
             fill_rgb=RGBColor(0x7f, 0x1d, 0x1d), line_rgb=RED)
    add_text_box(sl, '⚠ STATUS_LOCKED_BY_LABOR_LAW — 4時間連続作業のため休憩必須（全入力無効）',
                 Inches(0.3), Inches(2.62), Inches(7.4), Inches(0.3),
                 font_size=Pt(8), bold=True, color=RGBColor(0xfc, 0xa5, 0xa5), align=PP_ALIGN.CENTER)

    # 中央（配送リスト）
    add_rect(sl, Inches(0.2), Inches(3.1), Inches(7.6), Inches(4.2),
             fill_rgb=RGBColor(0x0f, 0x17, 0x2a), line_rgb=ORANGE)
    add_text_box(sl, '中央: 配送リスト（最適順路ソート済み）', Inches(0.3), Inches(3.12), Inches(6.0), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)

    pkg_samples = [
        ('1F-0101', '通常', RGBColor(0x1e, 0x29, 0x3b), ''),
        ('2F-0201', '冷凍', RGBColor(0x1a, 0x0a, 0x0a), '【ロッカー格納禁止：手渡し必須】'),
        ('2F-0203', '冷蔵', RGBColor(0x1a, 0x0a, 0x0a), '【ロッカー格納禁止：手渡し必須】'),
        ('2F-0205', '要注意', RGBColor(0x1e, 0x29, 0x3b), '（クレーム要注意 / 同フロア最後）'),
    ]
    for i, (room, flag, bg, note) in enumerate(pkg_samples):
        py = Inches(3.4) + i * Inches(0.9)
        add_rect(sl, Inches(0.3), py, Inches(7.4), Inches(0.82), fill_rgb=bg,
                 line_rgb=RED if '禁止' in note else DARK_NAVY)
        add_text_box(sl, room, Inches(0.4), py + Inches(0.05), Inches(1.5), Inches(0.35),
                     font_size=Pt(14), bold=True, color=WHITE)
        flag_color = RED if flag in ('冷凍', '冷蔵') else ORANGE
        add_rect(sl, Inches(2.0), py + Inches(0.08), Inches(0.8), Inches(0.25),
                 fill_rgb=flag_color, line_rgb=flag_color)
        add_text_box(sl, flag, Inches(2.0), py + Inches(0.1), Inches(0.8), Inches(0.2),
                     font_size=Pt(7), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if note:
            add_text_box(sl, note, Inches(0.4), py + Inches(0.4), Inches(7.0), Inches(0.25),
                         font_size=Pt(7.5), bold=True, color=RED)

    # 下部（動線マップ）
    add_rect(sl, Inches(0.2), Inches(7.35), Inches(7.6), Inches(1.15),
             fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=ORANGE)
    add_text_box(sl, '下部: 動線マップ', Inches(0.3), Inches(7.37), Inches(3.0), Inches(0.2),
                 font_size=Pt(8), bold=True, color=ORANGE)
    add_text_box(sl, 'EVホール → 廊下(8m) → 2F-0201', Inches(0.3), Inches(7.6), Inches(5.0), Inches(0.25),
                 font_size=Pt(9), color=RGBColor(0x60, 0xa5, 0xfa))
    add_text_box(sl, '推定到着残り時間: 1分23秒',
                 Inches(0.3), Inches(7.88), Inches(7.0), Inches(0.3),
                 font_size=Pt(11), bold=True, color=RGBColor(0x34, 0xd3, 0x99))
    add_text_box(sl, 'font-mono tabular-nums tracking-tight',
                 Inches(0.3), Inches(8.2), Inches(7.0), Inches(0.18),
                 font_size=Pt(6.5), color=ORANGE)
    slide_number_tag(sl, 3, TOTAL)

    # ─ Slide 4: バリデーション・エラー仕様 ─
    sl = blank_slide(prs)
    title_bar(sl, '入力バリデーション・エラー表示仕様', 'TMS-SET-001 / TMS-DRV-001 共通', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'TMS-SET-001 入力フィールド仕様', Inches(0.2), Inches(0.7), Inches(7.6), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    rows_set = [
        ['棟数', 'number', '1〜20', '必須', '範囲外: 赤枠 + エラーメッセージ'],
        ['階数', 'number', '1〜100', '必須', '範囲外: 赤枠 + エラーメッセージ'],
        ['居住者用EV', 'number', '0〜20', '必須', '範囲外: 赤枠'],
        ['業務用EV', 'number', '最低4', '必須', '4未満: 赤枠「最低4基必須」'],
        ['部屋番号', 'text', r'^\d{3,4}[A-Za-z]?$', '必須', '正規表現不一致: 赤枠'],
        ['専有面積', 'number', '>0', '必須', '0以下: 赤枠'],
        ['家賃', 'number', '≥0', '必須', '負値: 赤枠'],
        ['EV出口距離', 'number', '≥0', '必須', '負値: 赤枠'],
    ]
    from pptx.util import Inches as I2
    table = sl.shapes.add_table(len(rows_set) + 1, 5,
                                 I2(0.2), I2(1.0), I2(7.6), I2(3.0)).table
    headers_v = ['フィールド', '型', 'バリデーション', '必須/任意', 'エラー表示']
    for i, h in enumerate(headers_v):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    for ri, row in enumerate(rows_set):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(7.5)

    add_text_box(sl, 'TMS-DRV-001 警告バナー仕様', Inches(0.2), Inches(4.15), Inches(7.6), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    banners = [
        ('冷凍/冷蔵警告', '赤背景ボーダー + 赤字大文字【ロッカー格納禁止：手渡し必須】', RED),
        ('STATUS_LOCKED_BY_LABOR_LAW', '赤ダーク背景バナー全幅表示 + 全入力pointer-events:none', RED),
        ('クレーム要注意', 'オレンジフラグ + 同フロア最後尾ソート', ORANGE),
        ('1分前PULL通知', 'ETA残り60秒でLINE PUSH発火（UIは青字表示）', RGBColor(0x60, 0xa5, 0xfa)),
    ]
    for i, (name, desc, color) in enumerate(banners):
        y = Inches(4.45) + i * Inches(0.7)
        add_rect(sl, Inches(0.2), y, Inches(7.6), Inches(0.6),
                 fill_rgb=RGBColor(0xf0, 0xf4, 0xf8), line_rgb=color)
        add_text_box(sl, name, Inches(0.3), y + Inches(0.05), Inches(2.5), Inches(0.25),
                     font_size=Pt(9), bold=True, color=color)
        add_text_box(sl, desc, Inches(0.3), y + Inches(0.3), Inches(7.2), Inches(0.25),
                     font_size=Pt(8), color=DARK_NAVY)

    slide_number_tag(sl, 4, TOTAL)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ─────────────────────────────────────────────────────────────────
# PHASE: フェーズ管理書
# ─────────────────────────────────────────────────────────────────

def gen_phase(output='docs/SBDS/SBDS_PHASE_v1.0.pptx'):
    prs = new_prs_landscape()
    FOOTER = 'SBDS  |  フェーズ管理書 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # Slide 1: 文書管理・開発基本方針
    sl = blank_slide(prs)
    title_bar(sl, 'SBDS フェーズ管理書 v1.0 — 文書管理・開発基本方針', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, '文書管理表', Inches(0.3), Inches(0.7), Inches(5.0), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    mgmt_rows = [
        ['文書名', 'SBDS フェーズ管理書 v1.0'],
        ['作成日', '2026-06-05'],
        ['作成者', '自律COO（Claude Code）'],
        ['承認者', '代表取締役CEO 松浦 学'],
        ['根拠文書', 'BRD v1.0 / SRS v1.0 / PLAN-20260604-001 Rev.2'],
    ]
    t = sl.shapes.add_table(len(mgmt_rows), 2, Inches(0.3), Inches(1.0), Inches(5.0), Inches(1.8)).table
    for ri, (k, v) in enumerate(mgmt_rows):
        t.cell(ri, 0).text = k
        t.cell(ri, 1).text = v
        for ci in range(2):
            t.cell(ri, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(9)

    add_text_box(sl, '開発基本方針', Inches(5.5), Inches(0.7), Inches(5.8), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    principles = [
        '① Gate制厳守: G0→G1→G2→G3→G4の順。CEOへの承認なしにGateを飛ばさない',
        '② レイアウト固定: LAYOUT_MASTER.mdを唯一の正本とする。変更はCEO承認後',
        '③ GeminiのUI禁止: レイアウト・画面実装にGeminiを関与させない',
        '④ DevSepOps: APIキーはフロントエンドに一切露出しない',
        '⑤ 自律文書化: 全変更をCHANGELOG_JA.mdに即時コミット',
        '⑥ 不明点報告: 「松浦CEO要件定義待ち」として即座に報告。でっち上げ禁止',
    ]
    for i, p in enumerate(principles):
        add_text_box(sl, p, Inches(5.5), Inches(1.0) + i * Inches(0.42), Inches(5.8), Inches(0.38),
                     font_size=Pt(8.5), color=DARK_NAVY)
    slide_number_tag(sl, 1, TOTAL, is_landscape=True)

    # Slide 2: フェーズ全体マップ
    sl = blank_slide(prs)
    title_bar(sl, 'フェーズ全体マップ（Gate 0〜7）', 'Phase 1完了目標: 2027年2月（3万世帯）', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    gates = [
        ('G0', '2026/07末', 'GCP基盤構築\nLAYOUT_MASTER登録\nCI整備', GREEN),
        ('G1', '2026/09末', 'SBDS完成\nTMS-SET-001\nTMS-DRV-001', ORANGE),
        ('G2', '2026/11末', 'SURPLUS SHIFT v14.2\ni18n 15ヶ国語\nResearch完成', DARK_NAVY),
        ('G3', '2027/01末', 'Marketing-Sys\n朝夕スケジューラー\nGOV/S10完成', DARK_NAVY),
        ('G4', '2027/02末', 'Cloud Run proxy\n本番デプロイ\nGo-Live（3万世帯）', RED),
    ]
    gate_w = Inches(2.0)
    for i, (gate, date, content, color) in enumerate(gates):
        gx = Inches(0.3) + i * Inches(2.22)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(5.5), fill_rgb=LIGHT_BG, line_rgb=color)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(0.55), fill_rgb=color, line_rgb=color)
        add_text_box(sl, gate, gx, Inches(0.72), gate_w, Inches(0.3),
                     font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, date, gx, Inches(1.05), gate_w, Inches(0.25),
                     font_size=Pt(8), color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, content, gx + Inches(0.1), Inches(1.35), gate_w - Inches(0.2), Inches(4.5),
                     font_size=Pt(9), color=DARK_NAVY)

    # Phase 1 括弧
    add_text_box(sl, '◀─────────────── Phase 1（〜2027年2月） ───────────────▶',
                 Inches(0.3), Inches(6.4), Inches(10.8), Inches(0.3),
                 font_size=Pt(9), bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 2, TOTAL, is_landscape=True)

    # Slide 3: SBDS工数内訳
    sl = blank_slide(prs)
    title_bar(sl, 'G1 SBDS工数内訳・予算管理', 'Gate 1（〜2026/09末）スプリント', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'SBDS G1 工数内訳', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    work_rows = [
        ['TMS-SET-001 基本実装（BuildingMaster/RoomRecord）', '自律COO（Claude Code）', '完了', '¥0'],
        ['IndexedDB v142移行（v140→v142）', '自律COO', '完了', '¥0'],
        ['Jaro-Winkler名寄せエンジン（D_jw≥0.85）', '自律COO', '完了', '¥0'],
        ['STATUS_LOCKED_BY_LABOR_LAW（労働法ロック）', '自律COO', '完了', '¥0'],
        ['TMS-DRV-001 UI（冷凍警告・ETA・ルーティング）', '自律COO', '完了', '¥0'],
        ['DXFインポート（dxf-parser.js統合）', '自律COO', '未着手（G1内）', '¥0'],
        ['WebXR AR計測（±50cm許容）', '自律COO', '判断③待ち', '¥0'],
        ['LIFF→PWA導線設計（QRコード・案内文）', '自律COO', '未着手（G1内）', '¥0'],
        ['Cloud Run APIサーバー（FastAPI）', '自律COO', '未着手（G0）', '¥0'],
        ['Firestore スキーマ設計・デプロイ', '自律COO', '未着手（G0）', '¥0'],
        ['LINE Webhook + Redis Consumer Group', '自律COO', '完了（src/layer3）', '¥0'],
    ]
    t = sl.shapes.add_table(len(work_rows) + 1, 4, Inches(0.3), Inches(1.0), Inches(10.8), Inches(4.5)).table
    for i, h in enumerate(['タスク', '担当', '状態', '費用']):
        t.cell(0, i).text = h
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)
    for ri, row in enumerate(work_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8)

    add_text_box(sl, 'AI自律開発（Claude Code）により開発人件費 = ¥0。インフラコストのみ発生。',
                 Inches(0.3), Inches(5.7), Inches(10.8), Inches(0.3),
                 font_size=Pt(9), bold=True, color=GREEN)
    slide_number_tag(sl, 3, TOTAL, is_landscape=True)

    # Slide 4: 予算管理表
    sl = blank_slide(prs)
    title_bar(sl, '予算管理表・FinOpsコスト（SBDS / 全社）', '総開発費: ¥2,844万（補助金活用後 ¥1,293万）', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, '全社 総開発費内訳', Inches(0.3), Inches(0.7), Inches(5.5), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    budget_rows = [
        ['開発費（AI自律開発）', '¥0', '100%削減'],
        ['GCPインフラ（8ヶ月）', '¥138,000', ''],
        ['Claude API（本番12ヶ月）', '¥540,000', ''],
        ['CEO・COO活動費（8ヶ月）', '¥27,759,000', '補助金対象'],
        ['合計', '¥28,440,000（¥2,844万）', ''],
        ['補助金控除後', '¥12,930,000（¥1,293万）', '補助金¥15,510,000想定'],
    ]
    t = sl.shapes.add_table(len(budget_rows) + 1, 3, Inches(0.3), Inches(1.0), Inches(5.5), Inches(2.5)).table
    for i, h in enumerate(['区分', '金額', '備考']):
        t.cell(0, i).text = h
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    for ri, row in enumerate(budget_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(9)

    add_text_box(sl, 'GCP月次コスト（SBDS / 3万世帯）', Inches(6.0), Inches(0.7), Inches(5.0), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    gcp_rows = [
        ['Cloud Run', '¥750〜2,250'],
        ['Memorystore Redis', '¥2,250'],
        ['Firestore', '¥0〜450'],
        ['BigQuery', '¥300'],
        ['LINE Messaging', '¥750〜1,500'],
        ['月額合計', '¥4,050〜6,750'],
        ['1個口コスト', '¥0.10〜0.14 ✅'],
    ]
    t2 = sl.shapes.add_table(len(gcp_rows) + 1, 2, Inches(6.0), Inches(1.0), Inches(5.0), Inches(2.5)).table
    for i, h in enumerate(['サービス', '月額']):
        t2.cell(0, i).text = h
        t2.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t2.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    for ri, row in enumerate(gcp_rows):
        for ci, val in enumerate(row):
            t2.cell(ri + 1, ci).text = val
            t2.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(9)

    add_text_box(sl, '補助金活用計画', Inches(0.3), Inches(3.65), Inches(10.8), Inches(0.25),
                 font_size=Pt(10), bold=True, color=DARK_NAVY)
    subs = [
        ['IT導入補助金（デジタル化基盤）', '〜¥3,500,000', '申請予定'],
        ['事業再構築補助金', '〜¥6,000,000', '申請予定'],
        ['東京都DX推進助成', '〜¥2,000,000', '検討中'],
        ['人材開発支援助成金', '〜¥4,010,000', '申請予定'],
    ]
    t3 = sl.shapes.add_table(len(subs) + 1, 3, Inches(0.3), Inches(4.0), Inches(10.8), Inches(2.0)).table
    for i, h in enumerate(['補助金種別', '想定上限', '状態']):
        t3.cell(0, i).text = h
        t3.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t3.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    for ri, row in enumerate(subs):
        for ci, val in enumerate(row):
            t3.cell(ri + 1, ci).text = val
            t3.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(9)

    slide_number_tag(sl, 4, TOTAL, is_landscape=True)

    # Slide 5: リスク管理
    sl = blank_slide(prs)
    title_bar(sl, 'リスク管理表', 'SBDS Gate 1 スプリント', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    risks = [
        ['LINE Mini App審査遅延', '高', 'SBDS-G1遅延', 'PWA先行リリース、後からLIFF追加', 'G1'],
        ['WebXR on LINE LIFF 動作しない', '高', 'AR機能無効', 'PWA別URL（/ar）+ QRコード導線（判断⑥A確定済）', 'G1'],
        ['DXF品質が低く自動抽出精度不足', '高', '手動入力工数増加', 'フォールバック手動入力UI完備済', 'G1'],
        ['iOS WebXR未対応（iOS15以下）', '中', 'iOS一部ユーザーAR利用不可', 'A案採用時: 手動入力で代替、B案採用時: 追加工数3-4日', '判断③待ち'],
        ['Jaro-Winkler偽陰性（別人を同一判定）', '低', '誤配送リスク', '閾値0.85は実証済、500件テスト実施', 'G1'],
        ['GCP初期コスト（スケール前）', '中', '1個口コスト超過', 'Firestore軽量構成で初期対応', 'G0'],
        ['セッション間レイアウト崩壊', '低', 'UI一貫性破壊', 'LAYOUT_MASTER.md + layout-guard CI（実装済）', '全Gate'],
    ]
    t = sl.shapes.add_table(len(risks) + 1, 5, Inches(0.3), Inches(0.7), Inches(10.8), Inches(5.5)).table
    for i, h in enumerate(['リスク', '確率', '影響', '対策', 'Gate']):
        t.cell(0, i).text = h
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, i).text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    for ri, row in enumerate(risks):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = Pt(8.5)

    slide_number_tag(sl, 5, TOTAL, is_landscape=True)

    prs.save(output)
    print(f"✅ {output} 生成完了")


if __name__ == '__main__':
    gen_seq()
    gen_ui()
    gen_phase()
