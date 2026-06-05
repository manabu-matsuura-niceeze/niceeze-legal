"""RESEARCH部 5文書生成スクリプト"""
import os

# ─── python-docx ────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, Mm, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# ─── python-pptx ────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt as PPt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── カラー定義 ─────────────────────────────────────────────────────
DARK_NAVY = RGBColor(0x1a, 0x3a, 0x5c)
ORANGE    = RGBColor(0xf5, 0xa6, 0x23)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)

DARK_NAVY_D = DocxRGB(0x1a, 0x3a, 0x5c)
ORANGE_D    = DocxRGB(0xf5, 0xa6, 0x23)
WHITE_D     = DocxRGB(0xFF, 0xFF, 0xFF)
LIGHT_BG_D  = DocxRGB(0xF0, 0xF4, 0xF8)

OUT_DIR = "docs/RESEARCH"

# ═══════════════════════════════════════════════════════════════════
# DOCX ユーティリティ
# ═══════════════════════════════════════════════════════════════════

def set_margins(doc, top=20, bottom=20, left=20, right=20):
    sec = doc.sections[0]
    sec.top_margin    = Mm(top)
    sec.bottom_margin = Mm(bottom)
    sec.left_margin   = Mm(left)
    sec.right_margin  = Mm(right)


def add_heading(doc, text, level=1, color=None):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after  = Pt(2)
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(12) if level == 1 else Pt(10.5)
    run.font.color.rgb = color or DARK_NAVY_D
    return p


def add_body(doc, text):
    p = doc.add_paragraph(text)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after  = Pt(2)
    for run in p.runs:
        run.font.size = Pt(10.5)
    return p


def set_cell_bg(cell, color_hex):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), color_hex)
    tcPr.append(shd)


def set_table_width(table, width_dxa=9026):
    tbl = table._tbl
    tblPr = tbl.find(qn('w:tblPr'))
    if tblPr is None:
        tblPr = OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    tblW = OxmlElement('w:tblW')
    tblW.set(qn('w:w'), str(width_dxa))
    tblW.set(qn('w:type'), 'dxa')
    tblPr.append(tblW)


def add_full_table(doc, headers, rows, header_bg='1a3a5c', header_fg=None):
    if header_fg is None:
        header_fg = WHITE_D
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    set_table_width(table, 9026)
    table.style = 'Table Grid'

    hdr_row = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr_row.cells[i]
        set_cell_bg(cell, header_bg)
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(h)
        run.bold = True
        run.font.size = Pt(9)
        run.font.color.rgb = header_fg

    for ri, row_data in enumerate(rows):
        tr = table.rows[ri + 1]
        bg = 'F0F4F8' if ri % 2 == 0 else 'FFFFFF'
        for ci, val in enumerate(row_data):
            cell = tr.cells[ci]
            set_cell_bg(cell, bg)
            p = cell.paragraphs[0]
            run = p.add_run(str(val))
            run.font.size = Pt(9)

    return table


def add_footer_docx(doc, text):
    for section in doc.sections:
        footer = section.footer
        p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
        p.clear()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.size = Pt(8)
        run.font.color.rgb = DocxRGB(0x64, 0x74, 0x8b)


# ═══════════════════════════════════════════════════════════════════
# PPTX ユーティリティ
# ═══════════════════════════════════════════════════════════════════

A4_W_LAND = Inches(11.69)
A4_H_LAND = Inches(8.27)
A4_W_PORT = Inches(8.27)
A4_H_PORT = Inches(11.69)


def new_prs_landscape():
    prs = Presentation()
    prs.slide_width  = A4_W_LAND
    prs.slide_height = A4_H_LAND
    return prs


def new_prs_portrait():
    prs = Presentation()
    prs.slide_width  = A4_W_PORT
    prs.slide_height = A4_H_PORT
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=PPt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=PPt(9), bold=False, color=None,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color or DARK_NAVY
    return txBox


def add_bullet_text(slide, items, x, y, w, h, font_size=PPt(10), color=None):
    """Add a text box with bullet-point items."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f'• {item}'
        run.font.size = font_size
        run.font.color.rgb = color or DARK_NAVY
    return txBox


def title_bar(slide, title, subtitle='', is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W_PORT
    add_rect(slide, 0, 0, sw, Inches(0.6), fill_rgb=DARK_NAVY)
    add_text_box(slide, title, Inches(0.15), Inches(0.05), sw - Inches(0.3), Inches(0.35),
                 font_size=PPt(13), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.15), Inches(0.38), sw - Inches(0.3), Inches(0.2),
                     font_size=PPt(8), color=LIGHT_BG, align=PP_ALIGN.LEFT)


def footer_bar(slide, text, is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W_PORT
    sh = A4_H_LAND if is_landscape else A4_H_PORT
    add_rect(slide, 0, sh - Inches(0.25), sw, Inches(0.25), fill_rgb=DARK_NAVY)
    add_text_box(slide, text, Inches(0.1), sh - Inches(0.23), sw - Inches(0.2), Inches(0.2),
                 font_size=PPt(7), color=LIGHT_BG, align=PP_ALIGN.CENTER)


def slide_number_tag(slide, num, total, is_landscape=True):
    sw = A4_W_LAND if is_landscape else A4_W_PORT
    sh = A4_H_LAND if is_landscape else A4_H_PORT
    add_text_box(slide, f'{num} / {total}',
                 sw - Inches(0.8), sh - Inches(0.23), Inches(0.7), Inches(0.2),
                 font_size=PPt(7), color=LIGHT_BG, align=PP_ALIGN.RIGHT)


def pptx_table(slide, headers, rows, x, y, w, h):
    """Add a simple pptx table."""
    t = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    for ci, h in enumerate(headers):
        t.cell(0, ci).text = h
        run = t.cell(0, ci).text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = PPt(9)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = str(val)
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    return t


# ═══════════════════════════════════════════════════════════════════
# 1. RESEARCH_BRD_v1.0.docx
# ═══════════════════════════════════════════════════════════════════

def gen_brd():
    output = f"{OUT_DIR}/RESEARCH_BRD_v1.0.docx"
    doc = Document()
    set_margins(doc)
    add_footer_docx(doc, 'RESEARCH部 リサーチ情報システム  |  ビジネス要件定義書 v1.0  |  © 2026 株式会社NiceEze  Confidential')

    # 表紙
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('RESEARCH部 — ビジネス要件定義書 (BRD) v1.0')
    run.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_NAVY_D

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run('EC仕入れリサーチ自動化システム（RESEARCH）　v1.0').font.size = Pt(12)
    doc.add_paragraph()

    # 1. 文書管理表
    add_heading(doc, '1. 文書管理表')
    add_full_table(doc,
        ['項目', '内容'],
        [
            ['文書名', 'RESEARCH部 ビジネス要件定義書'],
            ['バージョン', 'v1.0'],
            ['作成日', '2026-06-05'],
            ['最終更新日', '2026-06-05'],
            ['作成者', 'NiceEze 自律COO'],
            ['承認者', '代表取締役CEO 松浦 学'],
            ['ステータス', '承認済（Gate 1スプリント適用中）'],
            ['関連文書', 'RESEARCH_SRS_v1.0.docx / RESEARCH_SEQ_v1.0.pptx'],
        ]
    )

    # 2. ビジネス背景
    add_heading(doc, '2. ビジネス背景（EC仕入れ情報格差問題）')
    add_body(doc,
        'EC仕入れ業務において、仕入れ担当者は複数プラットフォーム（Amazon・楽天・Yahoo等）の価格を手動で調査しており、'
        '1商品あたり30〜60分の調査工数が発生している。情報非対称性により最安仕入れ先の見逃しが頻発し、'
        '利益率の低下と意思決定の遅延が深刻な課題となっている。'
        'また、定番商品の選定はバイヤーの経験則に依存しており、データドリブンな判断基準が欠如している。'
        'RESEARCH部は8社固定プラットフォームの価格マトリクス自動取得・トレンドスコア算出・統合ダッシュボード提供により、'
        '価格調査時間80%削減・定番商品判定精度85%以上を目標とする。'
    )
    add_full_table(doc,
        ['課題', '現状値', 'RESEARCH目標値'],
        [
            ['価格調査工数（1商品）', '30〜60分/件', '6分以内（80%削減）'],
            ['調査対象プラットフォーム数', '担当者依存（2〜4社）', '8社固定（均一カバレッジ）'],
            ['定番商品判定精度', '経験則（定量基準なし）', '≥85%（retention_score）'],
            ['情報非対称性リスク', '高（最安先見逃し多数）', '最安amber強調で即時把握'],
        ]
    )

    # 3. 主要機能一覧
    add_heading(doc, '3. 主要機能一覧')
    add_full_table(doc,
        ['ID', '機能名', '概要', '優先度'],
        [
            ['RES-A01', '8社価格マトリクス',
             'Amazon/楽天/Yahoo/au PAY/Qoo10/ヨドバシ/ビック/ヤマダの8社固定。unit_price/case_price取得。SHA-256 cache_key管理',
             'Must'],
            ['RES-A02', 'トレンドスコア',
             'growth_score/bestseller_score/retention_score算出。RETENTION_THRESHOLD=0.6でis_staple判定。≥0.8でTODO起票',
             'Must'],
            ['RES-A03', '統合ダッシュボード',
             'RES-A01/A02タブ切替。モバイル対応。IndexedDB niceeze_cache_v142キャッシュ。LAYOUT_MASTER準拠',
             'Must'],
        ]
    )

    # 4. ステークホルダー
    add_heading(doc, '4. ステークホルダー定義')
    add_full_table(doc,
        ['ステークホルダー', '役割', '主要タッチポイント', '優先度'],
        [
            ['仕入れ担当者', '価格マトリクス参照・発注判断', 'RES-A01 価格マトリクス画面', 'Must'],
            ['EC運営者', 'トレンド分析・商品選定', 'RES-A02 トレンド分析画面', 'Must'],
            ['CEO（松浦 学）', 'KPI確認・Gate承認', '統合ダッシュボード / COO報告', 'Must'],
        ]
    )

    # 5. KPI
    add_heading(doc, '5. 成功指標（KPI）')
    add_full_table(doc,
        ['KPI', '目標値', '計測方法', 'Gate'],
        [
            ['価格調査時間削減率', '≥80%削減', '調査工数ビフォーアフター計測', 'G1'],
            ['定番商品判定精度', '≥85%', 'retention_score実績 vs 実売', 'G2'],
            ['APIレスポンス時間', '≤0.7秒', 'Cloud Run p95レイテンシ', 'G1'],
            ['APIキーフロント露出件数', '0件', 'bandit自動スキャン', '全Gate'],
            ['テスト全Pass数', '38テスト全Pass', 'unittest CI結果', 'G1'],
            ['月額インフラコスト（MVP）', '¥0（無料枠）', 'GCP請求 / 月次', 'G1'],
        ]
    )

    # 6. Gate制
    add_heading(doc, '6. Gate制（G0〜G4）/ FinOps 月額上限管理')
    add_full_table(doc,
        ['Gate', '完了条件', 'FinOps上限', '備考'],
        [
            ['G0', 'GCP環境構築・Secret Manager設定完了', '¥0（無料枠）', 'インフラ基盤のみ'],
            ['G1', 'RES-A01/A02基本動作・bandit 0件・38テスト全Pass', '¥0（無料枠内）', 'IndexedDB v142 / Cloud Run無料枠'],
            ['G2', 'RES-A03統合ダッシュボード完成・KPI初期計測', '¥2,250〜¥4,500/月', 'Keepa API / Google Trends実連携'],
            ['G3', 'AI強化（Claude API統合）・定番判定精度≥85%', '¥4,500/月', 'Claude API月次呼出'],
            ['G4', '全機能Go-Live・全KPI達成', '¥5,000/月上限', 'スケールアップ準備'],
        ]
    )
    add_body(doc,
        'FinOps原則: MVP（G1完了）時点での月額インフラコスト¥0（Cloud Run・Firestore無料枠活用）。'
        'G3以降のClaude API費用は月¥2,250〜¥4,500。月額上限¥5,000でGCP予算アラート設定。'
    )

    # 7. 制約条件
    add_heading(doc, '7. 制約条件・前提条件')
    add_full_table(doc,
        ['区分', '内容'],
        [
            ['セキュリティ', 'bandit 0件必須（CI自動チェック）'],
            ['プライバシー', 'PII不使用（個人情報を収集・処理しない）'],
            ['暗号化', 'AES-256暗号化 + RLS（Firestore）'],
            ['パフォーマンス', 'APIレスポンス0.7秒以下（Cloud Run p95）'],
            ['UI準拠', 'LAYOUT_MASTER準拠（font-mono tabular-nums tracking-tight）'],
            ['キャッシュ', 'IndexedDB niceeze_cache_v142 使用'],
            ['前提条件', 'GCP Secret ManagerにAPIキー設定済み（Gate 0）'],
            ['前提条件', 'research_dashboard.html がG1完了時点で実装済み'],
        ]
    )

    doc.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 2. RESEARCH_SRS_v1.0.docx
# ═══════════════════════════════════════════════════════════════════

def gen_srs():
    output = f"{OUT_DIR}/RESEARCH_SRS_v1.0.docx"
    doc = Document()
    set_margins(doc)
    add_footer_docx(doc, 'RESEARCH部 リサーチ情報システム  |  ソフトウェア要件仕様書 v1.0  |  © 2026 株式会社NiceEze  Confidential')

    # 表紙
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('RESEARCH部 — ソフトウェア要件仕様書 (SRS) v1.0')
    run.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_NAVY_D

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run('EC仕入れリサーチ自動化システム（RESEARCH）　v1.0').font.size = Pt(12)
    doc.add_paragraph()

    # 1. システム概要
    add_heading(doc, '1. システム概要')
    add_body(doc,
        'RESEARCH部システムはCloud Run（APIサーバー）+ Firestore（データストア）+ IndexedDB niceeze_cache_v142（クライアントキャッシュ）の'
        '3層構成で実装する。フロントエンドはresearch_dashboard.htmlとして提供し、LAYOUT_MASTER準拠のUIを実現する。'
        'APIキーはフロントエンドに一切露出せず、Cloud Run経由でSecret Managerから取得する。'
    )
    add_full_table(doc,
        ['コンポーネント', '技術スタック', '役割'],
        [
            ['APIサーバー', 'Cloud Run (Python / FastAPI)', '価格取得・トレンド計算・キャッシュ管理'],
            ['データストア', 'Firestore', 'PriceRecord / PriceMatrix / ProductTrend永続化'],
            ['クライアントキャッシュ', 'IndexedDB niceeze_cache_v142', 'オフライン対応・TTLキャッシュ'],
            ['フロントエンド', 'research_dashboard.html', 'LAYOUT_MASTER準拠ダッシュボード'],
            ['シークレット管理', 'GCP Secret Manager', 'APIキー安全管理（bandit 0件）'],
        ]
    )

    # 2. 機能要件
    add_heading(doc, '2. 機能要件')

    add_heading(doc, '2.1 RES-A01: 8社価格マトリクス', level=2)
    add_body(doc, '8社固定プラットフォームから unit_price（1個単価）および case_price（1ケース価格）を取得し、'
             'SHA-256ハッシュによるcache_keyで重複排除・キャッシュ管理を行う。')
    add_full_table(doc,
        ['プラットフォーム', '取得項目', '備考'],
        [
            ['Amazon', 'unit_price / case_price', 'SHA-256 cache_key管理'],
            ['楽天市場', 'unit_price / case_price', 'SHA-256 cache_key管理'],
            ['Yahoo!ショッピング', 'unit_price / case_price', 'SHA-256 cache_key管理'],
            ['au PAY マーケット', 'unit_price / case_price', 'SHA-256 cache_key管理'],
            ['Qoo10', 'unit_price / case_price', 'SHA-256 cache_key管理'],
            ['ヨドバシ.com', 'unit_price / case_price', 'SHA-256 cache_key管理'],
            ['ビックカメラ', 'unit_price / case_price', 'SHA-256 cache_key管理'],
            ['ヤマダ電機', 'unit_price / case_price', 'SHA-256 cache_key管理'],
        ]
    )

    add_heading(doc, '2.2 RES-A02: トレンドスコア', level=2)
    add_body(doc, 'growth_score（急成長度）/ bestseller_score（売れ筋度）/ retention_score（定番残存率）の3指標を算出する。'
             'RETENTION_THRESHOLD=0.6でis_staple（定番商品）フラグを付与。retention_score≥0.8の場合はTODO起票。')
    add_full_table(doc,
        ['スコア', '算出方法', 'しきい値', '出力'],
        [
            ['growth_score', '直近7日 vs 前週比成長率', '—', '0.0〜1.0'],
            ['bestseller_score', 'カテゴリ内売上ランク正規化', '—', '0.0〜1.0'],
            ['retention_score', '定番残存率（継続購買率）', 'RETENTION_THRESHOLD=0.6', 'is_staple: bool'],
            ['TODO起票', 'retention_score ≥ 0.8', '0.8', 'TODO自動生成'],
        ]
    )

    # 3. 非機能要件
    add_heading(doc, '3. 非機能要件')
    add_full_table(doc,
        ['区分', '要件', '計測基準'],
        [
            ['セキュリティ', 'bandit 0件（CI必須）', 'bandit自動スキャン'],
            ['プライバシー', 'PII不使用', 'コードレビュー'],
            ['暗号化', 'AES-256 + RLS（Firestore）', 'セキュリティ監査'],
            ['パフォーマンス', 'APIレスポンス≤0.7秒', 'Cloud Run p95レイテンシ'],
            ['可用性', 'Cloud Run自動スケール', 'GCP SLA準拠'],
            ['オフライン', 'IndexedDB niceeze_cache_v142', 'TTLキャッシュ管理'],
        ]
    )

    # 4. データモデル
    add_heading(doc, '4. データモデル')
    add_full_table(doc,
        ['モデル名', 'フィールド', '説明'],
        [
            ['PriceRecord', 'platform / unit_price / case_price / cache_key / fetched_at',
             '各プラットフォームの価格レコード（SHA-256 cache_key）'],
            ['PriceMatrix', 'keyword / category / records[PriceRecord] / created_at',
             '8社分PriceRecordを束ねたマトリクス'],
            ['ProductTrend', 'keyword / category / datapoints[TrendDataPoint] / is_staple / growth_score / bestseller_score / retention_score',
             'トレンド集計結果'],
            ['TrendDataPoint', 'date / value / source',
             '日次トレンドデータポイント'],
        ]
    )

    # 5. APIエンドポイント
    add_heading(doc, '5. APIエンドポイント')
    add_full_table(doc,
        ['メソッド', 'パス', 'パラメータ', '説明'],
        [
            ['GET', '/health', '—', 'ヘルスチェック（bandit 0件確認用）'],
            ['GET', '/price', 'keyword: str, category: str', '8社価格マトリクス取得（IndexedDB v142キャッシュ優先）'],
            ['GET', '/trend', 'keyword: str, category: str, days: int', 'トレンドスコア取得（growth/bestseller/retention）'],
        ]
    )

    # 6. テスト要件
    add_heading(doc, '6. テスト要件')
    add_body(doc, 'unittest による38テスト全Pass必須。CI/CDパイプラインでbandit 0件・テスト全Pass確認後にデプロイ許可。')
    add_full_table(doc,
        ['テストカテゴリ', 'テスト数', '内容'],
        [
            ['PriceFetcher単体テスト', '16件', '8社各2テスト（正常系・異常系）'],
            ['TrendFetcher単体テスト', '9件', 'growth/bestseller/retention各3テスト'],
            ['APIエンドポイントテスト', '6件', '/health・/price・/trend 各2テスト'],
            ['キャッシュ機能テスト', '4件', 'IndexedDB v142 TTL・HIT/MISS'],
            ['セキュリティテスト', '3件', 'bandit・PII不使用・AES-256'],
            ['合計', '38件', '全Pass必須（CI/CD Gate）'],
        ]
    )

    doc.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 3. RESEARCH_SEQ_v1.0.pptx
# ═══════════════════════════════════════════════════════════════════

def gen_seq():
    output = f"{OUT_DIR}/RESEARCH_SEQ_v1.0.pptx"
    prs = new_prs_landscape()
    FOOTER = 'RESEARCH部  |  シーケンス図 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # S1: タイトル
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, A4_W_LAND, A4_H_LAND, fill_rgb=DARK_NAVY)
    add_text_box(sl, 'RESEARCH部 シーケンス図 v1.0',
                 Inches(1.0), Inches(2.5), Inches(9.5), Inches(1.2),
                 font_size=PPt(36), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'EC仕入れリサーチ自動化システム — シーケンス図',
                 Inches(1.0), Inches(3.8), Inches(9.5), Inches(0.5),
                 font_size=PPt(16), color=ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, '2026-06-05  |  © 2026 株式会社NiceEze  Confidential',
                 Inches(1.0), Inches(4.5), Inches(9.5), Inches(0.4),
                 font_size=PPt(11), color=LIGHT_BG, align=PP_ALIGN.CENTER)
    footer_bar(sl, FOOTER, is_landscape=True)
    slide_number_tag(sl, 1, TOTAL, is_landscape=True)

    # S2: RES-A01 価格検索フロー
    sl = blank_slide(prs)
    title_bar(sl, 'RES-A01 価格検索フロー', '8社固定プラットフォーム価格マトリクス取得シーケンス', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    add_text_box(sl, 'RES-A01 価格検索フロー', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(12), bold=True, color=DARK_NAVY)
    add_bullet_text(sl,
        [
            'ユーザー入力（キーワード・カテゴリ）',
            'IndexedDB niceeze_cache_v142 確認（SHA-256 cache_key照合）',
            'MISS → Cloud Run /price エンドポイント呼出',
            'PriceFetcher → 8社データ並列取得（Amazon/楽天/Yahoo/au PAY/Qoo10/ヨドバシ/ビック/ヤマダ）',
            'unit_price / case_price 正規化・PriceMatrix構築',
            '最安プラットフォーム amber強調ハイライト',
            'IndexedDB v142 に結果キャッシュ（TTL管理）',
            '画面表示 — LAYOUT_MASTER準拠（font-mono tabular-nums tracking-tight）',
        ],
        Inches(0.5), Inches(1.1), Inches(10.5), Inches(5.8),
        font_size=PPt(11), color=DARK_NAVY
    )
    slide_number_tag(sl, 2, TOTAL, is_landscape=True)

    # S3: RES-A02 トレンド分析フロー
    sl = blank_slide(prs)
    title_bar(sl, 'RES-A02 トレンド分析フロー', 'growth/bestseller/retention スコア算出シーケンス', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    add_text_box(sl, 'RES-A02 トレンド分析フロー', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(12), bold=True, color=DARK_NAVY)
    add_bullet_text(sl,
        [
            'カテゴリ選択（8カテゴリチップから選択）',
            'モード切替（売れ筋 / 急成長 / 定番残存 3モードボタン）',
            'TrendFetcher → Cloud Run /trend エンドポイント呼出',
            'growth_score 算出（直近7日 vs 前週比成長率）',
            'bestseller_score 算出（カテゴリ内売上ランク正規化）',
            'retention_score 算出 → 定番判定（RETENTION_THRESHOLD=0.6 → is_staple）',
            'retention_score ≥ 0.8 → TODO自動起票',
            'S_retention font-mono tabular-nums表示',
        ],
        Inches(0.5), Inches(1.1), Inches(10.5), Inches(5.8),
        font_size=PPt(11), color=DARK_NAVY
    )
    slide_number_tag(sl, 3, TOTAL, is_landscape=True)

    # S4: G2以降実API連携
    sl = blank_slide(prs)
    title_bar(sl, 'G2以降実API連携', 'Keepa API / Google Trends / Cloud Run proxy / IndexedDB v142 TTL', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    add_text_box(sl, 'G2以降 実API連携アーキテクチャ', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(12), bold=True, color=DARK_NAVY)

    # 4ボックス配置
    boxes = [
        ('Keepa API', 'Amazon価格履歴取得\n商品ランク推移\nSHA-256 cache_key', Inches(0.3), Inches(1.2)),
        ('Google Trends', 'キーワード検索トレンド\nカテゴリ別急成長検知\ngrowth_score算出', Inches(3.2), Inches(1.2)),
        ('Cloud Run proxy', 'APIキー完全秘匿\nSecret Manager連携\nbandit 0件必須', Inches(6.1), Inches(1.2)),
        ('IndexedDB v142 TTL', 'niceeze_cache_v142\nTTL管理（価格:1h/トレンド:6h）\nオフライン対応', Inches(9.0), Inches(1.2)),
    ]
    for label, desc, bx, by in boxes:
        add_rect(sl, bx, by, Inches(2.7), Inches(2.0), fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
        add_rect(sl, bx, by, Inches(2.7), Inches(0.4), fill_rgb=DARK_NAVY, line_rgb=DARK_NAVY)
        add_text_box(sl, label, bx + Inches(0.05), by + Inches(0.05), Inches(2.6), Inches(0.3),
                     font_size=PPt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, desc, bx + Inches(0.1), by + Inches(0.45), Inches(2.5), Inches(1.4),
                     font_size=PPt(9), color=DARK_NAVY)

    add_text_box(sl, 'G2完了条件: 実APIデータによるPriceMatrix更新 / TrendDataPoint日次蓄積開始 / IndexedDB v142 TTL最適化完了',
                 Inches(0.3), Inches(3.4), Inches(11.0), Inches(0.35),
                 font_size=PPt(9), bold=True, color=ORANGE)

    # フロー図（テキスト）
    add_text_box(sl, 'データフロー: ユーザー → research_dashboard.html → Cloud Run proxy → [Keepa/Google Trends] → Firestore → IndexedDB v142',
                 Inches(0.3), Inches(3.85), Inches(11.0), Inches(0.35),
                 font_size=PPt(9), color=DARK_NAVY)

    slide_number_tag(sl, 4, TOTAL, is_landscape=True)

    # S5: エラーハンドリング
    sl = blank_slide(prs)
    title_bar(sl, 'エラーハンドリング', 'API障害 / IndexedDB失敗 / タイムアウト 各フォールバック', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)
    add_text_box(sl, 'エラーパターン別フォールバック一覧', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(12), bold=True, color=DARK_NAVY)

    errors = [
        ('API障害（価格取得失敗）',
         'モックデータでフォールバック表示 / エラーバナー表示 / 再試行ボタン提供'),
        ('IndexedDB失敗（v142 read error）',
         'IndexedDB スキップ → ネット直接取得 / audit_logにエラー記録 / TTL無効化'),
        ('タイムアウト（0.7秒超過）',
         'Cloud Run → 503レスポンス / フロント側でローディング解除 / エラーメッセージ表示'),
    ]
    RED_COLOR = RGBColor(0xEF, 0x44, 0x44)
    for i, (title, fallback) in enumerate(errors):
        ey = Inches(1.1) + i * Inches(2.0)
        add_rect(sl, Inches(0.3), ey, Inches(11.0), Inches(1.8), fill_rgb=LIGHT_BG, line_rgb=RED_COLOR)
        add_rect(sl, Inches(0.3), ey, Inches(11.0), Inches(0.4), fill_rgb=RGBColor(0x7f, 0x1d, 0x1d), line_rgb=RED_COLOR)
        add_text_box(sl, f'ERROR: {title}',
                     Inches(0.4), ey + Inches(0.06), Inches(10.5), Inches(0.28),
                     font_size=PPt(10), bold=True, color=RGBColor(0xfc, 0xa5, 0xa5))
        add_text_box(sl, f'フォールバック: {fallback}',
                     Inches(0.4), ey + Inches(0.5), Inches(10.5), Inches(1.0),
                     font_size=PPt(9.5), color=DARK_NAVY, word_wrap=True)

    slide_number_tag(sl, 5, TOTAL, is_landscape=True)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 4. RESEARCH_UI_v1.0.pptx
# ═══════════════════════════════════════════════════════════════════

def gen_ui():
    output = f"{OUT_DIR}/RESEARCH_UI_v1.0.pptx"
    prs = new_prs_portrait()
    FOOTER = 'RESEARCH部  |  UI設計 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # S1: タイトル
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, A4_W_PORT, A4_H_PORT, fill_rgb=DARK_NAVY)
    add_text_box(sl, 'RESEARCH部 UI設計 v1.0',
                 Inches(0.4), Inches(3.5), Inches(7.3), Inches(1.0),
                 font_size=PPt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'EC仕入れリサーチ自動化システム — UI設計書',
                 Inches(0.4), Inches(4.6), Inches(7.3), Inches(0.5),
                 font_size=PPt(14), color=ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'v1.0  |  2026-06-05  |  © 2026 株式会社NiceEze',
                 Inches(0.4), Inches(5.2), Inches(7.3), Inches(0.35),
                 font_size=PPt(10), color=LIGHT_BG, align=PP_ALIGN.CENTER)
    footer_bar(sl, FOOTER, is_landscape=False)
    slide_number_tag(sl, 1, TOTAL, is_landscape=False)

    # S2: RES-A01 価格マトリクス
    sl = blank_slide(prs)
    title_bar(sl, 'RES-A01 価格マトリクス', '8社固定テーブル / 最安amber強調 / 1ケース価格・1個単価 / LAYOUT_MASTER準拠', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'LAYOUT_MASTER.md 準拠 — 8社固定価格マトリクス', Inches(0.2), Inches(0.65), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=ORANGE)

    # 価格マトリクステーブル
    headers = ['プラットフォーム', '1個単価 (unit_price)', '1ケース価格 (case_price)', '最安']
    rows = [
        ['Amazon', '¥1,280', '¥12,800 (×10)', '★'],
        ['楽天市場', '¥1,350', '¥13,500 (×10)', ''],
        ['Yahoo!ショッピング', '¥1,290', '¥12,900 (×10)', ''],
        ['au PAY マーケット', '¥1,310', '¥13,100 (×10)', ''],
        ['Qoo10', '¥1,380', '¥13,800 (×10)', ''],
        ['ヨドバシ.com', '¥1,320', '¥13,200 (×10)', ''],
        ['ビックカメラ', '¥1,340', '¥13,400 (×10)', ''],
        ['ヤマダ電機', '¥1,360', '¥13,600 (×10)', ''],
    ]
    t = sl.shapes.add_table(len(rows) + 1, 4, Inches(0.2), Inches(0.92), Inches(7.6), Inches(4.0)).table
    for ci, h in enumerate(headers):
        t.cell(0, ci).text = h
        run = t.cell(0, ci).text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = PPt(8.5)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri + 1, ci)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
            else:
                run = para.add_run()
                run.text = val
            run.font.size = PPt(8)
            if ri == 0 and ci in (1, 2, 3):  # 最安行はamber色注記
                run.font.bold = True

    add_text_box(sl, '★ 最安プラットフォームは amber (#E8A020) 強調表示 / font-mono tabular-nums tracking-tight',
                 Inches(0.2), Inches(5.05), Inches(7.6), Inches(0.25),
                 font_size=PPt(8), color=ORANGE)
    add_text_box(sl, 'SHA-256 cache_key により重複リクエスト排除 / IndexedDB niceeze_cache_v142 TTLキャッシュ',
                 Inches(0.2), Inches(5.32), Inches(7.6), Inches(0.25),
                 font_size=PPt(8), color=DARK_NAVY)
    slide_number_tag(sl, 2, TOTAL, is_landscape=False)

    # S3: RES-A02 トレンド分析
    sl = blank_slide(prs)
    title_bar(sl, 'RES-A02 トレンド分析', '8カテゴリチップ / 3モードボタン / S_retention font-mono', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'LAYOUT_MASTER.md 準拠 — トレンド分析UI', Inches(0.2), Inches(0.65), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=ORANGE)

    # 8カテゴリチップ
    add_text_box(sl, '8カテゴリチップ', Inches(0.2), Inches(0.92), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=DARK_NAVY)
    categories = ['家電', 'PC周辺機器', '生活用品', '食品', 'ファッション', 'スポーツ', 'ゲーム', 'その他']
    for i, cat in enumerate(categories):
        cx = Inches(0.2) + (i % 4) * Inches(1.9)
        cy = Inches(1.18) + (i // 4) * Inches(0.4)
        add_rect(sl, cx, cy, Inches(1.75), Inches(0.32), fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
        add_text_box(sl, cat, cx, cy + Inches(0.04), Inches(1.75), Inches(0.24),
                     font_size=PPt(8.5), color=DARK_NAVY, align=PP_ALIGN.CENTER)

    # 3モードボタン
    add_text_box(sl, '3モードボタン', Inches(0.2), Inches(2.08), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=DARK_NAVY)
    modes = [('売れ筋', ORANGE), ('急成長', RGBColor(0x05, 0x96, 0x69)), ('定番残存', DARK_NAVY)]
    for i, (mode, color) in enumerate(modes):
        mx = Inches(0.2) + i * Inches(2.55)
        add_rect(sl, mx, Inches(2.34), Inches(2.3), Inches(0.45), fill_rgb=color, line_rgb=color)
        add_text_box(sl, mode, mx, Inches(2.37), Inches(2.3), Inches(0.38),
                     font_size=PPt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # スコア表示
    add_text_box(sl, 'スコア表示 (S_retention — font-mono tabular-nums)', Inches(0.2), Inches(2.9), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=DARK_NAVY)
    score_rows = [
        ['growth_score', '0.82', '急成長フラグ'],
        ['bestseller_score', '0.71', '売れ筋フラグ'],
        ['retention_score', '0.65', 'is_staple: True (≥0.6)'],
        ['retention_score', '0.83', 'is_staple: True + TODO起票 (≥0.8)'],
    ]
    t = sl.shapes.add_table(len(score_rows) + 1, 3, Inches(0.2), Inches(3.15), Inches(7.6), Inches(2.2)).table
    for ci, h in enumerate(['スコア種別', '値 (font-mono)', '判定結果']):
        t.cell(0, ci).text = h
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(score_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    add_text_box(sl, 'RETENTION_THRESHOLD=0.6 → is_staple判定 / ≥0.8 → TODO自動起票',
                 Inches(0.2), Inches(5.45), Inches(7.6), Inches(0.25),
                 font_size=PPt(8.5), color=ORANGE)
    slide_number_tag(sl, 3, TOTAL, is_landscape=False)

    # S4: 統合ダッシュボード
    sl = blank_slide(prs)
    title_bar(sl, '統合ダッシュボード', 'タブ切替(RES-A01/RES-A02) / モバイル対応 / IndexedDB v142', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'LAYOUT_MASTER.md 準拠 — 統合ダッシュボード (research_dashboard.html)', Inches(0.2), Inches(0.65), Inches(7.6), Inches(0.22),
                 font_size=PPt(9), bold=True, color=ORANGE)

    # タブUI
    add_rect(sl, Inches(0.2), Inches(0.9), Inches(7.6), Inches(0.5), fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
    tabs = [('RES-A01 価格マトリクス', True), ('RES-A02 トレンド分析', False)]
    for i, (tab, active) in enumerate(tabs):
        tx = Inches(0.2) + i * Inches(3.8)
        fill = DARK_NAVY if active else LIGHT_BG
        add_rect(sl, tx, Inches(0.9), Inches(3.75), Inches(0.5), fill_rgb=fill, line_rgb=DARK_NAVY)
        add_text_box(sl, tab, tx, Inches(0.95), Inches(3.75), Inches(0.4),
                     font_size=PPt(9), bold=active, color=WHITE if active else DARK_NAVY,
                     align=PP_ALIGN.CENTER)

    # コンテンツエリア
    add_rect(sl, Inches(0.2), Inches(1.42), Inches(7.6), Inches(4.5), fill_rgb=LIGHT_BG, line_rgb=DARK_NAVY)
    add_text_box(sl, '【RES-A01 アクティブ時】\n8社価格マトリクステーブル表示\n最安 amber 強調 / font-mono tabular-nums / unit_price・case_price\n\n【RES-A02 アクティブ時】\n8カテゴリチップ + 3モードボタン\ngrowth_score / bestseller_score / retention_score表示\nis_staple フラグ + TODO起票バッジ',
                 Inches(0.4), Inches(1.55), Inches(7.2), Inches(4.2),
                 font_size=PPt(10), color=DARK_NAVY)

    # モバイル対応
    add_text_box(sl, 'モバイル対応仕様', Inches(0.2), Inches(6.05), Inches(7.6), Inches(0.25),
                 font_size=PPt(9), bold=True, color=DARK_NAVY)
    add_full_table_pptx_note = [
        'ブレークポイント: sm(640px) / md(768px) / lg(1024px)',
        'モバイル: タブ縦スクロール / テーブル横スクロール / チップ折り返し',
        'IndexedDB niceeze_cache_v142: オフライン時はキャッシュデータ表示',
    ]
    for j, note in enumerate(add_full_table_pptx_note):
        add_text_box(sl, f'• {note}', Inches(0.3), Inches(6.35) + j * Inches(0.3), Inches(7.4), Inches(0.28),
                     font_size=PPt(8.5), color=DARK_NAVY)
    slide_number_tag(sl, 4, TOTAL, is_landscape=False)

    # S5: LAYOUT_MASTER準拠
    sl = blank_slide(prs)
    title_bar(sl, 'LAYOUT_MASTER準拠', 'font-mono tabular-nums tracking-tight / #1A2B4C/#E8A020 / 全数値DOM必須', is_landscape=False)
    footer_bar(sl, FOOTER, is_landscape=False)

    add_text_box(sl, 'フォント仕様', Inches(0.2), Inches(0.68), Inches(7.6), Inches(0.25),
                 font_size=PPt(10), bold=True, color=DARK_NAVY)
    font_rows = [
        ['数値表示全般', 'font-mono tabular-nums tracking-tight', '全価格・スコア・数量'],
        ['見出し', 'font-sans font-bold', 'タイトル・セクション見出し'],
        ['本文', 'font-sans', '説明文・ラベル'],
        ['S_retention', 'font-mono tabular-nums', 'retention_score専用クラス'],
    ]
    t = sl.shapes.add_table(len(font_rows) + 1, 3, Inches(0.2), Inches(0.95), Inches(7.6), Inches(1.8)).table
    for ci, h in enumerate(['適用箇所', 'CSSクラス', '備考']):
        t.cell(0, ci).text = h
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(font_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    add_text_box(sl, 'カラーパレット', Inches(0.2), Inches(2.9), Inches(7.6), Inches(0.25),
                 font_size=PPt(10), bold=True, color=DARK_NAVY)
    colors_def = [
        ('Navy #1A3A5C\n(= #1A2B4C系)', DARK_NAVY, 'ヘッダー・プライマリ'),
        ('Amber #E8A020\n(= #F5A623系)', ORANGE, '最安強調・アクション'),
        ('White #FFFFFF', WHITE, 'テキスト（ダーク背景）'),
        ('LightBG #F0F4F8', LIGHT_BG, 'カード背景・セカンダリ'),
    ]
    for i, (name, color, usage) in enumerate(colors_def):
        cx = Inches(0.2) + (i % 2) * Inches(3.85)
        cy = Inches(3.2) + (i // 2) * Inches(1.0)
        add_rect(sl, cx, cy, Inches(0.65), Inches(0.65), fill_rgb=color, line_rgb=DARK_NAVY)
        add_text_box(sl, name, cx + Inches(0.72), cy, Inches(2.8), Inches(0.35),
                     font_size=PPt(8), bold=True, color=DARK_NAVY)
        add_text_box(sl, usage, cx + Inches(0.72), cy + Inches(0.36), Inches(2.8), Inches(0.28),
                     font_size=PPt(7.5), color=RGBColor(0x47, 0x55, 0x69))

    add_text_box(sl, '全数値DOM必須仕様', Inches(0.2), Inches(5.35), Inches(7.6), Inches(0.25),
                 font_size=PPt(10), bold=True, color=DARK_NAVY)
    dom_rules = [
        '全数値（価格・スコア・数量）は必ずDOMテキストノードとして出力（スクリーンリーダー対応）',
        'aria-label属性で単位を明示（例: aria-label="unit_price: 1,280円"）',
        'font-mono tabular-nums により桁区切りの視覚的整列を保証',
    ]
    for j, rule in enumerate(dom_rules):
        add_text_box(sl, f'• {rule}', Inches(0.3), Inches(5.65) + j * Inches(0.33), Inches(7.4), Inches(0.3),
                     font_size=PPt(8.5), color=DARK_NAVY)
    slide_number_tag(sl, 5, TOTAL, is_landscape=False)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# 5. RESEARCH_PHASE_v1.0.pptx
# ═══════════════════════════════════════════════════════════════════

def gen_phase():
    output = f"{OUT_DIR}/RESEARCH_PHASE_v1.0.pptx"
    prs = new_prs_landscape()
    FOOTER = 'RESEARCH部  |  フェーズ計画 v1.0  |  © 2026 株式会社NiceEze  Confidential'
    TOTAL = 5

    # S1: タイトル
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, A4_W_LAND, A4_H_LAND, fill_rgb=DARK_NAVY)
    add_text_box(sl, 'RESEARCH部 フェーズ計画 v1.0',
                 Inches(1.0), Inches(2.5), Inches(9.5), Inches(1.2),
                 font_size=PPt(34), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, 'EC仕入れリサーチ自動化システム — フェーズ計画',
                 Inches(1.0), Inches(3.8), Inches(9.5), Inches(0.5),
                 font_size=PPt(16), color=ORANGE, align=PP_ALIGN.CENTER)
    add_text_box(sl, '2026-06-05  |  © 2026 株式会社NiceEze  Confidential',
                 Inches(1.0), Inches(4.5), Inches(9.5), Inches(0.4),
                 font_size=PPt(11), color=LIGHT_BG, align=PP_ALIGN.CENTER)
    footer_bar(sl, FOOTER, is_landscape=True)
    slide_number_tag(sl, 1, TOTAL, is_landscape=True)

    # S2: Gate制概要
    sl = blank_slide(prs)
    title_bar(sl, 'Gate制概要 G0〜G4', 'RESEARCH部 開発フェーズゲート', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    GREEN_C = RGBColor(0x05, 0x96, 0x69)
    RED_C   = RGBColor(0xEF, 0x44, 0x44)
    gates = [
        ('G0', '設計完了', '2026/07末', '詳細設計書完成\nDB設計確定\nGCP基盤構築', GREEN_C),
        ('G1', 'MVP完成', '2026/09末', 'RES-A01/A02実装\n38テスト全Pass\nbandit 0件\nresearch_dashboard.html', ORANGE),
        ('G2', '実API連携', '2026/11末', 'Keepa API連携\nGoogle Trends連携\nキャッシュTTL最適化', DARK_NAVY),
        ('G3', 'AI強化', '2027/01末', 'Claude API統合\n定番判定精度≥85%\nTODO自動起票高度化', RGBColor(0x6b, 0x21, 0xa8)),
        ('G4', 'Go-Live', '2027/02末', '全機能本番稼働\n全KPI達成\nSLA設定', RED_C),
    ]

    gate_w = Inches(2.0)
    for i, (gate, milestone, date, content, color) in enumerate(gates):
        gx = Inches(0.3) + i * Inches(2.22)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(5.8), fill_rgb=LIGHT_BG, line_rgb=color)
        add_rect(sl, gx, Inches(0.7), gate_w, Inches(0.8), fill_rgb=color, line_rgb=color)
        add_text_box(sl, gate, gx, Inches(0.72), gate_w, Inches(0.35),
                     font_size=PPt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, milestone, gx, Inches(1.08), gate_w, Inches(0.28),
                     font_size=PPt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, date, gx, Inches(1.4), gate_w, Inches(0.25),
                     font_size=PPt(8), color=LIGHT_BG if color != LIGHT_BG else DARK_NAVY, align=PP_ALIGN.CENTER)
        add_text_box(sl, content, gx + Inches(0.1), Inches(1.7), gate_w - Inches(0.2), Inches(4.5),
                     font_size=PPt(9), color=DARK_NAVY)

    add_text_box(sl, '◀─────────────── RESEARCH部 開発ロードマップ（〜2027年2月） ───────────────▶',
                 Inches(0.3), Inches(6.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    slide_number_tag(sl, 2, TOTAL, is_landscape=True)

    # S3: G1マイルストーン
    sl = blank_slide(prs)
    title_bar(sl, 'G1マイルストーン 2026/09末', 'MVP完成 / テスト38件Pass / bandit 0件 / research_dashboard.html', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'G1 完了条件 / 受け入れ基準', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=DARK_NAVY)

    g1_tasks = [
        ['RES-A01 8社価格マトリクス実装', '自律COO', '未着手', '2026/08末', '¥0'],
        ['SHA-256 cache_key管理・IndexedDB v142連携', '自律COO', '未着手', '2026/08末', '¥0'],
        ['RES-A02 トレンドスコア算出エンジン', '自律COO', '未着手', '2026/09末', '¥0'],
        ['RETENTION_THRESHOLD=0.6 is_staple判定', '自律COO', '未着手', '2026/09末', '¥0'],
        ['research_dashboard.html LAYOUT_MASTER準拠実装', '自律COO', '未着手', '2026/09末', '¥0'],
        ['unittest 38テスト実装・全Pass確認', '自律COO', '未着手', '2026/09末', '¥0'],
        ['bandit 0件CI/CD設定', '自律COO', '未着手', '2026/08末', '¥0'],
        ['Cloud Run /health /price /trend エンドポイント実装', '自律COO', '未着手', '2026/09末', '¥0'],
    ]
    t = sl.shapes.add_table(len(g1_tasks) + 1, 5, Inches(0.3), Inches(1.0), Inches(10.8), Inches(4.5)).table
    for ci, h in enumerate(['タスク', '担当', '状態', '期限', '費用']):
        t.cell(0, ci).text = h
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(g1_tasks):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    add_text_box(sl, 'G1完了基準: RES-A01/A02がステージング環境で動作確認済 / 38テスト全Pass / bandit 0件 / APIレスポンス≤0.7秒',
                 Inches(0.3), Inches(5.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=ORANGE)
    slide_number_tag(sl, 3, TOTAL, is_landscape=True)

    # S4: G2マイルストーン
    sl = blank_slide(prs)
    title_bar(sl, 'G2マイルストーン 2026/11末', 'Keepa API / Google Trends実連携 / キャッシュTTL最適化', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'G2 完了条件 / 受け入れ基準', Inches(0.3), Inches(0.7), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=DARK_NAVY)

    g2_tasks = [
        ['Keepa API実連携（Amazon価格履歴取得）', '自律COO', '未着手', '2026/10末', 'API費用試算中'],
        ['Google Trends API実連携（growth_score精度向上）', '自律COO', '未着手', '2026/10末', 'API費用試算中'],
        ['IndexedDB v142 TTL最適化（価格:1h/トレンド:6h）', '自律COO', '未着手', '2026/11末', '¥0'],
        ['PriceMatrix Firestore永続化・履歴蓄積', '自律COO', '未着手', '2026/10末', '¥0〜¥450/月'],
        ['TrendDataPoint日次バッチ処理', '自律COO', '未着手', '2026/11末', '¥0〜¥150/月'],
        ['定番商品判定精度KPI初期計測開始', '自律COO', '未着手', '2026/11末', '¥0'],
        ['RES-A03統合ダッシュボード完成', '自律COO', '未着手', '2026/11末', '¥0'],
        ['Cloud Run本番デプロイ（常時稼働設定）', '自律COO', '未着手', '2026/11末', '¥750〜¥2,250/月'],
    ]
    t = sl.shapes.add_table(len(g2_tasks) + 1, 5, Inches(0.3), Inches(1.0), Inches(10.8), Inches(4.5)).table
    for ci, h in enumerate(['タスク', '担当', '状態', '期限', '費用']):
        t.cell(0, ci).text = h
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)
    for ri, row in enumerate(g2_tasks):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8)

    add_text_box(sl, 'G2完了基準: 実APIデータによるPriceMatrix更新確認 / TrendDataPoint日次蓄積開始 / APIレスポンス≤0.7秒維持',
                 Inches(0.3), Inches(5.7), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=ORANGE)
    slide_number_tag(sl, 4, TOTAL, is_landscape=True)

    # S5: FinOps計画
    sl = blank_slide(prs)
    title_bar(sl, 'FinOps計画 — コスト管理', 'MVP¥0 / G2 API費用試算中 / G3 Claude API¥2,250〜¥4,500 / 月額上限¥5,000', is_landscape=True)
    footer_bar(sl, FOOTER, is_landscape=True)

    add_text_box(sl, 'フェーズ別月次コスト', Inches(0.3), Inches(0.7), Inches(5.5), Inches(0.25),
                 font_size=PPt(10), bold=True, color=DARK_NAVY)
    finops_rows = [
        ['G0〜G1（MVP）', '¥0/月', 'Cloud Run/Firestore無料枠', 'AI自律開発・人件費¥0'],
        ['G2（実API連携）', 'API費用試算中', 'Keepa / Google Trends', '従量課金・要確認'],
        ['G3（AI強化）', '¥2,250〜¥4,500/月', '+ Claude API', '月1,000回呼出想定'],
        ['G4（Go-Live）', '¥3,000〜¥5,000/月', '全サービス稼働', '月額上限¥5,000'],
    ]
    t = sl.shapes.add_table(len(finops_rows) + 1, 4, Inches(0.3), Inches(1.0), Inches(5.5), Inches(2.2)).table
    for ci, h in enumerate(['フェーズ', '月額', '内訳', '備考']):
        t.cell(0, ci).text = h
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(9)
    for ri, row in enumerate(finops_rows):
        for ci, val in enumerate(row):
            t.cell(ri + 1, ci).text = val
            t.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)

    # Claude API詳細
    add_text_box(sl, 'Claude API コスト試算（G3以降）', Inches(6.0), Inches(0.7), Inches(5.0), Inches(0.25),
                 font_size=PPt(10), bold=True, color=DARK_NAVY)
    api_rows = [
        ['定番商品AI判定', '月500回', '¥1,125'],
        ['トレンド分析補助', '月300回', '¥675'],
        ['TODO自動起票生成', '月200回', '¥450'],
        ['合計', '月1,000回', '¥2,250/月'],
        ['上限設定', '月2,000回', '¥4,500/月'],
    ]
    t2 = sl.shapes.add_table(len(api_rows) + 1, 3, Inches(6.0), Inches(1.0), Inches(5.0), Inches(2.2)).table
    for ci, h in enumerate(['用途', '呼出回数', '月額概算']):
        t2.cell(0, ci).text = h
        t2.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t2.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(9)
    for ri, row in enumerate(api_rows):
        for ci, val in enumerate(row):
            t2.cell(ri + 1, ci).text = val
            t2.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)

    # GCPコスト内訳
    add_text_box(sl, 'GCP月次コスト内訳（G4 本番稼働時）', Inches(0.3), Inches(3.4), Inches(10.8), Inches(0.25),
                 font_size=PPt(10), bold=True, color=DARK_NAVY)
    gcp_rows = [
        ['Cloud Run（APIサーバー）', '¥750〜¥2,250', '従量制・リクエスト数に比例'],
        ['Firestore（PriceMatrix/ProductTrend）', '¥0〜¥450', '読取/書込回数に依存'],
        ['IndexedDB v142', '¥0', 'クライアント側（費用なし）'],
        ['Keepa API（G2以降）', '試算中', '商品数・呼出頻度に依存'],
        ['Google Trends API（G2以降）', '試算中', 'クォータ制・無料枠あり'],
        ['月額合計上限', '¥5,000', '月額上限¥5,000で厳格管理'],
    ]
    t3 = sl.shapes.add_table(len(gcp_rows) + 1, 3, Inches(0.3), Inches(3.7), Inches(10.8), Inches(2.5)).table
    for ci, h in enumerate(['サービス', '月額概算', '備考']):
        t3.cell(0, ci).text = h
        t3.cell(0, ci).text_frame.paragraphs[0].runs[0].font.bold = True
        t3.cell(0, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(9)
    for ri, row in enumerate(gcp_rows):
        for ci, val in enumerate(row):
            t3.cell(ri + 1, ci).text = val
            t3.cell(ri + 1, ci).text_frame.paragraphs[0].runs[0].font.size = PPt(8.5)

    GREEN_C = RGBColor(0x05, 0x96, 0x69)
    add_text_box(sl, 'FinOps原則: AI自律開発により開発人件費¥0 / MVP¥0 / G3以降Claude API¥2,250〜¥4,500 / 月額上限¥5,000で厳格管理',
                 Inches(0.3), Inches(6.3), Inches(10.8), Inches(0.3),
                 font_size=PPt(9), bold=True, color=GREEN_C)
    slide_number_tag(sl, 5, TOTAL, is_landscape=True)

    prs.save(output)
    print(f"✅ {output} 生成完了")


# ═══════════════════════════════════════════════════════════════════
# メイン
# ═══════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    os.makedirs(OUT_DIR, exist_ok=True)
    gen_brd()
    gen_srs()
    gen_seq()
    gen_ui()
    gen_phase()
    print("\n✅ RESEARCH部 5文書 生成完了")
